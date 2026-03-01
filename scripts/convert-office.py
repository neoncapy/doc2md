#!/usr/bin/env python3
"""
convert-office.py - Convert PPTX/DOCX/TXT files to Markdown with image extraction.

Part of the document conversion pipeline v3.0.
Handles:
  - PPTX: python-pptx for text + images (recursive GROUP shape extraction)
  - DOCX: pandoc for text, python-docx for images
  - TXT: direct read with YAML frontmatter wrapping

Usage:
  python3 convert-office.py <input-file> [--output-dir <dir>] [--skip-vision]

Output:
  {basename}.md          - Draft markdown with YAML frontmatter
  {basename}_images/     - Extracted images
  {basename}_manifest.json - Image manifest with metadata

Registry:
  Writes to ~/.claude/pipeline/conversion_registry.json on success.
"""

import argparse
import hashlib
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(levelname)s: %(message)s",
)
log = logging.getLogger("convert-office")

VERSION = "3.2.0"
PIPELINE_VERSION = "3.2.0"
REGISTRY_PATH = Path.home() / ".claude" / "pipeline" / "conversion_registry.json"

# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------

def sha256_file(path: Path) -> str:
    """Return SHA-256 hex digest of a file."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def _ensure_max_dimension(img_path, max_dim=2000):
    """Resize image in-place if either dimension exceeds max_dim.

    Uses PIL thumbnail() which preserves aspect ratio and uses LANCZOS
    resampling for quality.  Images already within limits are untouched.

    Returns True if image was resized, False otherwise.
    """
    try:
        from PIL import Image as _PilResize
        with _PilResize.open(img_path) as _img:
            if max(_img.size) <= max_dim:
                return False
            _orig = _img.size
            _img.thumbnail((max_dim, max_dim), _PilResize.LANCZOS)
            _img.save(img_path)
            log.info(f"  [resize] {Path(img_path).name}: "
                     f"{_orig[0]}x{_orig[1]} -> "
                     f"{_img.size[0]}x{_img.size[1]}")
            return True
    except Exception as _e:
        log.warning(f"Could not resize {Path(img_path).name}: {_e}")
        return False


# RC7: PUA (Private Use Area) → Unicode mapping for Symbol/Wingdings fonts.
# PowerPoint embeds Greek letters and math symbols as PUA codepoints when
# the source font is Symbol or Wingdings.  python-pptx extracts the raw
# codepoint, which renders as a blank box.  This table maps the most common
# PUA characters to their correct Unicode equivalents.
_PUA_TO_UNICODE: dict[str, str] = {
    # Greek lowercase
    "\uF061": "\u03B1",  # α
    "\uF062": "\u03B2",  # β
    "\uF067": "\u03B3",  # γ
    "\uF064": "\u03B4",  # δ
    "\uF065": "\u03B5",  # ε
    "\uF07A": "\u03B6",  # ζ
    "\uF068": "\u03B7",  # η
    "\uF071": "\u03B8",  # θ
    "\uF069": "\u03B9",  # ι
    "\uF06B": "\u03BA",  # κ
    "\uF06C": "\u03BB",  # λ
    "\uF06D": "\u03BC",  # μ
    "\uF06E": "\u03BD",  # ν
    "\uF078": "\u03BE",  # ξ
    "\uF070": "\u03C0",  # π
    "\uF072": "\u03C1",  # ρ
    "\uF073": "\u03C3",  # σ
    "\uF074": "\u03C4",  # τ
    "\uF075": "\u03C5",  # υ
    "\uF066": "\u03C6",  # φ
    "\uF063": "\u03C7",  # χ
    "\uF079": "\u03C8",  # ψ
    "\uF077": "\u03C9",  # ω
    # Greek uppercase
    "\uF044": "\u0394",  # Δ
    "\uF046": "\u03A6",  # Φ
    "\uF047": "\u0393",  # Γ
    "\uF04C": "\u039B",  # Λ
    "\uF050": "\u03A0",  # Π
    "\uF051": "\u0398",  # Θ
    "\uF053": "\u03A3",  # Σ
    "\uF055": "\u03A5",  # Υ
    "\uF057": "\u03A9",  # Ω
    "\uF058": "\u039E",  # Ξ
    "\uF059": "\u03A8",  # Ψ
    # Arrows
    "\uF0DE": "\u2192",  # →
    "\uF0DF": "\u2190",  # ←
    "\uF0E0": "\u2191",  # ↑
    "\uF0E1": "\u2193",  # ↓
    "\uF0DB": "\u21D4",  # ⇔
    "\uF0DC": "\u21D5",  # ⇕
    # Math operators and symbols
    "\uF0B1": "\u00B1",  # ±
    "\uF0B3": "\u2265",  # ≥
    "\uF0A3": "\u2264",  # ≤
    "\uF0B9": "\u2260",  # ≠
    "\uF0BB": "\u2248",  # ≈
    "\uF0B4": "\u00D7",  # ×
    "\uF0B7": "\u2022",  # •
    "\uF0AE": "\u00AE",  # ®
    "\uF0D3": "\u00A9",  # ©
    "\uF0E4": "\u2122",  # ™
    "\uF0A5": "\u221E",  # ∞
    "\uF0D9": "\u2329",  # 〈 (left angle bracket)
    "\uF0F1": "\u232A",  # 〉 (right angle bracket)
    "\uF0C5": "\u2215",  # ∕ (division slash)
    "\uF0D6": "\u2212",  # − (minus sign)
    "\uF0E5": "\u2202",  # ∂
    "\uF0C6": "\u0192",  # ƒ
    "\uF0C8": "\u2666",  # ♦
    "\uF0D0": "\u2014",  # — (em dash)
}

# Pre-compiled translation table for str.translate() — fast single-pass.
_PUA_TRANS_TABLE = str.maketrans(
    {ord(k): v for k, v in _PUA_TO_UNICODE.items()}
)


def clean_text(text: str) -> str:
    """Clean vertical tabs, normalize whitespace, and decode PUA chars."""
    if not text:
        return ""
    # RC7: Decode PUA characters from Symbol/Wingdings fonts
    text = text.translate(_PUA_TRANS_TABLE)
    # Replace vertical tabs with newlines
    text = text.replace("\x0b", "\n")
    # Collapse excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _sanitize_comment_text(text: str) -> str:
    """Collapse newlines/whitespace to single space for use inside HTML comments.

    Fix 3.11: Prevents IMAGE NOTEs from being trapped inside multi-line
    HTML comment blocks.  When slide titles or context strings contain
    newlines, <!-- IMAGE: ... --> comments span multiple lines, causing
    merge-image-notes.py to insert notes INSIDE the comment block.
    """
    if not text:
        return ""
    text = text.replace("\n", " ").replace("\r", " ")
    text = re.sub(r"\s+", " ", text).strip()
    return text


def ext_from_content_type(content_type: str) -> str:
    """Map content_type to file extension."""
    mapping = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/tiff": "tiff",
        "image/x-wmf": "wmf",
        "image/x-emf": "emf",
        "image/svg+xml": "svg",
    }
    return mapping.get(content_type, "png")


def convert_wmf_to_png(wmf_path: Path, output_dir: Path) -> Path | None:
    """Convert WMF/EMF to PNG using PIL first, then soffice fallback."""
    png_path = wmf_path.with_suffix(".png")

    # Try PIL first
    try:
        from PIL import Image
        img = Image.open(wmf_path)
        img.save(png_path)
        _ensure_max_dimension(png_path)
        log.info(f"  PIL converted {wmf_path.name} -> {png_path.name}")
        return png_path
    except Exception:
        pass

    # Fallback: soffice (use PATH lookup first, hardcoded fallback second)
    soffice_which = shutil.which("soffice")
    soffice = Path(soffice_which) if soffice_which else Path("/opt/homebrew/bin/soffice")
    if not soffice.exists():
        log.warning(f"  soffice not found, cannot convert {wmf_path.name}")
        return None

    try:
        result = subprocess.run(
            [str(soffice), "--headless", "--convert-to", "png",
             str(wmf_path), "--outdir", str(output_dir)],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and png_path.exists():
            _ensure_max_dimension(png_path)
            log.info(f"  soffice converted {wmf_path.name} -> {png_path.name}")
            return png_path
        else:
            log.warning(f"  soffice conversion failed for {wmf_path.name}")
            return None
    except Exception as e:
        log.warning(f"  soffice error for {wmf_path.name}: {e}")
        return None


def get_image_dimensions(img_path: Path) -> list[int]:
    """Return [width, height] using PIL, or [0, 0] on failure."""
    try:
        from PIL import Image
        with Image.open(img_path) as img:
            return list(img.size)
    except Exception:
        return [0, 0]


def is_blank_image(image_path: Path, threshold: float = 0.99) -> bool:
    """Check if image is blank (near-white / single solid color).

    Uses a sampling grid approach with pure PIL (no numpy dependency).
    Samples a 10x10 grid of pixels and checks if all are near-white (>240).
    Also checks per-channel extrema as a fast pre-filter.

    Returns True if the image is considered blank (conversion failure).

    F3 SAFEGUARDS (false-positive prevention):
      Safeguard 2 — Large file escape: files > 50KB are never blank.
        Prevents false positives on soffice chart renders (e.g. 113KB PNG
        with proportional hazards curves on white background).
      Safeguard 1 — Line drawing escape: if the image has significant pixel
        variation (grayscale std > 20) AND many unique colors (> 100), it
        contains real content and is not blank, regardless of white fraction.
        Prevents false positives on sparse line-art diagrams (e.g. economic
        step-function diagrams with labeled axes, 6.5KB, 773 unique colors).
    """
    try:
        import os as _os_blank
        _file_size = _os_blank.path.getsize(image_path)
    except OSError:
        _file_size = 0

    # Safeguard 2: large files are never blank conversion artifacts.
    # Real blanks from WMF failures are small (< 10KB). A 50KB+ image
    # always has genuine content even if it is mostly white (e.g. a chart
    # with two colored curves on a white canvas).
    if _file_size > 50000:
        return False

    try:
        from PIL import Image, ImageStat
        with Image.open(image_path) as img:
            rgb = img.convert("RGB")

            # Safeguard 1: line drawing / sparse diagram escape.
            # Compute grayscale std and unique color count using PIL only.
            # If the image has significant variation (std > 20) AND many
            # unique colors (> 100), it contains real content — thin lines,
            # axes, labels — and must not be marked blank even if the grid
            # sample lands almost entirely on the white background.
            try:
                gray = rgb.convert("L")
                _stat = ImageStat.Stat(gray)
                _gray_std = _stat.stddev[0]
                # Sample unique colors: resize to 250x200 for performance
                _w_orig, _h_orig = rgb.size
                if _w_orig * _h_orig > 50000:
                    _rgb_small = rgb.resize(
                        (min(_w_orig, 250), min(_h_orig, 200)),
                        Image.NEAREST)
                else:
                    _rgb_small = rgb
                _unique_colors = len(set(_rgb_small.getdata()))
                if _gray_std > 20 and _unique_colors > 100:
                    return False
            except Exception:
                pass  # If stat fails, continue to grid sampling below

            # Fast pre-filter: if ALL channels have minimum > 240,
            # every pixel is near-white — blank for sure.
            extrema = rgb.getextrema()  # [(min, max), (min, max), (min, max)]
            if all(lo > 240 for lo, hi in extrema):
                return True

            # Secondary check: sample a grid of pixels.
            # If > threshold fraction are near-white, treat as blank.
            w, h = rgb.size
            if w == 0 or h == 0:
                return True

            # RC6: Large-image safeguard for WMF decision trees.
            # Images larger than 500x500 can have ~98-99% white
            # background yet contain real content (decision trees,
            # flowcharts).  For these, count non-white pixels across
            # the full image — if > 5000 non-white pixels exist,
            # the image has real content regardless of white fraction.
            if w > 500 and h > 500:
                try:
                    import numpy as _np_blank
                    _arr = _np_blank.array(rgb)
                    # Pixel is non-white if ANY channel <= 240
                    _non_white = _np_blank.any(_arr <= 240, axis=2)
                    _non_white_count = int(_np_blank.sum(_non_white))
                    if _non_white_count > 5000:
                        return False
                except ImportError:
                    # numpy unavailable: use stricter threshold instead
                    threshold = 0.995

            sample_cols = max(1, w // 10)
            sample_rows = max(1, h // 10)
            white_count = 0
            total_count = 0
            for row_i in range(10):
                for col_i in range(10):
                    px_x = min(col_i * sample_cols, w - 1)
                    px_y = min(row_i * sample_rows, h - 1)
                    r, g, b = rgb.getpixel((px_x, px_y))
                    total_count += 1
                    if r > 240 and g > 240 and b > 240:
                        white_count += 1
            return (white_count / total_count) > threshold
    except Exception:
        return False


def update_registry(source_path: Path, output_path: Path, extractor: str,
                    image_index_meta: dict | None = None):
    """Add entry to conversion_registry.json.

    Deduplicates by source_hash (replaces existing entry for same hash).
    Uses atomic write (temp file + replace) to prevent registry corruption.
    Registry write failure is WARN only — does not abort pipeline.

    Args:
        source_path: Path to source file.
        output_path: Path to output .md file.
        extractor: Extractor identifier string.
        image_index_meta: Optional dict with R21 image index fields:
            image_index_path, image_index_generated_at, total_pages,
            pages_with_images, total_images_detected, substantive_images,
            has_testable_images.
    """
    try:
        REGISTRY_PATH.parent.mkdir(parents=True, exist_ok=True)

        registry = {"conversions": []}
        if REGISTRY_PATH.exists():
            try:
                with open(REGISTRY_PATH, "r") as f:
                    registry = json.load(f)
                if not isinstance(registry, dict) or "conversions" not in registry:
                    raise ValueError("malformed registry")
            except (json.JSONDecodeError, ValueError) as e:
                backup = REGISTRY_PATH.with_suffix(".json.bak")
                log.warning(f"Registry corrupt ({e}). Backing up to {backup} "
                            "and starting fresh.")
                try:
                    shutil.copy2(REGISTRY_PATH, backup)
                except Exception:
                    pass
                registry = {"conversions": []}

        source_hash = sha256_file(source_path)
        hash_key = f"sha256:{source_hash}"

        # Remove existing entry with same hash (dedup) — M1 fix
        existing = [c for c in registry.get("conversions", [])
                    if c.get("source_hash") != hash_key]
        entry = {
            "source_hash": hash_key,
            "source_path": str(source_path),
            "output_path": str(output_path),
            "extractor_used": extractor,
            "converted_at": datetime.now(timezone.utc).isoformat(),
            "pipeline_version": PIPELINE_VERSION,
        }

        # R21: Add image index metadata if available
        if image_index_meta and isinstance(image_index_meta, dict):
            for field in ("image_index_path", "image_index_generated_at",
                          "total_pages", "pages_with_images",
                          "total_images_detected", "substantive_images",
                          "has_testable_images"):
                if field in image_index_meta:
                    entry[field] = image_index_meta[field]

        existing.append(entry)
        registry["conversions"] = existing

        # Atomic write: write to temp file then replace (M3 fix)
        tmp_path = REGISTRY_PATH.with_suffix(".json.tmp")
        with open(tmp_path, "w") as f:
            json.dump(registry, f, indent=2)
        tmp_path.replace(REGISTRY_PATH)

        log.info(f"Registry updated: {REGISTRY_PATH}")
    except Exception as e:
        log.warning(f"Failed to update registry: {e}")


# ===========================================================================
# R19: Image Index Generation (PPTX / DOCX)
# ===========================================================================

def _is_decorative_image(page_num: int, total_pages: int,
                         width: int, height: int,
                         context: str, filename: str,
                         image_count_on_page: int,
                         is_chart: bool = False,
                         is_repeated: bool = False) -> bool:
    """Apply decorative image filtering heuristics.

    Conservative: when uncertain, classify as substantive.
    Charts are always substantive.

    Args:
        page_num: 1-indexed page/slide number.
        total_pages: Total pages/slides in document.
        width: Image width in pixels (0 if unknown).
        height: Image height in pixels (0 if unknown).
        context: Text context from the page/slide.
        filename: Source filename (for title/cover heuristics).
        image_count_on_page: Number of images on this page.
        is_chart: Whether this is a chart shape.
        is_repeated: Whether this image appears on >50% of pages
            (watermark/logo detection, heuristic #4).

    Returns:
        True if image is likely decorative.
    """
    # Charts are always substantive
    if is_chart:
        return False

    # Heuristic 1: Small images (< 50x50) are decorative (icons, bullets, tiny logos)
    if width > 0 and height > 0 and width < 50 and height < 50:
        return True

    # Fix 3.6/M4: Full-slide section dividers in PPTX
    # Very large images (>= 7000x4000) with minimal text (< 15 words)
    # are typically PPTX section divider backgrounds, not substantive content.
    if width >= 7000 and height >= 4000:
        if len(context.split()) < 15:
            return True

    # Heuristic 4: Repeated image on >50% of pages -> watermark/header/logo
    if is_repeated:
        return True

    # Heuristic 2: Page 1 with title/cover in filename -> decorative (logo/branding)
    if page_num == 1:
        fn_lower = filename.lower()
        if "title" in fn_lower or "cover" in fn_lower:
            return True

    # Heuristic 3: Last page with "thank" or "question" in context -> decorative
    if page_num == total_pages and context:
        ctx_lower = context.lower()
        if "thank" in ctx_lower or "question" in ctx_lower:
            return True

    # Heuristic 7: Figure-related keywords in context -> substantive
    if context:
        ctx_lower = context.lower()
        for kw in ("figure", "table", "diagram", "model", "chart", "graph",
                    "plot", "curve", "analysis"):
            if kw in ctx_lower:
                return False

    # Heuristic 5: Large images with text context -> substantive
    if width > 200 and height > 200 and context and len(context) > 20:
        return False

    # Default: substantive (conservative)
    return False


def _shape_has_image(shape) -> bool:
    """Check if a shape has an extractable image.

    Returns True for PICTURE shapes (type 13) and PLACEHOLDER shapes
    (type 14) that contain embedded images (e.g. PlaceholderPicture).
    python-pptx PlaceholderPicture inherits from both _InheritsDimensions
    (which sets shape_type=PLACEHOLDER) and Picture (which provides .image).
    Without this, ~59% of visual content in some PPTX files is missed.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        try:
            _ = shape.image
            return True
        except (AttributeError, ValueError):
            return False
    return False


def scan_pptx_for_image_index(pptx_path: Path) -> dict:
    """Scan a PPTX file and return per-slide image data for the image index.

    Returns dict with keys: total_pages, pages (list of per-slide dicts),
    and summary counts.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        log.warning(f"Cannot open PPTX for image indexing: {e}")
        return {"error": str(e), "total_pages": 0, "pages": []}

    total_slides = len(prs.slides)
    pages = []

    def _scan_shapes_recursive(shapes, images_list, charts_list):
        """Scan shapes for images/charts, recursing into GROUP shapes."""
        for shape in shapes:
            # GROUP shape: recurse into sub-shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _scan_shapes_recursive(shape.shapes, images_list, charts_list)
                continue

            if _shape_has_image(shape):
                w_px = 0
                h_px = 0
                try:
                    # python-pptx dimensions are in EMU; convert to approx px
                    # 1 inch = 914400 EMU, assume 96 DPI
                    if shape.width:
                        w_px = int(shape.width / 914400 * 96)
                    if shape.height:
                        h_px = int(shape.height / 914400 * 96)
                except Exception:
                    pass
                # Compute blob hash for watermark/repeated logo detection
                blob_hash = None
                try:
                    blob_hash = hashlib.sha256(shape.image.blob).hexdigest()
                except Exception:
                    pass
                images_list.append({
                    "width": w_px,
                    "height": h_px,
                    "name": shape.name,
                    "blob_hash": blob_hash,
                })
            elif hasattr(shape, "has_chart") and shape.has_chart:
                charts_list.append({"name": shape.name})
            # Also check for SmartArt / diagram types
            elif shape.shape_type == MSO_SHAPE_TYPE.IGX_GRAPHIC:
                charts_list.append({"name": shape.name, "type": "smartart"})
            elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
                charts_list.append({"name": shape.name, "type": "diagram"})

    for slide_num, slide in enumerate(prs.slides, 1):
        images = []
        charts = []
        text_parts = []

        _scan_shapes_recursive(slide.shapes, images, charts)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        text_parts.append(t)

        # Use slide title or first text as context
        context = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            context = slide.shapes.title.text.strip()[:150]
        elif text_parts:
            context = text_parts[0][:150]

        total_on_slide = len(images) + len(charts)
        pages.append({
            "page": slide_num,
            "image_count": total_on_slide,
            "picture_count": len(images),
            "chart_count": len(charts),
            "context": context,
            "image_details": images,
            "chart_details": charts,
        })

    return {
        "total_pages": total_slides,
        "pages": pages,
        "error": None,
    }


def scan_docx_for_image_index(docx_path: Path) -> dict:
    """Scan a DOCX file and return image data for the image index.

    DOCX has no native page concept. Images are grouped by their nearest
    preceding heading. We treat each heading section as a virtual 'page'.
    """
    from docx import Document

    try:
        doc = Document(str(docx_path))
    except Exception as e:
        log.warning(f"Cannot open DOCX for image indexing: {e}")
        return {"error": str(e), "total_pages": 0, "pages": []}

    # Count total images via relationships and collect blob hashes
    total_rels_images = 0
    # Map rel target_ref -> blob_hash for watermark detection
    rel_blob_hashes: dict[str, str] = {}  # target_ref -> blob_hash
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            total_rels_images += 1
            try:
                blob_hash = hashlib.sha256(rel.target_part.blob).hexdigest()
                rel_blob_hashes[rel.target_ref] = blob_hash
            except Exception:
                pass

    # Walk paragraphs to find images with heading context
    current_heading = "Document start"
    section_num = 0
    sections: dict[int, dict] = {}  # section_num -> {heading, images, context}

    # Ensure we have at least section 0
    sections[0] = {
        "heading": current_heading,
        "image_count": 0,
        "image_details": [],
        "context": "",
    }

    for para in doc.paragraphs:
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            section_num += 1
            current_heading = para.text.strip() or f"Section {section_num}"
            if section_num not in sections:
                sections[section_num] = {
                    "heading": current_heading,
                    "image_count": 0,
                    "image_details": [],
                    "context": current_heading[:150],
                }

        # Check for images in runs (inline images)
        if para.runs:
            for run in para.runs:
                drawing_els = run._element.findall(
                    ".//{http://schemas.openxmlformats.org/"
                    "wordprocessingml/2006/main}drawing"
                )
                if not drawing_els:
                    # Also check drawingML namespace
                    drawing_els = run._element.findall(
                        ".//{http://schemas.openxmlformats.org/"
                        "drawingml/2006/wordprocessingDrawing}inline"
                    )
                if drawing_els:
                    if section_num not in sections:
                        sections[section_num] = {
                            "heading": current_heading,
                            "image_count": 0,
                            "image_details": [],
                            "context": current_heading[:150],
                        }
                    sections[section_num]["image_count"] += len(drawing_els)
                    # Extract blob hashes from drawing elements for
                    # watermark/repeated image detection (heuristic #4).
                    for draw_el in drawing_els:
                        # Find the blip element which has the embed rId
                        blips = draw_el.findall(
                            ".//{http://schemas.openxmlformats.org/"
                            "drawingml/2006/main}blip"
                        )
                        for blip in blips:
                            # The r:embed attribute holds the relationship ID
                            embed_rid = blip.get(
                                "{http://schemas.openxmlformats.org/"
                                "officeDocument/2006/relationships}embed"
                            )
                            if embed_rid:
                                try:
                                    rel_part = doc.part.rels[embed_rid]
                                    blob_hash = hashlib.sha256(
                                        rel_part.target_part.blob
                                    ).hexdigest()
                                    sections[section_num]["image_details"].append({
                                        "blob_hash": blob_hash,
                                    })
                                except Exception:
                                    pass
                    # Try to get paragraph text as context
                    para_text = para.text.strip()[:150] if para.text.strip() else ""
                    if para_text and not sections[section_num]["context"]:
                        sections[section_num]["context"] = para_text

    # Convert sections dict to pages list, using section_num as virtual page
    pages = []
    for sec_num in sorted(sections.keys()):
        sec = sections[sec_num]
        pages.append({
            "page": sec_num + 1,  # 1-indexed
            "image_count": sec["image_count"],
            "picture_count": sec["image_count"],
            "chart_count": 0,
            "context": sec["context"] or sec["heading"],
            "image_details": sec.get("image_details", []),
            "chart_details": [],
        })

    # Total pages = number of heading-based sections (virtual pages).
    # KNOWN LIMITATION (m5): DOCX has no native page concept. This count
    # represents heading-based sections, NOT physical pages. Physical page
    # count requires full document rendering (e.g., via Word or LibreOffice)
    # which is not feasible in a headless pipeline. python-docx's
    # core_properties.pages is unreliable (often None or 0 unless the
    # document was last saved after a full repagination in Word).
    # For R21 registry, total_pages for DOCX = section count.
    total_pages = len(pages)

    return {
        "total_pages": total_pages,
        "pages": pages,
        "total_rels_images": total_rels_images,
        "error": None,
    }


def _build_pptx_pages_from_manifest(manifest_data: dict,
                                     total_slides: int) -> list[dict]:
    """Build per-slide page data from a manifest JSON (Root Cause B fix).

    Consumes the manifest that convert_pptx produced, using its comprehensive
    metadata (blank flags, duplicate flags, render PNGs, dimensions) instead
    of re-scanning the PPTX from scratch. This ensures the image index
    reflects exactly what was extracted, including:
    - Chart/SmartArt renders (type==chart_render)
    - Blank image flags (WMF conversion failures)
    - Duplicate image flags (is_duplicate)

    Args:
        manifest_data: Parsed JSON from {basename}_manifest.json.
        total_slides: Total slide count (for initializing empty slides).

    Returns:
        List of per-slide dicts compatible with the scan_result["pages"] format.
    """
    # Initialize all slides (even those with no images)
    slides: dict[int, dict] = {}
    for s in range(1, total_slides + 1):
        slides[s] = {
            "page": s,
            "image_count": 0,
            "unique_image_count": 0,
            "picture_count": 0,
            "chart_count": 0,
            "context": "",
            "image_details": [],
            "chart_details": [],
            "has_blank": False,
            "all_duplicate": True,  # Will be set False if any non-duplicate found
        }

    images = manifest_data.get("images", [])
    for img in images:
        slide_num = img.get("page")
        if slide_num is None or slide_num not in slides:
            continue

        slide_data = slides[slide_num]
        slide_data["image_count"] += 1

        is_blank = img.get("blank", False)
        is_duplicate = img.get("is_duplicate", False)
        is_decorative = img.get("decorative", False)
        source_format = img.get("source_format", "")
        is_render = source_format == "pptx_soffice_render"
        content_type = img.get("content_type", "")

        # Track blank slides
        if is_blank:
            slide_data["has_blank"] = True

        # Track whether ALL images on this slide are duplicates
        if not is_duplicate:
            slide_data["all_duplicate"] = False

        # Count by type
        if is_render or content_type in ("chart", "smartart"):
            slide_data["chart_count"] += 1
            slide_data["chart_details"].append({
                "name": img.get("source_shape", ""),
                "type": img.get("type_guess") or "chart",
                "blank": is_blank,
                "is_duplicate": is_duplicate,
            })
        else:
            slide_data["picture_count"] += 1

        # Count unique (non-duplicate) images
        if not is_duplicate:
            slide_data["unique_image_count"] += 1

        # Build image detail entry
        dims = img.get("dimensions", [0, 0])
        w = dims[0] if isinstance(dims, list) and len(dims) >= 2 else 0
        h = dims[1] if isinstance(dims, list) and len(dims) >= 2 else 0
        slide_data["image_details"].append({
            "width": w,
            "height": h,
            "name": img.get("source_shape", ""),
            "blob_hash": None,  # Not needed; manifest has richer data
            "blank": is_blank,
            "is_duplicate": is_duplicate,
            "is_decorative_manifest": is_decorative,
            "is_render": is_render,
            "filename": img.get("filename", ""),
            "size_bytes": img.get("size_bytes", 0),
        })

        # Use slide context from manifest
        section_ctx = img.get("section_context") or {}
        nearby = img.get("nearby_text") or ""
        heading = section_ctx.get("heading", "")
        if heading and not slide_data["context"]:
            slide_data["context"] = heading[:150]
        elif nearby and not slide_data["context"]:
            slide_data["context"] = nearby[:150]

    # Fix all_duplicate for slides with no images
    for s_data in slides.values():
        if s_data["image_count"] == 0:
            s_data["all_duplicate"] = False

    return [slides[s] for s in sorted(slides.keys())]


def generate_image_index(source_path: Path, output_dir: Path,
                         source_format: str,
                         manifest_path: Path | None = None) -> dict | None:
    """Generate a per-file image index manifest (R19).

    For PPTX: uses the manifest JSON (if available) as the PRIMARY data
    source. The manifest contains extracted images, render PNGs, blank
    flags, and duplicate flags — all of which are invisible to a raw
    PPTX re-scan. Falls back to PPTX scan only if no manifest exists.

    For DOCX: scans the source document directly (no manifest consumption
    yet — DOCX manifest structure differs).

    Args:
        source_path: Path to the source PPTX or DOCX file.
        output_dir: Directory where the image index will be written.
        source_format: Either 'pptx' or 'docx'.
        manifest_path: Optional path to the manifest JSON produced by
            convert_pptx/convert_docx. When provided for PPTX, used as
            primary data source instead of re-scanning the source file.

    Returns:
        Dict with image index metadata for R21 registry integration,
        or None if indexing failed entirely.
    """
    stem = source_path.stem
    index_path = output_dir / f"{stem}-image-index.md"

    log.info(f"Generating image index for: {source_path.name}")

    # Determine data source
    manifest_data = None
    used_manifest = False
    total_slides_from_pptx = 0

    if source_format == "pptx" and manifest_path and manifest_path.exists():
        try:
            with open(manifest_path, "r", encoding="utf-8") as f:
                manifest_data = json.load(f)
            log.info(f"Using manifest as primary data source: {manifest_path.name}")
            used_manifest = True
        except Exception as e:
            log.warning(f"Failed to load manifest ({e}), falling back to PPTX scan")

    # Get total slide count from PPTX (needed for manifest-based indexing)
    if source_format == "pptx" and used_manifest:
        try:
            from pptx import Presentation
            prs = Presentation(str(source_path))
            total_slides_from_pptx = len(prs.slides)
        except Exception as e:
            log.warning(f"Cannot open PPTX for slide count: {e}")
            total_slides_from_pptx = 0

    # Build scan result
    if used_manifest and total_slides_from_pptx > 0:
        pages = _build_pptx_pages_from_manifest(manifest_data, total_slides_from_pptx)
        scan_result = {
            "total_pages": total_slides_from_pptx,
            "pages": pages,
            "error": None,
        }
    elif source_format == "pptx":
        scan_result = scan_pptx_for_image_index(source_path)
    elif source_format == "docx":
        scan_result = scan_docx_for_image_index(source_path)
    else:
        log.warning(f"Image indexing not supported for format: {source_format}")
        return None

    if scan_result.get("error"):
        log.warning(f"Image index scan failed: {scan_result['error']}")
        # Write a minimal error manifest
        now_str = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M")
        error_content = (
            f"# Image Index: {stem}\n\n"
            f"Source: {source_path}\n"
            f"Generated: {now_str}\n"
            f"Pipeline version: v{PIPELINE_VERSION}\n\n"
            f"**Error:** Could not scan: {scan_result['error']}\n"
        )
        try:
            with open(index_path, "w", encoding="utf-8") as f:
                f.write(error_content)
            log.info(f"Image index written (error state): {index_path}")
        except Exception as e:
            log.warning(f"Failed to write image index error file: {e}")
        return None

    total_pages = scan_result["total_pages"]
    pages = scan_result["pages"]
    filename = source_path.name

    # Heuristic #4 pre-pass: build cross-page image frequency map for
    # watermark/repeated logo detection. Images (by blob_hash) appearing
    # on >50% of pages are classified as decorative.
    # (Only used when NOT consuming manifest — manifest has its own dedup)
    blob_hash_page_count: dict[str, int] = {}  # blob_hash -> number of pages
    if not used_manifest:
        for page_data in pages:
            page_hashes = set()
            for img_detail in page_data.get("image_details", []):
                bh = img_detail.get("blob_hash")
                if bh:
                    page_hashes.add(bh)
            for bh in page_hashes:
                blob_hash_page_count[bh] = blob_hash_page_count.get(bh, 0) + 1

    # Identify repeated blob hashes (>50% of total pages)
    repeated_hashes: set[str] = set()
    if not used_manifest and total_pages > 0:
        for bh, count in blob_hash_page_count.items():
            if count > total_pages * 0.5:
                repeated_hashes.add(bh)
    if repeated_hashes:
        log.info(f"Watermark detection: {len(repeated_hashes)} image(s) "
                 f"appear on >50% of pages — classified as decorative")

    # Apply decorative filtering per page
    pages_with_images = 0
    total_images_detected = 0
    total_manifest_decorative = 0  # M3/m3: track manifest-decorative images
    total_substantive_pages = 0
    total_decorative_pages = 0
    content_lost_pages = 0

    page_rows = []  # For the full page-by-page table
    substantive_rows = []  # For the substantive-only table

    for page_data in pages:
        page_num = page_data["page"]
        img_count = page_data["image_count"]
        context = page_data.get("context", "")

        if img_count == 0:
            continue

        pages_with_images += 1
        total_images_detected += img_count

        # Determine if page images are substantive or decorative
        is_chart = page_data.get("chart_count", 0) > 0
        has_substantive = False
        page_has_blank = page_data.get("has_blank", False)
        page_all_duplicate = page_data.get("all_duplicate", False)

        if used_manifest and source_format == "pptx":
            # Manifest-based classification (Root Cause B/C/D fix)
            # Skip slides where ALL images are duplicates (Root Cause D)
            if page_all_duplicate:
                # All images on this slide are duplicates of images on
                # other slides. Do not count as substantive.
                pass
            else:
                # Check each image from manifest data
                for img_detail in page_data.get("image_details", []):
                    # Skip duplicates — they don't contribute substantive status
                    if img_detail.get("is_duplicate", False):
                        continue

                    # Blank images = content lost (Root Cause C)
                    if img_detail.get("blank", False):
                        # Content exists but render failed — still counts
                        # as substantive (content is there, just lost)
                        has_substantive = True
                        break

                    # Renders from soffice are always substantive
                    if img_detail.get("is_render", False):
                        has_substantive = True
                        break

                    # Manifest decorative flag (size < 5KB in extraction)
                    if img_detail.get("is_decorative_manifest", False):
                        continue

                    # Apply dimension-based heuristics
                    w = img_detail.get("width", 0)
                    h = img_detail.get("height", 0)
                    if not _is_decorative_image(
                        page_num, total_pages, w, h, context,
                        filename, img_count, is_chart=False,
                        is_repeated=False,
                    ):
                        has_substantive = True
                        break

                # Charts/SmartArt from manifest are always substantive
                for chart_detail in page_data.get("chart_details", []):
                    if chart_detail.get("is_duplicate", False):
                        continue
                    has_substantive = True
                    break

        elif source_format == "pptx":
            # Fallback: PPTX scan-based classification (original logic)
            for img_detail in page_data.get("image_details", []):
                w = img_detail.get("width", 0)
                h = img_detail.get("height", 0)
                bh = img_detail.get("blob_hash")
                img_is_repeated = bh in repeated_hashes if bh else False
                if not _is_decorative_image(
                    page_num, total_pages, w, h, context,
                    filename, img_count, is_chart=False,
                    is_repeated=img_is_repeated,
                ):
                    has_substantive = True
                    break
            # Charts are always substantive
            if is_chart:
                has_substantive = True
            # SmartArt/diagrams in chart_details are also substantive
            for chart_detail in page_data.get("chart_details", []):
                has_substantive = True
                break
        elif source_format == "docx":
            # DOCX: check if all images in this section are repeated
            # (watermark/logo detection). Per-image dimensions not available
            # from scan, so use context-based heuristics (width/height = 0).
            img_details = page_data.get("image_details", [])
            if img_details:
                # Check each image individually
                for img_detail in img_details:
                    bh = img_detail.get("blob_hash")
                    img_is_repeated = bh in repeated_hashes if bh else False
                    if not _is_decorative_image(
                        page_num, total_pages, 0, 0, context,
                        filename, img_count, is_chart=False,
                        is_repeated=img_is_repeated,
                    ):
                        has_substantive = True
                        break
            else:
                # No per-image details available, use page-level heuristics
                if not _is_decorative_image(
                    page_num, total_pages, 0, 0, context,
                    filename, img_count, is_chart=False,
                ):
                    has_substantive = True

        # Determine labels for manifest-aware output
        substantive_label = "Yes" if has_substantive else "No"
        notes = []
        if used_manifest:
            if page_has_blank and has_substantive:
                notes.append("content_lost (WMF render failure)")
            if page_all_duplicate:
                substantive_label = "No (duplicate)"
                notes.append("all images are duplicates of earlier slides")

        # Use unique image count for manifest-based counting (Root Cause D)
        # M3 fix: subtract manifest-decorative images from the count so
        # that decorative icons are not tallied as substantive images.
        display_count = img_count
        manifest_dec_on_page = 0
        if used_manifest:
            unique_count = page_data.get("unique_image_count", img_count)
            # Count manifest-decorative images on this page
            for img_detail in page_data.get("image_details", []):
                if img_detail.get("is_duplicate", False):
                    continue  # already excluded by unique_count
                if img_detail.get("is_decorative_manifest", False):
                    manifest_dec_on_page += 1
            display_count = unique_count - manifest_dec_on_page
            # m-1 FIX: floor clamp to prevent negative display counts
            # if manifest data is inconsistent (more decorative than unique)
            display_count = max(0, display_count)
            total_manifest_decorative += manifest_dec_on_page

        context_display = context[:150] + "..." if len(context) > 150 else context
        if not context_display:
            context_display = "[no text context]"

        note_str = ""
        if notes:
            note_str = " | " + "; ".join(notes)

        page_rows.append({
            "page": page_num,
            "images": display_count,
            "substantive": substantive_label,
            "context": context_display,
            "note": note_str,
        })

        if has_substantive:
            total_substantive_pages += 1
            if page_has_blank and used_manifest:
                content_lost_pages += 1
            substantive_rows.append({
                "page": page_num,
                "images": display_count,
                "context": context_display,
                "note": note_str,
            })
        else:
            total_decorative_pages += 1

    # Estimate substantive image count (conservative: unique images on
    # substantive pages count as substantive)
    substantive_images = sum(r["images"] for r in substantive_rows)

    # Build the markdown index file
    now_str = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M")
    md_path = output_dir / f"{stem}.md"

    lines = []
    lines.append(f"# Image Index: {stem}")
    lines.append("")
    lines.append(f"Source: {source_path}")
    lines.append(f"Converted: {md_path}")
    lines.append(f"Generated: {now_str}")
    lines.append(f"Pipeline version: v{PIPELINE_VERSION}")
    if used_manifest:
        lines.append(f"Data source: manifest ({manifest_path.name})")
    lines.append("")
    # m1 fix: compute page-by-page sum for header consistency.
    # total_images_detected includes duplicates; table_sum excludes them.
    table_image_sum = sum(r["images"] for r in page_rows)
    decorative_images = total_manifest_decorative  # M3/m3: from manifest flags

    lines.append(f"Total pages: {total_pages}")
    lines.append(f"Pages with images: {pages_with_images}")
    lines.append(f"Total images detected: {total_images_detected}")
    # m1: show the table-consistent count when it differs from raw total
    if used_manifest and table_image_sum != total_images_detected:
        lines.append(f"Unique non-decorative images: {table_image_sum} "
                     "(excludes duplicates and manifest-decorative)")
    lines.append(f"Estimated substantive images: {substantive_images} "
                 "(after filtering)")
    # m2/m3: show per-image decorative count alongside per-page count
    if decorative_images > 0:
        lines.append(f"Manifest-decorative images: {decorative_images} "
                     "(small icons/badges flagged during extraction)")
    if content_lost_pages > 0:
        lines.append(f"Pages with content lost (WMF failure): {content_lost_pages}")
    lines.append("")
    lines.append("---")
    lines.append("")
    lines.append("## Page-by-Page Index")
    lines.append("")

    if used_manifest:
        lines.append("| Page | Images | Substantive | Context (first 150 chars) | Notes |")
        lines.append("|------|--------|-------------|---------------------------|-------|")
    else:
        lines.append("| Page | Images | Substantive | Context (first 150 chars) |")
        lines.append("|------|--------|-------------|---------------------------|")

    for row in page_rows:
        ctx_escaped = row["context"].replace("|", "\\|")
        if used_manifest:
            note_escaped = row.get("note", "").replace("|", "\\|")
            lines.append(
                f"| {row['page']} | {row['images']} | "
                f"{row['substantive']} | {ctx_escaped} | {note_escaped} |"
            )
        else:
            lines.append(
                f"| {row['page']} | {row['images']} | "
                f"{row['substantive']} | {ctx_escaped} |"
            )

    if not page_rows:
        if used_manifest:
            lines.append("| — | — | — | No images found in document | |")
        else:
            lines.append("| — | — | — | No images found in document |")

    lines.append("")
    lines.append("---")
    lines.append("")
    lines.append("## Substantive Images Only")
    lines.append("")

    if used_manifest:
        lines.append("| Page | Images | Context | Notes |")
        lines.append("|------|--------|---------|-------|")
    else:
        lines.append("| Page | Images | Context |")
        lines.append("|------|--------|---------|")

    for row in substantive_rows:
        ctx_escaped = row["context"].replace("|", "\\|")
        if used_manifest:
            note_escaped = row.get("note", "").replace("|", "\\|")
            lines.append(
                f"| {row['page']} | {row['images']} | {ctx_escaped} | {note_escaped} |"
            )
        else:
            lines.append(
                f"| {row['page']} | {row['images']} | {ctx_escaped} |"
            )

    if not substantive_rows:
        if used_manifest:
            lines.append("| — | — | No substantive images found | |")
        else:
            lines.append("| — | — | No substantive images found |")

    lines.append("")
    lines.append("---")
    lines.append("")
    lines.append("## Filtering Summary")
    lines.append("")
    lines.append(f"- Pages scanned: {total_pages}")
    lines.append(f"- Pages with images: {pages_with_images}")
    lines.append(f"- Pages classified as decorative: {total_decorative_pages}")
    lines.append(f"- Pages classified as substantive: {total_substantive_pages}")
    # m2/m3 fix: show per-image decorative count for clarity
    if total_manifest_decorative > 0:
        lines.append(f"- Individual images flagged decorative (manifest): "
                     f"{total_manifest_decorative}")
    if content_lost_pages > 0:
        lines.append(f"- Pages with content lost (WMF render failure): {content_lost_pages}")
    if used_manifest:
        lines.append("- Data source: manifest JSON (includes renders, blank/duplicate flags)")
        lines.append("- Filtering criteria applied: manifest blank flag, "
                     "manifest duplicate flag, manifest decorative flag, "
                     "dimensions (<50x50px), title/cover page, "
                     "last-page decorations, figure keywords")
    else:
        lines.append("- Filtering criteria applied: dimensions (<50x50px), "
                     "repeated image (>50% pages), title/cover page, "
                     "last-page decorations, figure keywords")

    # Write the index file
    try:
        with open(index_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines) + "\n")
        log.info(f"Image index written: {index_path} "
                 f"({pages_with_images} pages with images, "
                 f"{substantive_images} substantive)")
    except Exception as e:
        log.warning(f"Failed to write image index: {e}")
        return None

    # Return metadata for R21 registry integration
    return {
        "image_index_path": str(index_path),
        "image_index_generated_at": datetime.now(timezone.utc).isoformat(),
        "total_pages": total_pages,
        "pages_with_images": pages_with_images,
        "total_images_detected": total_images_detected,
        "substantive_images": substantive_images,
        "has_testable_images": substantive_images > 0,
    }


# ===========================================================================
# PPTX Conversion
# ===========================================================================

def convert_pptx(input_path: Path, output_dir: Path, skip_vision: bool) -> dict:
    """Convert PPTX to markdown with image extraction."""
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    log.info(f"Opening PPTX: {input_path.name}")
    prs = Presentation(str(input_path))

    basename = input_path.stem
    images_dir = output_dir / f"{basename}_images"
    images_dir.mkdir(parents=True, exist_ok=True)
    md_path = output_dir / f"{basename}.md"
    manifest_path = output_dir / f"{basename}_manifest.json"

    # Slide dimensions for position calculation
    slide_w = prs.slide_width or Emu(9144000)   # default 10 inches
    slide_h = prs.slide_height or Emu(6858000)  # default 7.5 inches

    total_slides = len(prs.slides)
    manifest_images = []
    md_lines = []
    global_img_counter = [0]  # mutable counter for recursive use

    # -----------------------------------------------------------------------
    # Deduplication tracking: maps SHA-256 hex digest -> first img_filename
    # QC-4: prevents sending byte-identical images to Opus vision multiple times.
    # -----------------------------------------------------------------------
    seen_hashes: dict[str, str] = {}  # hash -> first filename that used this hash

    # -----------------------------------------------------------------------
    # Recursive image extraction
    # -----------------------------------------------------------------------
    def extract_images_recursive(shapes, slide_num, slide_title, slide_text):
        """Extract images from shapes, recursing into GROUP shapes."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        slide_images = []

        for shape in shapes:
            # GROUP shape: recurse
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                sub_images = extract_images_recursive(
                    shape.shapes, slide_num, slide_title, slide_text,
                )
                slide_images.extend(sub_images)
                continue

            # CHART shape: note in manifest but cannot extract as image
            if hasattr(shape, "has_chart") and shape.has_chart:
                global_img_counter[0] += 1
                idx = global_img_counter[0]
                img_id = f"s{slide_num:02d}-chart{idx:02d}"
                manifest_images.append({
                    "id": img_id,
                    "file": None,
                    "slide": slide_num,
                    "content_type": "chart",
                    "size_bytes": 0,
                    "dimensions": [0, 0],
                    "position": _shape_position(shape, slide_w, slide_h),
                    "decorative": False,
                    "source_shape": shape.name,
                    "context": {
                        "slide_title": slide_title,
                        "nearby_text": slide_text[:100] if slide_text else "",
                    },
                    "type": "chart",
                    "note": "Chart rendered by PowerPoint; not extractable as image via python-pptx.",
                })
                slide_images.append({
                    "id": img_id,
                    "type": "chart",
                })
                continue

            # SMARTART shape: note in manifest, render via soffice fallback
            # python-pptx uses IGX_GRAPHIC (24) for SmartArt. Also check
            # DIAGRAM (21) and XML namespace as fallback for edge cases.
            is_smartart = False
            if shape.shape_type == MSO_SHAPE_TYPE.IGX_GRAPHIC:
                is_smartart = True
            elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
                is_smartart = True
            elif hasattr(shape, "_element"):
                # XML fallback: check for diagram namespace in shape XML
                shape_xml = shape._element.xml if hasattr(shape._element, "xml") else ""
                if "http://schemas.openxmlformats.org/drawingml/2006/diagram" in shape_xml:
                    is_smartart = True

            if is_smartart:
                global_img_counter[0] += 1
                idx = global_img_counter[0]
                img_id = f"s{slide_num:02d}-smartart{idx:02d}"
                log.info(f"  SmartArt detected: {shape.name} (type={shape.shape_type})")
                manifest_images.append({
                    "id": img_id,
                    "file": None,
                    "slide": slide_num,
                    "content_type": "smartart",
                    "size_bytes": 0,
                    "dimensions": [0, 0],
                    "position": _shape_position(shape, slide_w, slide_h),
                    "decorative": False,
                    "source_shape": shape.name,
                    "context": {
                        "slide_title": slide_title,
                        "nearby_text": slide_text[:100] if slide_text else "",
                    },
                    "type": "smartart",
                    "note": "SmartArt not extractable as image via python-pptx.",
                })
                slide_images.append({
                    "id": img_id,
                    "type": "smartart",
                })
                continue

            # PICTURE or PLACEHOLDER-with-image shape: extract image
            # PlaceholderPicture reports shape_type=PLACEHOLDER (14) but has
            # .image from its Picture parent class. _shape_has_image() checks
            # both types — fixes missing extraction of KM curves, extrapolation
            # plots, and other images embedded in placeholder shapes.
            if _shape_has_image(shape):
                try:
                    image = shape.image
                    content_type = image.content_type
                    ext = ext_from_content_type(content_type)
                    blob = image.blob

                    global_img_counter[0] += 1
                    idx = global_img_counter[0]
                    img_id = f"s{slide_num:02d}-img{idx:02d}"
                    img_filename = f"{img_id}.{ext}"
                    img_path = images_dir / img_filename

                    # QC-4: Deduplication — hash the blob before writing.
                    blob_hash = hashlib.sha256(blob).hexdigest()
                    is_duplicate = blob_hash in seen_hashes
                    original_filename = seen_hashes.get(blob_hash)
                    if not is_duplicate:
                        seen_hashes[blob_hash] = img_filename

                    # Write raw image (always, to preserve slide-to-image mapping)
                    with open(img_path, "wb") as f:
                        f.write(blob)

                    # Handle WMF/EMF conversion
                    final_path = img_path
                    if ext in ("wmf", "emf"):
                        converted = convert_wmf_to_png(img_path, images_dir)
                        if converted:
                            final_path = converted
                            img_filename = final_path.name
                            # Update dedup tracking to point at the PNG filename
                            if not is_duplicate:
                                seen_hashes[blob_hash] = img_filename

                    # Handle TIFF→PNG conversion for Opus vision compatibility
                    if final_path.suffix.lower() in (".tiff", ".tif"):
                        try:
                            from PIL import Image as PilImage
                            pil_img = PilImage.open(final_path)
                            png_path = Path(str(final_path).rsplit(".", 1)[0] + ".png")
                            pil_img.save(png_path, "PNG")
                            os.remove(final_path)
                            final_path = png_path
                            img_filename = final_path.name
                            # Update dedup tracking to the PNG filename
                            if not is_duplicate:
                                seen_hashes[blob_hash] = img_filename
                            log.info(f"  TIFF converted to PNG: {img_filename}")
                        except Exception as e:
                            log.warning(f"  TIFF→PNG conversion failed for {final_path.name}: {e}")

                    # Fix 1.8: resize oversized images before analysis
                    _ensure_max_dimension(final_path)

                    size_bytes = final_path.stat().st_size
                    decorative = size_bytes < 5000
                    dims = get_image_dimensions(final_path)

                    # QC-3: Blank image detection — catches WMF→PNG failures.
                    blank = False
                    if not decorative and final_path.suffix.lower() in (".png", ".jpg", ".jpeg", ".bmp", ".gif"):
                        blank = is_blank_image(final_path)
                        if blank:
                            log.warning(
                                f"WARNING: Blank image detected (likely WMF conversion failure): "
                                f"{final_path}"
                            )

                    pos = _shape_position(shape, slide_w, slide_h)

                    entry = {
                        "id": img_id,
                        "file": img_filename,
                        "slide": slide_num,
                        "content_type": content_type,
                        "size_bytes": size_bytes,
                        "dimensions": dims,
                        "position": pos,
                        "decorative": decorative,
                        "source_shape": shape.name,
                        "context": {
                            "slide_title": slide_title,
                            "nearby_text": slide_text[:100] if slide_text else "",
                        },
                        # QC-3: blank detection flag
                        "blank": blank,
                        # QC-4: deduplication flags
                        "is_duplicate": is_duplicate,
                        "original": original_filename,
                    }
                    manifest_images.append(entry)
                    slide_images.append({
                        "id": img_id,
                        "file": img_filename,
                        "decorative": decorative,
                        "blank": blank,
                        "is_duplicate": is_duplicate,
                        "dims": dims,
                    })

                    if is_duplicate:
                        log.info(
                            f"  {img_filename} ({size_bytes}B) - duplicate of {original_filename}"
                        )
                    elif decorative:
                        log.info(f"  {img_filename} ({size_bytes}B) - decorative")
                    else:
                        log.info(f"  {img_filename} ({size_bytes}B)")

                except Exception as e:
                    log.warning(f"  Failed to extract image from shape '{shape.name}': {e}")
                    continue

        return slide_images

    # -----------------------------------------------------------------------
    # Process each slide
    # -----------------------------------------------------------------------
    for slide_idx, slide in enumerate(prs.slides, start=1):
        log.info(f"Processing slide {slide_idx}/{total_slides}")

        # Extract slide title
        slide_title = "Untitled"
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = clean_text(slide.shapes.title.text)

        # Extract all text from slide
        text_parts = []
        ole_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                # m13: Determine if this shape is the slide title — headings
                # get plain text; body shapes preserve run-level bold/italic.
                _is_title_shape = (
                    slide.shapes.title is not None
                    and shape is slide.shapes.title
                )
                for para in shape.text_frame.paragraphs:
                    if _is_title_shape:
                        # Title shapes: use plain text (F11 strips ** later)
                        para_text = clean_text(para.text)
                    else:
                        # Body shapes: preserve bold/italic from runs
                        run_parts = []
                        for run in para.runs:
                            rtext = run.text
                            if rtext.strip():
                                if run.font.bold:
                                    rtext = f"**{rtext}**"
                                if run.font.italic:
                                    rtext = f"*{rtext}*"
                            run_parts.append(rtext)
                        para_text = clean_text("".join(run_parts))
                    if para_text:
                        text_parts.append(para_text)
            # F4: TABLE shape extraction (MSO_SHAPE_TYPE=19)
            # TABLE shapes are silently skipped by has_text_frame — extract
            # their content as a markdown pipe table.
            elif shape.has_table:
                table = shape.table
                rows_md = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_md.append("| " + " | ".join(cells) + " |")
                if rows_md:
                    separator = "| " + " | ".join(
                        ["---"] * len(table.columns)
                    ) + " |"
                    table_md = (
                        rows_md[0] + "\n"
                        + separator + "\n"
                        + "\n".join(rows_md[1:])
                    )
                    text_parts.append(table_md)
            # F5: OLE/embedded object placeholder (MSO_SHAPE_TYPE=7)
            # python-pptx cannot render OLE content (e.g. Equation.3 objects).
            # Insert a visible placeholder so the reader knows something exists.
            elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                prog_id = ""
                ole_elem = shape.element
                for child in ole_elem.iter():
                    if "progId" in child.attrib:
                        prog_id = child.attrib["progId"]
                        break
                obj_type = prog_id if prog_id else "embedded object"
                ole_parts.append(
                    f"[{obj_type}: not rendered — see source slide {slide_idx}]"
                )

        slide_text = "\n".join(text_parts)
        if ole_parts:
            ole_block = "\n".join(ole_parts)
            slide_text = (slide_text + "\n\n" + ole_block).strip()

        # F10: Extract speaker notes from notesSlide XML
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text = (
                        slide_text + f"\n\n> **Speaker notes:** {notes_text}\n"
                    ).strip()
        except Exception:
            pass  # Slide has no notes or notes are inaccessible

        # Extract images (recursive)
        slide_images = extract_images_recursive(
            slide.shapes, slide_idx, slide_title, slide_text,
        )

        # Build markdown for this slide
        # F11: Strip bold markers from heading text (multi-run source formatting
        # can embed ** inside heading text; headings are plain text by definition)
        slide_title_clean = re.sub(r'\*\*', '', slide_title)
        md_lines.append(f"## Slide {slide_idx}: {slide_title_clean}")
        md_lines.append("")

        if slide_text:
            md_lines.append(slide_text)
            md_lines.append("")

        for img_info in slide_images:
            if img_info.get("type") == "chart":
                md_lines.append(
                    f"<!-- CHART: {img_info['id']} | "
                    f"Note: Chart rendered by PowerPoint, not extractable -->"
                )
            elif img_info.get("type") == "smartart":
                md_lines.append(
                    f"<!-- SMARTART: {img_info['id']} | "
                    f"Note: SmartArt rendered by PowerPoint, not extractable -->"
                )
            else:
                fname = img_info["file"]
                # Issue 6 fix: use subdirectory-relative path instead of bare filename
                img_rel_path = f"{basename}_images/{fname}"
                dims = img_info.get("dims", [0, 0])
                dec = " | Decorative: yes" if img_info.get("decorative") else ""
                # Fix 3.6/M18: Propagate duplicate flag to IMAGE comment
                dup = " | Duplicate: yes" if img_info.get("is_duplicate") else ""
                if img_info.get("is_duplicate") and not dec:
                    dec = " | Decorative: yes"
                # Fix 3.11: sanitize context to prevent multi-line comments
                _ctx = _sanitize_comment_text(slide_title)
                md_lines.append(
                    f"<!-- IMAGE: {img_rel_path} | "
                    f"Size: {dims[0]}x{dims[1]}{dec}{dup} | "
                    f"Context: {_ctx} -->"
                )
                md_lines.append(f"![Image {img_info['id']}]({img_rel_path})")
            md_lines.append("")

    # -----------------------------------------------------------------------
    # LibreOffice slide-render fallback for chart/SmartArt slides
    # -----------------------------------------------------------------------
    # Identify slides with chart/SmartArt placeholder entries (file=None).
    # These shapes were detected but could not be extracted as images by
    # python-pptx. Render those slides via soffice and add the PNGs to the
    # manifest so the Opus vision pipeline can describe them.
    chart_smartart_entries = [
        entry for entry in manifest_images
        if entry.get("type") in ("chart", "smartart") and entry.get("file") is None
    ]
    slides_needing_render = sorted(set(
        entry["slide"] for entry in chart_smartart_entries
    ))

    if slides_needing_render:
        log.info(
            f"Chart/SmartArt fallback: {len(slides_needing_render)} slide(s) "
            f"need soffice rendering: {slides_needing_render}"
        )
        # Issue 7: Use shutil.which() with hardcoded fallback for portability
        soffice_which = shutil.which("soffice")
        soffice = Path(soffice_which) if soffice_which else Path("/opt/homebrew/bin/soffice")
        if not soffice.exists():
            log.warning(
                f"soffice not found (tried PATH lookup and /opt/homebrew/bin/soffice) — "
                "cannot render chart/SmartArt slides. Skipping fallback."
            )
        else:
            render_tmpdir = tempfile.mkdtemp(prefix="soffice-chart-render-")
            render_tmpdir_path = Path(render_tmpdir)
            # Isolated user installation to avoid locking conflicts with
            # any running LibreOffice instance.
            env_dir = f"/tmp/soffice-chart-{os.getpid()}"
            env_dir_path = Path(env_dir)

            try:
                # Issue 1: Two-step approach — PPTX->PDF (soffice) then
                # PDF->PNG per slide (pdftoppm). LibreOffice --convert-to png
                # on PPTX only renders slide 1 (Impress limitation).
                # Step 1: Convert PPTX to PDF (all slides)
                pdf_cmd = [
                    str(soffice),
                    "--headless",
                    f"-env:UserInstallation=file:///{env_dir}",
                    "--convert-to", "pdf",
                    "--outdir", str(render_tmpdir_path),
                    str(input_path),
                ]
                log.info(f"Running soffice PPTX->PDF: {' '.join(pdf_cmd[:4])} ...")
                pdf_result = subprocess.run(
                    pdf_cmd, capture_output=True, text=True, timeout=120,
                )
                if pdf_result.returncode != 0:
                    log.warning(
                        f"soffice PDF conversion returned non-zero exit code "
                        f"({pdf_result.returncode}). stderr: {pdf_result.stderr[:500]}"
                    )

                # Find the generated PDF
                pdf_path = render_tmpdir_path / f"{input_path.stem}.pdf"
                if not pdf_path.exists():
                    log.warning(
                        f"soffice produced no PDF file at {pdf_path}. "
                        "Skipping chart/SmartArt fallback."
                    )
                else:
                    log.info(f"soffice produced PDF: {pdf_path.name}")

                    # Step 2: Render specific slides from PDF using pdftoppm
                    # pdftoppm uses 1-based page numbers: -f N -l N for page N.
                    # Issue 8: 60s timeout per pdftoppm call (single slide).
                    pdftoppm_path = shutil.which("pdftoppm")
                    if not pdftoppm_path:
                        pdftoppm_path = "/opt/homebrew/bin/pdftoppm"
                    if not Path(pdftoppm_path).exists():
                        log.warning(
                            f"pdftoppm not found (tried PATH lookup and "
                            "/opt/homebrew/bin/pdftoppm) — cannot render "
                            "individual slides. Skipping fallback."
                        )
                    else:
                        # Build slide_png_map: slide_number -> rendered PNG path
                        # 300 DPI balances render quality for Opus vision
                        # analysis against file size. Produces ~3-5 MB per
                        # slide. 200 DPI is an alternative for large decks.
                        RENDER_DPI = "300"
                        slide_png_map: dict[int, Path] = {}
                        for slide_idx in slides_needing_render:
                            # Use per-slide unique prefix to eliminate glob
                            # ambiguity (Issue A: prefix collision fix).
                            slide_prefix = f"slide-{slide_idx:03d}"
                            output_prefix = str(
                                render_tmpdir_path / slide_prefix
                            )
                            render_cmd = [
                                pdftoppm_path,
                                "-png", "-r", RENDER_DPI,
                                "-f", str(slide_idx),
                                "-l", str(slide_idx),
                                str(pdf_path),
                                output_prefix,
                            ]
                            log.info(
                                f"Rendering slide {slide_idx} via pdftoppm "
                                f"({RENDER_DPI} DPI) ..."
                            )
                            try:
                                render_result = subprocess.run(
                                    render_cmd, capture_output=True,
                                    text=True, timeout=60,
                                )
                                if render_result.returncode != 0:
                                    log.warning(
                                        f"pdftoppm failed for slide {slide_idx}: "
                                        f"{render_result.stderr[:300]}"
                                    )
                                    continue
                            except subprocess.TimeoutExpired:
                                log.warning(
                                    f"pdftoppm timed out (60s) for slide "
                                    f"{slide_idx}. Skipping this slide."
                                )
                                continue

                            # pdftoppm names output as {prefix}-{padded_page}.png
                            # With unique per-slide prefix, glob for just this
                            # slide's output. Single -f N -l N call produces
                            # exactly one file.
                            rendered_pngs = sorted(
                                render_tmpdir_path.glob(f"{slide_prefix}-*.png")
                            )
                            if rendered_pngs:
                                slide_png_map[slide_idx] = rendered_pngs[0]
                            else:
                                log.warning(
                                    f"pdftoppm produced no PNG for slide "
                                    f"{slide_idx}."
                                )

                        if not slide_png_map:
                            log.warning(
                                "pdftoppm produced no PNG files for any "
                                "needed slides. Skipping fallback."
                            )
                        else:
                            log.info(
                                f"pdftoppm rendered {len(slide_png_map)} "
                                f"slide(s): {list(slide_png_map.keys())}"
                            )

                        # Build a lookup: slide_idx -> list of chart/smartart
                        # IDs on that slide, for post-processing md_lines.
                        slide_to_ids: dict[int, list[str]] = {}
                        for entry in chart_smartart_entries:
                            sid = entry["slide"]
                            slide_to_ids.setdefault(sid, []).append(entry["id"])

                        # Determine type of original entries per slide for
                        # type_guess (Issue 3)
                        slide_to_types: dict[int, set[str]] = {}
                        for entry in chart_smartart_entries:
                            sid = entry["slide"]
                            slide_to_types.setdefault(sid, set()).add(
                                entry.get("type", "chart")
                            )

                        # Process each slide that needs rendering
                        rendered_count = 0
                        for slide_idx in slides_needing_render:
                            src_png = slide_png_map.get(slide_idx)
                            if not src_png or not src_png.exists():
                                log.warning(
                                    f"No rendered PNG found for slide "
                                    f"{slide_idx}. Available: "
                                    f"{list(slide_png_map.keys())}"
                                )
                                continue

                            # Determine destination filename and copy
                            global_img_counter[0] += 1
                            render_idx = global_img_counter[0]
                            render_id = (
                                f"s{slide_idx:02d}-render{render_idx:02d}"
                            )
                            render_filename = f"{render_id}.png"
                            dest_path = images_dir / render_filename
                            shutil.copy2(src_png, dest_path)
                            _ensure_max_dimension(dest_path)

                            # Get dimensions
                            dims = get_image_dimensions(dest_path)
                            size_bytes = dest_path.stat().st_size

                            # Blank detection
                            blank = False
                            if size_bytes >= 5000:
                                blank = is_blank_image(dest_path)
                                if blank:
                                    log.warning(
                                        f"WARNING: Blank rendered slide "
                                        f"detected: {render_filename} "
                                        f"(slide {slide_idx})"
                                    )

                            # SHA-256 dedup against previously seen images
                            file_hash = sha256_file(dest_path)
                            is_duplicate = file_hash in seen_hashes
                            original_filename = seen_hashes.get(file_hash)
                            if not is_duplicate:
                                seen_hashes[file_hash] = render_filename

                            # Retrieve context from original entries
                            original_entries = [
                                e for e in chart_smartart_entries
                                if e["slide"] == slide_idx
                            ]
                            slide_title_ctx = ""
                            nearby_text_ctx = ""
                            if original_entries:
                                ctx = original_entries[0].get("context", {})
                                slide_title_ctx = ctx.get("slide_title", "")
                                nearby_text_ctx = ctx.get("nearby_text", "")

                            # Issue 3: Set type_guess based on original
                            # entry types. "chart" -> "chart" for richer
                            # persona activation; "smartart" -> "diagram".
                            slide_types = slide_to_types.get(
                                slide_idx, set()
                            )
                            if "chart" in slide_types:
                                render_type_guess = "chart"
                            elif "smartart" in slide_types:
                                render_type_guess = "diagram"
                            else:
                                render_type_guess = "chart"

                            # Add to manifest_images (will be picked up
                            # by normalization loop since file is not None)
                            manifest_images.append({
                                "id": render_id,
                                "file": render_filename,
                                "slide": slide_idx,
                                "content_type": "image/png",
                                "size_bytes": size_bytes,
                                "dimensions": dims,
                                "position": {
                                    "left": 0, "top": 0,
                                    "width": 100, "height": 100,
                                },
                                "decorative": False,
                                "source_shape": "soffice_render",
                                "context": {
                                    "slide_title": slide_title_ctx,
                                    "nearby_text": nearby_text_ctx,
                                },
                                "type": "chart_render",
                                "type_guess": render_type_guess,
                                "note": (
                                    f"Full slide {slide_idx} rendered via "
                                    f"LibreOffice (contains chart/SmartArt "
                                    f"not extractable by python-pptx)."
                                ),
                                "blank": blank,
                                "is_duplicate": is_duplicate,
                                "original": original_filename,
                            })

                            rendered_count += 1
                            if is_duplicate:
                                log.info(
                                    f"  {render_filename} (slide "
                                    f"{slide_idx}, {size_bytes}B) - "
                                    f"duplicate of {original_filename}"
                                )
                            elif blank:
                                log.info(
                                    f"  {render_filename} (slide "
                                    f"{slide_idx}, {size_bytes}B) - blank"
                                )
                            else:
                                log.info(
                                    f"  {render_filename} (slide "
                                    f"{slide_idx}, {dims[0]}x{dims[1]}, "
                                    f"{size_bytes}B)"
                                )

                            # Issue 4: Only replace placeholders if render
                            # is NOT blank. Blank renders get a failure
                            # comment.
                            # Issue 2: For multi-chart slides, replace only
                            # the FIRST placeholder with the image ref;
                            # subsequent placeholders get a cross-reference.
                            chart_ids = slide_to_ids.get(slide_idx, [])
                            if blank:
                                # Replace placeholders with failure note
                                for chart_id in chart_ids:
                                    c_pat = f"<!-- CHART: {chart_id} |"
                                    s_pat = f"<!-- SMARTART: {chart_id} |"
                                    for i, line in enumerate(md_lines):
                                        if c_pat in line:
                                            md_lines[i] = (
                                                f"<!-- CHART: {chart_id}"
                                                f" — could not render "
                                                f"(blank output) -->"
                                            )
                                            break
                                        elif s_pat in line:
                                            md_lines[i] = (
                                                f"<!-- SMARTART: "
                                                f"{chart_id} — could not"
                                                f" render (blank output)"
                                                f" -->"
                                            )
                                            break
                            else:
                                # Replace placeholders with image refs
                                first_replacement = True
                                for chart_id in chart_ids:
                                    c_pat = f"<!-- CHART: {chart_id} |"
                                    s_pat = f"<!-- SMARTART: {chart_id} |"
                                    for i, line in enumerate(md_lines):
                                        if c_pat in line or s_pat in line:
                                            if first_replacement:
                                                # Issue 6 fix: use subdirectory-relative path
                                                render_rel_path = f"{basename}_images/{render_filename}"
                                                # Fix 3.11: sanitize context
                                                _ctx2 = _sanitize_comment_text(
                                                    slide_title_ctx)
                                                md_lines[i] = (
                                                    f"<!-- IMAGE: "
                                                    f"{render_rel_path}"
                                                    f" | Size: "
                                                    f"{dims[0]}x"
                                                    f"{dims[1]} | "
                                                    f"Source: soffice "
                                                    f"slide render | "
                                                    f"Context: "
                                                    f"{_ctx2}"
                                                    f" -->\n"
                                                    f"![Figure "
                                                    f"{render_id}: "
                                                    f"Slide {slide_idx}"
                                                    f" chart/SmartArt "
                                                    f"render]"
                                                    f"({render_rel_path}"
                                                    f")"
                                                )
                                                first_replacement = False
                                            else:
                                                md_lines[i] = (
                                                    f"<!-- See "
                                                    f"{render_filename}"
                                                    f" above (covers all"
                                                    f" chart/SmartArt on"
                                                    f" slide "
                                                    f"{slide_idx}) -->"
                                                )
                                            break

                        log.info(
                            f"Chart/SmartArt fallback complete: "
                            f"{rendered_count}/"
                            f"{len(slides_needing_render)} "
                            f"slide(s) rendered successfully."
                        )

                    # Clean up intermediate PDF
                    try:
                        pdf_path.unlink()
                    except Exception as e:
                        log.debug(
                            f"Could not remove intermediate PDF: {e}"
                        )

            except subprocess.TimeoutExpired:
                log.warning(
                    "soffice PDF conversion timed out after 120s. "
                    "Skipping chart/SmartArt fallback."
                )
            except Exception as e:
                log.warning(
                    f"soffice slide render failed: {e}. "
                    f"Skipping chart/SmartArt fallback."
                )
            finally:
                # Clean up temp directories
                # Issue 6: Log cleanup exceptions at debug level
                try:
                    shutil.rmtree(render_tmpdir_path, ignore_errors=True)
                except Exception as e:
                    log.debug(f"Cleanup warning: {e}")
                try:
                    if env_dir_path.exists():
                        shutil.rmtree(env_dir_path, ignore_errors=True)
                except Exception as e:
                    log.debug(f"Cleanup warning: {e}")

    # -----------------------------------------------------------------------
    # Determine document title
    # -----------------------------------------------------------------------
    doc_title = basename.replace("_", " ").replace("-", " ")
    # Try to get title from first slide
    if prs.slides and prs.slides[0].shapes.title:
        first_title = clean_text(prs.slides[0].shapes.title.text)
        if first_title:
            doc_title = first_title

    # -----------------------------------------------------------------------
    # Write markdown file
    # -----------------------------------------------------------------------
    frontmatter = f"""---
title: "{doc_title}"
source_file: "{input_path.name}"
source_format: "pptx"
conversion_tool: "convert-office.py v{VERSION}"
conversion_date: "{datetime.now().strftime('%Y-%m-%dT%H:%M:%S')}"
document_type: "presentation"
fidelity_standard: "visual_content"
slides: {total_slides}
total_images: {len([i for i in manifest_images if i.get('file')])}
images_directory: "{basename}_images/"
pipeline_version: "{PIPELINE_VERSION}"
image_notes: {"pending" if len([i for i in manifest_images if i.get('file')]) > 0 else "none"}
---
"""

    # Issue 5 fix: Strip redundant backref comments from chart/SmartArt dedup.
    # These "<!-- See sXX-renderYYY.png above (...) -->" lines carry no
    # information (the render is already embedded once per slide) and add
    # significant noise (e.g. 178 lines / 28% in DOC-2).
    md_body = "\n".join(md_lines)
    md_body = re.sub(r'\s*<!-- *See.*?above.*?-->\s*', '\n', md_body)
    # Collapse any triple+ blank lines left after stripping
    md_body = re.sub(r'\n{3,}', '\n\n', md_body)

    with open(md_path, "w", encoding="utf-8") as f:
        f.write(frontmatter)
        f.write("\n")
        f.write(f"# {doc_title}\n\n")
        f.write(md_body)

    log.info(f"Markdown written: {md_path}")

    # -----------------------------------------------------------------------
    # Write manifest
    # Normalized to match convert-paper.py format so that prepare-image-analysis.py
    # (Step 3 in run-pipeline.py) can consume it without format-specific branches.
    # Required fields per image: filename, figure_num, width, height, page.
    # Optional but expected: type_guess, section_context, nearby_text,
    #                        detected_caption, file_path.
    # -----------------------------------------------------------------------
    normalized_images = []
    figure_counter = 0
    for entry in manifest_images:
        if not entry.get("file"):
            # Skip chart/SmartArt placeholder entries (no image file) —
            # these have type=="chart" or type=="smartart" and file=None
            continue
        figure_counter += 1
        dims = entry.get("dimensions", [0, 0])
        w = dims[0] if isinstance(dims, list) and len(dims) >= 2 else 0
        h = dims[1] if isinstance(dims, list) and len(dims) >= 2 else 0
        img_filename = entry["file"]
        img_abs_path = str(images_dir / img_filename)
        context = entry.get("context", {}) or {}
        section_context = {
            "heading": context.get("slide_title") or "Unknown Section",
            "heading_level": 2,
        }
        nearby_text = context.get("nearby_text") or None
        normalized_images.append({
            # Fields required by prepare-image-analysis.py (no .get() fallback)
            "figure_num": figure_counter,
            "filename": img_filename,
            "width": w,
            "height": h,
            # Optional fields used by prepare-image-analysis.py
            "page": entry.get("slide"),
            "type_guess": entry.get("type_guess") or None,  # Issue 3: chart_render entries set this explicitly
            "section_context": section_context,
            "nearby_text": nearby_text,
            "detected_caption": entry.get("detected_caption") or None,
            "file_path": img_abs_path,
            # Original fields preserved for any other downstream consumers
            "id": entry.get("id"),
            "content_type": entry.get("content_type"),
            "size_bytes": entry.get("size_bytes", 0),
            "dimensions": dims,
            "position": entry.get("position"),
            "decorative": entry.get("decorative", False),
            "source_shape": entry.get("source_shape"),
            "analysis_status": "pending",
            "description": (
                f"Figure {figure_counter} from slide {entry.get('slide', '?')}"
                f" (chart/SmartArt render)"
                if entry.get("type") == "chart_render"
                else (
                    f"Figure {figure_counter} from slide {entry.get('slide', '?')}"
                    if entry.get("slide")
                    else f"Figure {figure_counter} from document"
                )
            ),
            "source_format": (
                "pptx_soffice_render"
                if entry.get("type") == "chart_render"
                else "pptx_embedded"
            ),
            # QC-3: blank image detection (WMF conversion failure indicator)
            "blank": entry.get("blank", False),
            # QC-4: deduplication fields
            "is_duplicate": entry.get("is_duplicate", False),
            "original": entry.get("original"),
        })

    manifest = {
        # Top-level fields required by prepare-image-analysis.py
        "md_file": str(md_path),
        "images_dir": str(images_dir),
        "image_count": len(normalized_images),
        "generated": datetime.now().isoformat(),
        # Schema alignment: document_domain (populated by run-pipeline.py F14)
        "document_domain": "general",
        # Legacy fields preserved for backward compatibility
        "source_file": input_path.name,
        "source_format": "pptx",
        "total_images": len(normalized_images),
        "images": normalized_images,
    }
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)

    log.info(f"Manifest written: {manifest_path}")

    return {
        "md_path": md_path,
        "images_dir": images_dir,
        "manifest_path": manifest_path,
        "total_slides": total_slides,
        "total_images": len(manifest_images),
        "image_files": len(normalized_images),
        "charts": len([i for i in manifest_images if i.get("type") == "chart"]),
        "smartart": len([i for i in manifest_images if i.get("type") == "smartart"]),
        "chart_renders": len([i for i in manifest_images if i.get("type") == "chart_render"]),
        "decorative": len([i for i in manifest_images if i.get("decorative")]),
    }


def _shape_position(shape, slide_w, slide_h) -> dict:
    """Calculate shape position as percentage of slide dimensions."""
    try:
        left_pct = round((shape.left / slide_w) * 100, 1) if shape.left else 0
        top_pct = round((shape.top / slide_h) * 100, 1) if shape.top else 0
        w_pct = round((shape.width / slide_w) * 100, 1) if shape.width else 0
        h_pct = round((shape.height / slide_h) * 100, 1) if shape.height else 0
        return {
            "left": left_pct,
            "top": top_pct,
            "width": w_pct,
            "height": h_pct,
        }
    except Exception:
        return {"left": 0, "top": 0, "width": 0, "height": 0}


# ===========================================================================
# DOCX Table Repair (fallback for pandoc-dropped tables)
# ===========================================================================

def _docx_table_to_pipe(table) -> str:
    """Convert a python-docx Table object to a GFM pipe table string.

    Handles:
      - Multi-row tables with header row
      - Merged cells (uses the cell's text; python-docx repeats merged cell
        refs so we deduplicate by grid position)
      - Empty cells (preserved as empty)
      - Pipes in cell content (escaped as ``\\|``)
    """
    if not table.rows:
        return ""

    # Build a 2D grid of cell texts, deduplicating merged cells
    grid = []
    for row in table.rows:
        row_texts = []
        for cell in row.cells:
            text = cell.text.strip().replace("\n", " ").replace("|", "\\|")
            row_texts.append(text)
        grid.append(row_texts)

    if not grid:
        return ""

    # Determine the column count (max across all rows)
    n_cols = max(len(r) for r in grid)

    # Pad short rows
    for r in grid:
        while len(r) < n_cols:
            r.append("")

    # Build pipe table
    lines = []
    for i, row_texts in enumerate(grid):
        line = "| " + " | ".join(row_texts) + " |"
        lines.append(line)
        if i == 0:
            # Separator row after header
            sep = "| " + " | ".join(["---"] * n_cols) + " |"
            lines.append(sep)

    return "\n".join(lines)


def _is_empty_pipe_table(table_text: str) -> bool:
    """Return True if a pipe table block contains no real content.

    Detects tables where every data cell is empty, e.g.:
        | | |
        |---|---|
        | | |
    These are stub tables that pandoc emits when it fails to extract content.
    """
    for line in table_text.strip().split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        # Skip separator rows
        if re.match(r'^\|[\s\-:]+\|$', stripped):
            continue
        # Check if all cells are empty
        cells = stripped.split('|')[1:-1]  # drop leading/trailing empty
        if any(cell.strip() for cell in cells):
            return False
    return True


def _repair_dropped_tables(md_content: str, docx_path: str) -> str:
    """Repair tables that pandoc silently dropped or rendered as empty stubs.

    Strategy:
      1. Extract ALL tables from the DOCX via python-docx.
      2. Find ALL table blocks in pandoc output: both GFM pipe tables AND
         HTML <table>...</table> blocks.
      3. For each position (pipe or HTML block):
         - HTML block: advance docx_idx only (table is already rendered)
         - Pipe table, non-empty: advance docx_idx (table is fine)
         - Pipe table, empty/stub: replace with DOCX-extracted pipe table
      4. If DOCX has MORE tables than pandoc emitted positions, append the
         remaining DOCX tables at the end.

    This is a FALLBACK — when pandoc tables have content, they are left alone.
    """
    from docx import Document as DocxDocument

    try:
        doc = DocxDocument(str(docx_path))
    except Exception as e:
        log.warning(f"_repair_dropped_tables: cannot open DOCX: {e}")
        return md_content

    docx_tables = doc.tables
    if not docx_tables:
        return md_content

    # Convert all DOCX tables to pipe-table strings
    docx_pipe_tables = []
    for tbl in docx_tables:
        pipe = _docx_table_to_pipe(tbl)
        if pipe.strip():
            docx_pipe_tables.append(pipe)

    if not docx_pipe_tables:
        return md_content

    # Find ALL table positions in the markdown: both GFM pipe tables
    # AND HTML <table>...</table> blocks. Each entry is a dict with
    # kind ("pipe" or "html"), start, end, and is_empty flag.
    lines = md_content.split('\n')
    all_table_positions = []
    current_pipe_start = None
    in_html_table = False
    html_table_start = None

    for i, line in enumerate(lines):
        stripped = line.strip()

        # Detect HTML table blocks (opening and closing)
        if not in_html_table and re.match(r'<table\b', stripped, re.IGNORECASE):
            # Close any open pipe table block first
            if current_pipe_start is not None:
                block_text = '\n'.join(lines[current_pipe_start:i])
                all_table_positions.append({
                    "kind": "pipe",
                    "start": current_pipe_start,
                    "end": i,
                    "is_empty": _is_empty_pipe_table(block_text),
                })
                current_pipe_start = None
            in_html_table = True
            html_table_start = i
            continue

        if in_html_table:
            if re.match(r'</table>', stripped, re.IGNORECASE):
                all_table_positions.append({
                    "kind": "html",
                    "start": html_table_start,
                    "end": i + 1,
                    "is_empty": False,  # HTML tables are never empty stubs
                })
                in_html_table = False
                html_table_start = None
            continue

        # Detect GFM pipe table blocks
        if stripped.startswith('|') and stripped.endswith('|'):
            if current_pipe_start is None:
                current_pipe_start = i
        else:
            if current_pipe_start is not None:
                block_text = '\n'.join(lines[current_pipe_start:i])
                all_table_positions.append({
                    "kind": "pipe",
                    "start": current_pipe_start,
                    "end": i,
                    "is_empty": _is_empty_pipe_table(block_text),
                })
                current_pipe_start = None

    # Close any open blocks at end of file
    if current_pipe_start is not None:
        block_text = '\n'.join(lines[current_pipe_start:len(lines)])
        all_table_positions.append({
            "kind": "pipe",
            "start": current_pipe_start,
            "end": len(lines),
            "is_empty": _is_empty_pipe_table(block_text),
        })
    if in_html_table and html_table_start is not None:
        all_table_positions.append({
            "kind": "html",
            "start": html_table_start,
            "end": len(lines),
            "is_empty": False,
        })

    html_count = sum(1 for p in all_table_positions if p["kind"] == "html")
    pipe_count = sum(1 for p in all_table_positions if p["kind"] == "pipe")
    if html_count > 0:
        log.info(
            f"_repair_dropped_tables: detected {pipe_count} pipe table(s) "
            f"and {html_count} HTML table block(s) in pandoc output"
        )

    # Walk positions and build replacements.
    # Rule: every detected position (pipe or HTML) consumes one docx_idx slot.
    docx_idx = 0
    replacements = []  # list of (start, end, replacement_text)

    for pos in all_table_positions:
        if docx_idx >= len(docx_pipe_tables):
            break
        if pos["kind"] == "html":
            # Already rendered as HTML — advance index, do not replace
            docx_idx += 1
        elif pos["is_empty"]:
            # Empty pipe stub — replace with DOCX extraction
            replacements.append((pos["start"], pos["end"],
                                 docx_pipe_tables[docx_idx]))
            docx_idx += 1
        else:
            # Non-empty pipe table — fine as-is, advance index
            docx_idx += 1

    # Apply replacements in reverse order to preserve line indices
    for start, end, replacement in reversed(replacements):
        replacement_lines = replacement.split('\n')
        lines[start:end] = replacement_lines

    repaired = '\n'.join(lines)

    # Append only truly missing tables (those with no corresponding pandoc
    # position at all). This count is now accurate because HTML blocks
    # were tracked in the alignment loop above.
    remaining_count = len(docx_pipe_tables) - docx_idx
    if remaining_count > 0:
        log.info(
            f"_repair_dropped_tables: appending {remaining_count} table(s) "
            f"that pandoc omitted entirely"
        )
        extra_tables = []
        for tbl in docx_pipe_tables[docx_idx:]:
            extra_tables.append(f"\n\n{tbl}")
        repaired += ''.join(extra_tables) + '\n'

    replaced_count = len(replacements)
    if replaced_count > 0 or remaining_count > 0:
        log.info(
            f"_repair_dropped_tables: repaired {replaced_count} empty stub table(s), "
            f"appended {remaining_count} missing table(s)"
        )

    return repaired


# ===========================================================================
# DOCX Conversion
# ===========================================================================

def convert_docx(input_path: Path, output_dir: Path, skip_vision: bool) -> dict:
    """Convert DOCX to markdown using pandoc + image extraction."""
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    basename = input_path.stem
    images_dir = output_dir / f"{basename}_images"
    images_dir.mkdir(parents=True, exist_ok=True)
    md_path = output_dir / f"{basename}.md"
    manifest_path = output_dir / f"{basename}_manifest.json"

    log.info(f"Opening DOCX: {input_path.name}")

    # -----------------------------------------------------------------------
    # Step 1: Extract text via pandoc
    # -----------------------------------------------------------------------
    log.info("Extracting text via pandoc...")
    try:
        result = subprocess.run(
            ["pandoc", str(input_path), "-t", "gfm", "--wrap=none"],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode != 0:
            log.warning(f"pandoc text extraction warning: {result.stderr}")
        md_text = result.stdout
    except Exception as e:
        log.error(f"pandoc failed: {e}. Falling back to python-docx text extraction.")
        doc = Document(str(input_path))
        paragraphs = [clean_text(p.text) for p in doc.paragraphs if p.text.strip()]
        md_text = "\n\n".join(paragraphs)

    # -----------------------------------------------------------------------
    # Step 1b: Repair tables that pandoc dropped or rendered as empty stubs
    # -----------------------------------------------------------------------
    log.info("Checking for pandoc-dropped tables...")
    md_text = _repair_dropped_tables(md_text, str(input_path))

    # -----------------------------------------------------------------------
    # Step 2: Extract images via python-docx relationships
    # -----------------------------------------------------------------------
    log.info("Extracting images via python-docx...")
    manifest_images = []
    img_counter = 0
    # QC-4: Deduplication tracking for DOCX images
    docx_seen_hashes: dict[str, str] = {}  # hash -> first img_filename
    # Bug 2 fix: mapping from DOCX-internal media paths to extracted filenames
    docx_media_to_extracted: dict[str, str] = {}

    try:
        doc = Document(str(input_path))

        for rel_id, rel in doc.part.rels.items():
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    content_type = image_part.content_type
                    ext = ext_from_content_type(content_type)
                    blob = image_part.blob

                    img_counter += 1
                    img_id = f"img{img_counter:02d}"
                    img_filename = f"{img_id}.{ext}"
                    img_path = images_dir / img_filename

                    # QC-4: Deduplication — hash the blob before writing.
                    blob_hash = hashlib.sha256(blob).hexdigest()
                    is_duplicate = blob_hash in docx_seen_hashes
                    original_filename = docx_seen_hashes.get(blob_hash)
                    if not is_duplicate:
                        docx_seen_hashes[blob_hash] = img_filename

                    # Write image (always — preserves correct rel mapping)
                    with open(img_path, "wb") as f:
                        f.write(blob)

                    # Bug 2 fix: capture the DOCX-internal path for this
                    # relationship so we can rewrite pandoc media/ references.
                    # rel.target_ref is e.g. "media/image1.png" or
                    # "../media/image1.png" in some DOCX files.
                    _rel_target = getattr(rel, "target_ref", None)
                    if _rel_target is None:
                        try:
                            _rel_target = rel._element.get("Target", "")
                        except Exception:
                            _rel_target = ""
                    if _rel_target:
                        # Normalize: strip leading "./" or "../" prefixes
                        _rel_norm = _rel_target
                        while _rel_norm.startswith("../"):
                            _rel_norm = _rel_norm[3:]
                        while _rel_norm.startswith("./"):
                            _rel_norm = _rel_norm[2:]
                        # Ensure it starts with "media/" for consistent lookup
                        if not _rel_norm.startswith("media/"):
                            _rel_norm = "media/" + _rel_norm.split("/")[-1]
                        # Map to img_filename (will be updated below if WMF/TIFF conversion happens)
                        docx_media_to_extracted[_rel_norm] = img_filename

                    # Handle WMF/EMF
                    final_path = img_path
                    if ext in ("wmf", "emf"):
                        converted = convert_wmf_to_png(img_path, images_dir)
                        if converted:
                            final_path = converted
                            img_filename = final_path.name
                            # Update dedup tracking to the converted PNG filename
                            if not is_duplicate:
                                docx_seen_hashes[blob_hash] = img_filename
                            # Bug 2 fix: update media mapping to converted filename
                            if _rel_target and _rel_norm in docx_media_to_extracted:
                                docx_media_to_extracted[_rel_norm] = img_filename

                    # Handle TIFF→PNG conversion for Opus vision compatibility
                    if final_path.suffix.lower() in (".tiff", ".tif"):
                        try:
                            from PIL import Image as PilImage
                            pil_img = PilImage.open(final_path)
                            png_path = Path(str(final_path).rsplit(".", 1)[0] + ".png")
                            pil_img.save(png_path, "PNG")
                            os.remove(final_path)
                            final_path = png_path
                            img_filename = final_path.name
                            # Update dedup tracking to the PNG filename
                            if not is_duplicate:
                                docx_seen_hashes[blob_hash] = img_filename
                            # Bug 2 fix: update media mapping to converted filename
                            if _rel_target and _rel_norm in docx_media_to_extracted:
                                docx_media_to_extracted[_rel_norm] = img_filename
                            log.info(f"  TIFF converted to PNG: {img_filename}")
                        except Exception as e:
                            log.warning(f"  TIFF→PNG conversion failed for {final_path.name}: {e}")

                    # Fix 1.8: resize oversized images before analysis
                    _ensure_max_dimension(final_path)

                    size_bytes = final_path.stat().st_size
                    decorative = size_bytes < 5000
                    dims = get_image_dimensions(final_path)

                    # QC-3: Blank image detection — catches WMF→PNG failures.
                    blank = False
                    if not decorative and final_path.suffix.lower() in (".png", ".jpg", ".jpeg", ".bmp", ".gif"):
                        blank = is_blank_image(final_path)
                        if blank:
                            log.warning(
                                f"WARNING: Blank image detected (likely WMF conversion failure): "
                                f"{final_path}"
                            )

                    entry = {
                        "id": img_id,
                        "file": img_filename,
                        "slide": None,
                        "content_type": content_type,
                        "size_bytes": size_bytes,
                        "dimensions": dims,
                        "position": None,
                        "decorative": decorative,
                        "source_shape": rel_id,
                        "context": {
                            "slide_title": None,
                            "nearby_text": "",
                        },
                        # QC-3: blank detection flag
                        "blank": blank,
                        # QC-4: deduplication flags
                        "is_duplicate": is_duplicate,
                        "original": original_filename,
                    }
                    manifest_images.append(entry)

                    if is_duplicate:
                        log.info(
                            f"  {img_filename} ({size_bytes}B) - duplicate of {original_filename}"
                        )
                    elif decorative:
                        log.info(f"  {img_filename} ({size_bytes}B) - decorative")
                    else:
                        log.info(f"  {img_filename} ({size_bytes}B)")

                except Exception as e:
                    log.warning(f"  Failed to extract image from rel '{rel_id}': {e}")
    except Exception as e:
        log.warning(f"python-docx image extraction failed: {e}")

    # -----------------------------------------------------------------------
    # Step 2b: Rewrite pandoc media/imageN paths to extracted imgNN paths
    # -----------------------------------------------------------------------
    # Pandoc preserves DOCX-internal paths like <img src="media/image1.png">.
    # These do not exist on disk. Replace them with the actual extracted paths
    # using the {basename}_images/imgNN.ext convention.
    if docx_media_to_extracted:
        _rewrite_count = 0

        def _replace_html_img_src(m):
            """Replace media/imageN.ext in <img src="..."> with extracted path."""
            nonlocal _rewrite_count
            full_tag = m.group(0)
            src_val = m.group(1)
            # Normalize the src value for lookup
            src_norm = src_val
            while src_norm.startswith("../"):
                src_norm = src_norm[3:]
            while src_norm.startswith("./"):
                src_norm = src_norm[2:]
            if not src_norm.startswith("media/"):
                src_norm = "media/" + src_norm.split("/")[-1]
            extracted_name = docx_media_to_extracted.get(src_norm)
            if extracted_name:
                new_src = f"{basename}_images/{extracted_name}"
                _rewrite_count += 1
                return full_tag.replace(src_val, new_src)
            return full_tag

        def _replace_md_img_src(m):
            """Replace media/imageN.ext in ![alt](path) with extracted path."""
            nonlocal _rewrite_count
            alt = m.group(1)
            path = m.group(2)
            # Normalize the path for lookup
            path_norm = path
            while path_norm.startswith("../"):
                path_norm = path_norm[3:]
            while path_norm.startswith("./"):
                path_norm = path_norm[2:]
            if not path_norm.startswith("media/"):
                path_norm = "media/" + path_norm.split("/")[-1]
            extracted_name = docx_media_to_extracted.get(path_norm)
            if extracted_name:
                new_path = f"{basename}_images/{extracted_name}"
                _rewrite_count += 1
                return f"![{alt}]({new_path})"
            return m.group(0)

        # Pattern 1: HTML img tags — <img src="media/imageN.ext" ...>
        md_text = re.sub(
            r'<img\b[^>]*\bsrc="([^"]+)"[^>]*/?>',
            _replace_html_img_src,
            md_text,
        )
        # Pattern 2: Markdown image syntax — ![alt](media/imageN.ext)
        md_text = re.sub(
            r'!\[([^\]]*)\]\(([^)]+)\)',
            _replace_md_img_src,
            md_text,
        )
        if _rewrite_count > 0:
            log.info(
                f"Rewrote {_rewrite_count} media/ image path(s) "
                f"in pandoc markdown output "
                f"(mapping: {len(docx_media_to_extracted)} entries)"
            )

    # -----------------------------------------------------------------------
    # F11: Strip bold markers (**) from heading lines only.
    # Heading text should be plain; ** from multi-run DOCX formatting is noise.
    # Applied per-line so body text bold formatting is preserved.
    # -----------------------------------------------------------------------
    def _strip_bold_from_heading(m):
        return re.sub(r'\*\*', '', m.group(0))

    md_text = re.sub(r'^#{1,4} .+', _strip_bold_from_heading, md_text, flags=re.MULTILINE)

    # -----------------------------------------------------------------------
    # RC11: Extract inline images from heading lines.
    # Pandoc sometimes embeds <img> tags inside heading lines, e.g.:
    #   #### <img src="...img36.png" .../> Generic or disease-specific
    # This breaks heading semantics. Fix: move the img tag to its own line
    # BEFORE the heading, and strip it from the heading text.
    # Handles HTML <img> tags embedded in heading lines.
    # -----------------------------------------------------------------------
    _HEADING_WITH_IMG_RE = re.compile(
        r'^(#{1,6}\s+)'                   # heading prefix (group 1)
        r'(.*?)'                           # optional text before img (group 2)
        r'(<img\s[^>]*?/?>)'              # HTML img tag (group 3)
        r'(.*?)$',                         # text after img (group 4)
        re.MULTILINE,
    )

    def _extract_img_from_heading(m):
        prefix = m.group(1)      # "#### "
        before = m.group(2).strip()
        img_tag = m.group(3).strip()
        after = m.group(4).strip()
        heading_text = f'{before} {after}'.strip()
        if heading_text:
            return f'{img_tag}\n\n{prefix}{heading_text}'
        else:
            # Heading was ONLY the image — keep image, drop empty heading
            return img_tag

    md_text = _HEADING_WITH_IMG_RE.sub(_extract_img_from_heading, md_text)

    # -----------------------------------------------------------------------
    # F13: Associate orphaned figure captions with preceding images.
    # Pandoc emits Caption-styled paragraphs as plain text. When a line that
    # starts with "Figure" or "Fig." immediately follows an <img> or ![...]()
    # line (with at most 1 blank line between them), wrap it in <figcaption>
    # tags so downstream consumers know it belongs to the image above.
    #
    # RC10 enhancement: expanded to handle more caption patterns:
    #   - Bold-wrapped captions like "**Figure 1:**"
    #   - Captions with leading whitespace
    #   - Forward-looking: also scan from each image line forward
    #     (catches cases where backward scan fails due to intervening content)
    #   - Skip lines already wrapped in <figcaption>
    # -----------------------------------------------------------------------
    _CAPTION_RE = re.compile(
        r'^\s*(?:\*\*)?'           # optional leading whitespace + optional bold
        r'(?:Figure|Fig\.)\s+\d+'  # "Figure N" or "Fig. N"
        r'(?:\*\*)?',              # optional closing bold
    )

    def _associate_captions(md_text_in):
        lines = md_text_in.split('\n')
        result = list(lines)  # work on a copy

        def _is_image_line(ln):
            s = ln.strip()
            return '<img' in s or bool(re.match(r'!\[', s))

        # Pass 1: backward search from caption lines (original logic, relaxed)
        for i, line in enumerate(result):
            stripped = line.strip()
            if not _CAPTION_RE.match(stripped):
                continue
            # Already wrapped — skip
            if '<figcaption>' in stripped:
                continue
            # Search backward through up to 3 preceding lines
            # (allows up to 2 blank lines between image and caption)
            for j in range(i - 1, max(i - 4, -1), -1):
                prev = result[j]
                if _is_image_line(prev):
                    result[i] = f'<figcaption>{stripped}</figcaption>'
                    break
                elif prev.strip():
                    # Non-empty, non-image line — stop
                    break

        # Pass 2: forward search from image lines
        # Catches captions that Pass 1 missed (e.g. image followed by
        # a non-blank non-image line like <!-- comment -->, then caption)
        for i, line in enumerate(result):
            if not _is_image_line(line):
                continue
            # Look forward up to 3 lines for a caption
            for j in range(i + 1, min(i + 4, len(result))):
                fwd = result[j]
                fwd_stripped = fwd.strip()
                if not fwd_stripped:
                    continue  # skip blank lines
                if '<figcaption>' in fwd_stripped:
                    break  # already wrapped — done
                if _CAPTION_RE.match(fwd_stripped):
                    result[j] = f'<figcaption>{fwd_stripped}</figcaption>'
                    break
                else:
                    break  # non-blank, non-caption — stop

        return '\n'.join(result)

    md_text = _associate_captions(md_text)

    # -----------------------------------------------------------------------
    # Step 3: Build final markdown with frontmatter
    # -----------------------------------------------------------------------
    # QC-5: Title extraction — prefer core_properties.title (stripped),
    # fall back to filename-derived title. Never use author name as title.
    doc_title = basename.replace("_", " ").replace("-", " ")

    try:
        # doc may already be open from image extraction above; open fresh to be safe
        doc_for_title = Document(str(input_path))
        raw_title = (doc_for_title.core_properties.title or "").strip()
        if raw_title:
            doc_title = raw_title
    except Exception:
        pass

    total_images_with_files = len([i for i in manifest_images if i.get("file")])

    frontmatter = f"""---
title: "{doc_title}"
source_file: "{input_path.name}"
source_format: "docx"
conversion_tool: "convert-office.py v{VERSION}"
conversion_date: "{datetime.now().strftime('%Y-%m-%dT%H:%M:%S')}"
document_type: "document"
fidelity_standard: "text_content"
total_images: {total_images_with_files}
images_directory: "{basename}_images/"
pipeline_version: "{PIPELINE_VERSION}"
image_notes: {"pending" if total_images_with_files > 0 else "none"}
---
"""

    # Append image references at the end of the markdown
    image_section = ""
    if manifest_images:
        image_section = "\n\n---\n\n## Extracted Images\n\n"
        for img in manifest_images:
            if img.get("file"):
                # Issue 6 fix: use subdirectory-relative path instead of bare filename
                img_rel_path = f"{basename}_images/{img['file']}"
                # MINOR-2: Standardized to PPTX format for consistency
                dec = " | Decorative: yes" if img.get("decorative") else ""
                # Fix 3.6/M18: Propagate duplicate flag to IMAGE comment
                dup = " | Duplicate: yes" if img.get("is_duplicate") else ""
                if img.get("is_duplicate") and not dec:
                    dec = " | Decorative: yes"
                image_section += (
                    f"<!-- IMAGE: {img_rel_path} | "
                    f"Size: {img['dimensions'][0]}x{img['dimensions'][1]}"
                    f"{dec}{dup} -->\n"
                    f"![Image {img['id']}]({img_rel_path})\n\n"
                )

    with open(md_path, "w", encoding="utf-8") as f:
        f.write(frontmatter)
        f.write("\n")
        f.write(md_text)
        f.write(image_section)

    log.info(f"Markdown written: {md_path}")

    # -----------------------------------------------------------------------
    # Step 4: Write manifest
    # Normalized to match convert-paper.py format so that prepare-image-analysis.py
    # (Step 3 in run-pipeline.py) can consume it without format-specific branches.
    # Required fields per image: filename, figure_num, width, height, page.
    # Optional but expected: type_guess, section_context, nearby_text,
    #                        detected_caption, file_path.
    # -----------------------------------------------------------------------
    normalized_images = []
    figure_counter = 0
    for entry in manifest_images:
        if not entry.get("file"):
            continue
        figure_counter += 1
        dims = entry.get("dimensions", [0, 0])
        w = dims[0] if isinstance(dims, list) and len(dims) >= 2 else 0
        h = dims[1] if isinstance(dims, list) and len(dims) >= 2 else 0
        img_filename = entry["file"]
        img_abs_path = str(images_dir / img_filename)
        normalized_images.append({
            # Fields required by prepare-image-analysis.py (no .get() fallback)
            "figure_num": figure_counter,
            "filename": img_filename,
            "width": w,
            "height": h,
            # Optional fields used by prepare-image-analysis.py
            "page": None,  # DOCX has no reliable page numbers
            "type_guess": None,
            "section_context": None,
            "nearby_text": None,
            "detected_caption": None,
            "file_path": img_abs_path,
            # Original fields preserved for any other downstream consumers
            "id": entry.get("id"),
            "content_type": entry.get("content_type"),
            "size_bytes": entry.get("size_bytes", 0),
            "dimensions": dims,
            "position": entry.get("position"),
            "decorative": entry.get("decorative", False),
            "source_shape": entry.get("source_shape"),
            "analysis_status": "pending",
            "description": f"Figure {figure_counter} from document",
            "source_format": "docx_embedded",
            # QC-3: blank image detection (WMF conversion failure indicator)
            "blank": entry.get("blank", False),
            # QC-4: deduplication fields
            "is_duplicate": entry.get("is_duplicate", False),
            "original": entry.get("original"),
        })

    manifest = {
        # Top-level fields required by prepare-image-analysis.py
        "md_file": str(md_path),
        "images_dir": str(images_dir),
        "image_count": len(normalized_images),
        "generated": datetime.now().isoformat(),
        # Schema alignment: document_domain (populated by run-pipeline.py F14)
        "document_domain": "general",
        # Legacy fields preserved for backward compatibility
        "source_file": input_path.name,
        "source_format": "docx",
        "total_images": len(normalized_images),
        "images": normalized_images,
    }
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)

    log.info(f"Manifest written: {manifest_path}")

    return {
        "md_path": md_path,
        "images_dir": images_dir,
        "manifest_path": manifest_path,
        "total_slides": None,
        "total_images": len(manifest_images),
        "image_files": len(normalized_images),
        "charts": 0,
        "smartart": 0,
        "chart_renders": 0,
        "decorative": len([i for i in manifest_images if i.get("decorative")]),
    }


# ===========================================================================
# TXT Conversion
# ===========================================================================

def convert_txt(input_path: Path, output_dir: Path, skip_vision: bool) -> dict:
    """Convert TXT to markdown with YAML frontmatter."""
    basename = input_path.stem
    md_path = output_dir / f"{basename}.md"

    log.info(f"Opening TXT: {input_path.name}")

    # Read content
    try:
        with open(input_path, "r", encoding="utf-8") as f:
            content = f.read()
    except UnicodeDecodeError:
        with open(input_path, "r", encoding="latin-1") as f:
            content = f.read()

    doc_title = basename.replace("_", " ").replace("-", " ")

    frontmatter = f"""---
title: "{doc_title}"
source_file: "{input_path.name}"
source_format: "txt"
conversion_tool: "convert-office.py v{VERSION}"
conversion_date: "{datetime.now().strftime('%Y-%m-%dT%H:%M:%S')}"
document_type: "text"
fidelity_standard: "text_content"
total_images: 0
pipeline_version: "{PIPELINE_VERSION}"
image_notes: none
---
"""

    with open(md_path, "w", encoding="utf-8") as f:
        f.write(frontmatter)
        f.write("\n")
        f.write(f"# {doc_title}\n\n")
        f.write(content)

    log.info(f"Markdown written: {md_path}")

    return {
        "md_path": md_path,
        "images_dir": None,
        "manifest_path": None,
        "total_slides": None,
        "total_images": 0,
        "image_files": 0,
        "charts": 0,
        "smartart": 0,
        "chart_renders": 0,
        "decorative": 0,
    }


# ===========================================================================
# Main
# ===========================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Convert PPTX/DOCX/TXT to Markdown with image extraction.",
    )
    parser.add_argument("input_file", help="Path to input file (.pptx, .docx, .txt)")
    parser.add_argument(
        "--output-dir",
        help="Output directory (default: same as input file)",
    )
    parser.add_argument(
        "--skip-vision",
        action="store_true",
        help="Skip vision placeholders (extract text + images only)",
    )
    args = parser.parse_args()

    input_path = Path(args.input_file).resolve()
    if not input_path.exists():
        log.error(f"Input file not found: {input_path}")
        sys.exit(1)

    # Determine output directory
    if args.output_dir:
        output_dir = Path(args.output_dir).resolve()
    else:
        output_dir = input_path.parent
    output_dir.mkdir(parents=True, exist_ok=True)

    ext = input_path.suffix.lower()
    log.info(f"Input: {input_path}")
    log.info(f"Output dir: {output_dir}")
    log.info(f"Format: {ext}")

    # Dispatch to appropriate converter
    if ext == ".pptx":
        result = convert_pptx(input_path, output_dir, args.skip_vision)
    elif ext == ".docx":
        result = convert_docx(input_path, output_dir, args.skip_vision)
    elif ext == ".txt":
        result = convert_txt(input_path, output_dir, args.skip_vision)
    elif ext == ".xlsx":
        log.error(
            "XLSX is not supported by convert-office.py. "
            "Use openpyxl to read Excel files directly, or export to CSV/TXT first. "
            "XLSX->MD conversion is lossy; openpyxl preserves formulas and structure."
        )
        sys.exit(1)
    else:
        log.error(f"Unsupported format: {ext}. Supported: .pptx, .docx, .txt")
        sys.exit(1)

    # R19: Generate image index for PPTX and DOCX
    # Pass manifest_path so generate_image_index can use manifest data
    # as primary source (Root Cause B fix: manifest-index reconciliation)
    image_index_meta = None
    if ext in (".pptx", ".docx"):
        try:
            m_path = result.get("manifest_path")
            image_index_meta = generate_image_index(
                input_path, output_dir, ext.lstrip("."),
                manifest_path=Path(m_path) if m_path else None,
            )
        except Exception as e:
            log.warning(f"Image index generation failed (non-fatal): {e}")

    # Update registry (R21: pass image index metadata)
    update_registry(
        input_path, result["md_path"],
        f"convert-office-{ext.lstrip('.')}",
        image_index_meta=image_index_meta,
    )

    # Print summary
    print()
    print("=" * 50)
    print("CONVERSION COMPLETE")
    print("=" * 50)
    print(f"  Markdown:   {result['md_path']}")
    if result.get("images_dir"):
        print(f"  Images dir: {result['images_dir']}")
    if result.get("manifest_path"):
        print(f"  Manifest:   {result['manifest_path']}")
    if result.get("total_slides") is not None:
        print(f"  Slides:     {result['total_slides']}")
    print(f"  Images:     {result['image_files']} extracted")
    if result.get("charts"):
        print(f"  Charts:     {result['charts']} (detected, not directly extractable)")
    if result.get("smartart"):
        print(f"  SmartArt:   {result['smartart']} (detected, not directly extractable)")
    if result.get("chart_renders"):
        print(f"  Rendered:   {result['chart_renders']} (via LibreOffice slide render)")
    if result.get("decorative"):
        print(f"  Decorative: {result['decorative']} (< 5KB)")
    if image_index_meta:
        print(f"  Image index: {image_index_meta.get('image_index_path', 'N/A')}")
        print(f"    Pages with images: {image_index_meta.get('pages_with_images', 0)}")
        print(f"    Substantive: {image_index_meta.get('substantive_images', 0)}")
        print(f"    Testable: {image_index_meta.get('has_testable_images', False)}")
    print("=" * 50)


if __name__ == "__main__":
    main()
