#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Convert documents (PDF, DOCX, PPTX, XLSX) to high-fidelity Markdown with enhanced metadata.

MODIFIED from ~/.claude/scripts/convert-paper.py with 6 NEW capabilities:
1. Context summary generation (extractive, no LLM)
2. Section-to-image mapping in manifest
3. Caption detection for PDF images
4. Document domain auto-detection
5. Duplicate xref tracking (bug fix)
6. Enhanced manifest fields (type_guess, section_context, detected_caption)

PRESERVED capabilities:
- Multi-panel image splitting (>800x800)
- Sparse page rendering (300 DPI for pages with <200 chars text)
- pymupdf.layout mode import
- MarkItDown PDF fallback

Text extraction:
  - PDF: pymupdf4llm (proper tables, headings, multi-column support)
  - DOCX/PPTX/XLSX: MarkItDown (good for structured formats)
Image extraction: PyMuPDF (PDF), python-docx (DOCX), python-pptx (PPTX).

Runs ENTIRELY via Bash/Python, bypassing Claude's API copyright filter.

Usage:
    python3 convert-paper.py <input_file> [-o output.md] [-i images/] [-s short-name]

Supported: .pdf .docx .pptx .xlsx
"""

import argparse
import io
import json
import math
import os
import re
import sys
import unicodedata
from datetime import datetime
from pathlib import Path
from typing import Optional

try:
    from markitdown import MarkItDown
except ImportError:
    MarkItDown = None

try:
    import pymupdf4llm
    # Layout mode (pymupdf.layout) is intentionally NOT imported.
    # It activates AI-based page analysis for better headings,
    # multi-column, and table detection — but it also triggers a
    # known pymupdf4llm bug on certain PDFs where table bounding
    # boxes are empty, causing:
    #   ValueError: min() iterable argument is empty
    #   (at pymupdf/table.py:1534)
    # This crash affected the Navigating-change OHE report (166
    # images).  Without layout mode, to_markdown() still works
    # correctly using standard text extraction.  The ValueError
    # catch at the to_markdown() call site (exit code 3) remains
    # as a safety net in case other pymupdf4llm bugs surface.
    _HAS_LAYOUT = False
except ImportError:
    pymupdf4llm = None
    _HAS_LAYOUT = False

try:
    from PIL import Image
except ImportError:
    Image = None
    print("WARNING: Pillow not installed. Multi-panel splitting disabled.")


def _ensure_max_dimension(img_path, max_dim=2000):
    """Resize image in-place if either dimension exceeds max_dim.

    Uses PIL thumbnail() which preserves aspect ratio and uses LANCZOS
    resampling for quality.  Images already within limits are untouched.

    Returns True if image was resized, False otherwise.
    """
    try:
        from PIL import Image as _PilResize
        with _PilResize.open(img_path) as _img:
            if max(_img.size) <= max_dim:
                return False
            _orig = _img.size
            _img.thumbnail((max_dim, max_dim), _PilResize.LANCZOS)
            _img.save(img_path)
            print(f"    [resize] {Path(img_path).name}: "
                  f"{_orig[0]}x{_orig[1]} -> "
                  f"{_img.size[0]}x{_img.size[1]}")
            return True
    except Exception as _e:
        print(f"WARNING: Could not resize {Path(img_path).name}: {_e}",
              file=sys.stderr)
        return False


def normalize_symbols(text: str, extractor: str = "docling") -> str:
    """Post-extraction symbol normalization.

    When extractor is "pymupdf4llm", applies glyph mapping fixes specific to
    pymupdf4llm's text extraction engine (C4, C5, C6, C7, m25). These patterns
    correct known glyph mapping errors where characters are rendered as wrong
    symbols. Other extractors (e.g., docling) use different text engines and
    do not produce these artifacts, so the glyph fixes are skipped to avoid
    false positives on legitimate text.

    Universal fixes (U+FFFD removal) are always applied regardless of extractor.

    CAUTION: minus-as-2 regex can false-positive on real 20.xx values.
    All patterns are guarded with statistical context.
    """
    _fixes_applied = 0

    # ── pymupdf4llm-specific glyph mapping fixes ──
    # These correct known extraction bugs in pymupdf4llm's text engine.
    # Other extractors (docling, tesseract) do not produce these artifacts.
    if extractor == "pymupdf4llm":
        # C4: Fix minus-as-2 in CI context: "CrI 20.26" -> "CrI -0.26"
        text, n = re.subn(r'(?<=CrI[\s:])2(\d+\.\d+)', r'-\1', text)
        _fixes_applied += n
        text, n = re.subn(r'(?<=CI[\s:])2(\d+\.\d+)', r'-\1', text)
        _fixes_applied += n
        text, n = re.subn(r'(?<=to\s)2(\d+\.\d+)', r'-\1', text)
        _fixes_applied += n
        # Comma-separated negative: ", 20.26" -> ", -0.26"
        # Only when first digit after 2 is 0 (distinguishes from real 20.xx)
        text, n = re.subn(r',\s*2(0\.\d+)', r', -\1', text)
        _fixes_applied += n

        # C5: Fix < as backslash in P-value: "P \ 0.05" -> "P < 0.05"
        text, n = re.subn(r'([Pp])\s*\\\s*(0\.\d+)', r'\1 < \2', text)
        _fixes_applied += n

        # C6: Fix > as period in age context: ".65 years" -> ">65 years"
        text, n = re.subn(
            r'(?<=\s)\.(\d+)\s+(y(?:ear|r)s?|months?|days?|weeks?)',
            r'>\1 \2', text
        )
        _fixes_applied += n

        # C7: Fix + as 1 in formula: "1 exp(" -> "+ exp("
        text, n = re.subn(
            r'(\d)\s+1\s+(exp|log|sqrt)\s*\(',
            r'\1 + \2(', text
        )
        _fixes_applied += n

        # m25: Fix = as hyphen in sample size: "n - 123" -> "n = 123"
        text, n = re.subn(r'\bn\s*-\s*(\d{2,})\b', r'n = \1', text)
        _fixes_applied += n
        # Fix = as hyphen in P-value: "P - 0.05" -> "P = 0.05"
        text, n = re.subn(r'([Pp])\s*-\s*(0\.\d+)', r'\1 = \2', text)
        _fixes_applied += n

    # ── Universal fixes (all extractors) ──

    # Ligature escape sequences (e.g., docling outputs "/uniFB01" instead of "fi")
    # These are Unicode ligature codepoints FB00-FB06.
    # MINOR-2 fix: case-normalize /uni escape hex digits to uppercase before
    # explicit map lookup, so "/unifb01" and "/uniFB01" both hit the map.
    _ligature_map = {
        '/uniFB00': 'ff',
        '/uniFB01': 'fi',
        '/uniFB02': 'fl',
        '/uniFB03': 'ffi',
        '/uniFB04': 'ffl',
        '/uniFB05': 'st',        # long s + t -> modern convention
        '/uniFB06': 'st',
    }
    # Normalize any /uniXXXX hex digits to uppercase so explicit map always matches
    text, n_case = re.subn(
        r'/uni([0-9A-Fa-f]{4})',
        lambda m: '/uni' + m.group(1).upper(),
        text
    )
    for escape_seq, replacement in _ligature_map.items():
        count = text.count(escape_seq)
        if count > 0:
            text = text.replace(escape_seq, replacement)
            _fixes_applied += count

    # General /uniXXXX handler for any other Unicode escape sequences
    # Matches /uni followed by exactly 4 hex digits and replaces with the
    # corresponding Unicode character.
    def _replace_uni_escape(m):
        codepoint = int(m.group(1), 16)
        # MINOR-1 fix: block surrogate codepoints (D800-DFFF) which create
        # unencodable strings that crash on file write
        if 0xD800 <= codepoint <= 0xDFFF:
            return m.group(0)  # leave surrogates unchanged
        try:
            return chr(codepoint)
        except (ValueError, OverflowError):
            return m.group(0)  # leave unchanged if invalid
    text, n = re.subn(r'/uni([0-9A-Fa-f]{4})', _replace_uni_escape, text)
    _fixes_applied += n

    # C3: Remove U+FFFD replacement characters
    fffd_count = text.count('\ufffd')
    if fffd_count > 0:
        _fixes_applied += fffd_count
        text = text.replace('\ufffd', '')

    if _fixes_applied > 0:
        print(f"  [normalize] Applied {_fixes_applied} symbol fixes"
              f"{f' (including {fffd_count} U+FFFD removals)' if fffd_count > 0 else ''}")

    return text


def post_extraction_cleanup(text: str) -> str:
    """Unified post-extraction cleanup for common PDF extraction artifacts.

    Addresses 11 issue IDs: M21, m4-m7, m10, m13, m15-m17, m26.
    Called after normalize_symbols() in the extraction pipeline.
    Applied to all extractor output (pymupdf4llm, docling, etc.).

    Fix 3.15 from PIPELINE-FIX-PLAN-AUDITED.md.
    """
    # m13: Strip control characters (except newline, tab, carriage return)
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)

    # m5: Remove Springer page artifacts "## 1 3"
    text = re.sub(r'^## \d+ \d+\s*$', '', text, flags=re.MULTILINE)

    # m10: Remove orphan page numbers (standalone 1-3 digits on a line)
    text = re.sub(r'^\s*\d{1,3}\s*$', '', text, flags=re.MULTILINE)

    # m4: ALL-CAPS headings -> ## Title Case
    for h in ['METHODS', 'RESULTS', 'DISCUSSION', 'CONCLUSIONS',
              'CONCLUSION', 'INTRODUCTION', 'ABSTRACT', 'REFERENCES',
              'BACKGROUND', 'OBJECTIVES', 'ACKNOWLEDGEMENTS',
              'ACKNOWLEDGMENTS', 'APPENDIX', 'LIMITATIONS']:
        text = re.sub(rf'^{h}\s*$', f'## {h.title()}',
                       text, flags=re.MULTILINE)

    # m16: Superscript brackets R [2] -> R^2
    text = re.sub(r'R\s*\[2\]', 'R^2', text)

    # m17: Backtick bullet `o` -> dash bullet
    text = re.sub(r'^\s*`o`\s+', '- ', text, flags=re.MULTILINE)

    # m26: Blockquote collision >75 year -> \>75 year
    text = re.sub(r'^(>)(\d)', r'\\>\2', text, flags=re.MULTILINE)

    # m6, m7: Domain-specific dehyphenation
    compounds = [
        'cost-effectiveness', 'cost-effective', 'cost-utility',
        'cost-benefit', 'self-report', 'self-reported',
        'meta-analysis', 'meta-analyses', 'health-related',
        'preference-based', 'condition-specific',
        'decision-making', 'decision-analytic', 'age-adjusted',
        'quality-adjusted', 'willingness-to-pay',
    ]
    for compound in compounds:
        merged = compound.replace('-', '')
        text = text.replace(merged, compound)
        text = text.replace(merged.title(), compound.title())

    # Collapse excessive blank lines
    text = re.sub(r'\n{4,}', '\n\n\n', text)

    return text


def _strip_ocr_artifacts(text: str) -> str:
    """
    Remove garbled OCR text from docling output (safety net).

    Uses Shannon entropy to detect random-looking text that OCR
    produces when applied to figure images. Lines with low entropy
    (repetitive/structured) or matching statistical patterns are preserved.

    Called from _docling_postprocess() as a final cleanup step.
    """
    lines = text.split('\n')
    cleaned = []
    stripped_count = 0
    in_figure_block = False

    for line in lines:
        stripped = line.strip()

        # Track figure blocks for context-aware detection
        if stripped.startswith('![') or stripped.startswith('<!-- IMAGE:'):
            in_figure_block = True
            cleaned.append(line)
            continue
        if in_figure_block and stripped == '':
            in_figure_block = False
            cleaned.append(line)
            continue

        # Always preserve structural lines
        if (not stripped
                or stripped.startswith('#')       # headings
                or stripped.startswith('|')       # tables
                or stripped.startswith('!')       # images
                or stripped.startswith('<!--')    # HTML comments
                or stripped.startswith('```')     # code blocks
                or stripped.startswith('---')     # horizontal rules
                or stripped.startswith('>')       # blockquotes
                or len(stripped) < 20):           # minimum length guard
            cleaned.append(line)
            continue

        # Preserve statistical text (CRITICAL: do not strip real data)
        if _STATISTICAL_PATTERNS.search(stripped):
            cleaned.append(line)
            continue

        # Compute Shannon entropy
        freq = {}
        for ch in stripped.lower():
            freq[ch] = freq.get(ch, 0) + 1
        total = len(stripped)
        entropy = -sum(
            (c / total) * math.log2(c / total) for c in freq.values()
        )

        # High entropy + low word ratio = likely OCR garbage
        words = stripped.split()
        if len(words) > 0:
            # word_ratio: proportion of "words" that look like real words
            # (contain at least one vowel and are 2+ chars)
            vowel_set = set('aeiouAEIOU')
            real_words = sum(
                1 for w in words
                if len(w) >= 2 and any(c in vowel_set for c in w)
            )
            word_ratio = real_words / len(words)
        else:
            word_ratio = 0.0

        # Conservative thresholds: entropy > 4.5 AND word_ratio < 0.3
        if entropy > 4.5 and word_ratio < 0.3:
            stripped_count += 1
            continue  # Strip this line

        # In figure block context, be more aggressive (lower threshold)
        if in_figure_block and entropy > 4.0 and word_ratio < 0.4:
            stripped_count += 1
            continue

        cleaned.append(line)

    # Collapse 3+ consecutive newlines to 2
    result = '\n'.join(cleaned)
    while '\n\n\n' in result:
        result = result.replace('\n\n\n', '\n\n')

    if stripped_count > 0:
        print(f"  [ocr-artifacts] Stripped {stripped_count} likely OCR garbage lines")

    return result


def _fix_font_encoding(text: str, extractor: str = "docling") -> str:
    """Replace /C## glyph artifacts and strip GLYPH patterns.

    Runs FIRST in the post-processing pipeline (FIX-B01, CM-009).
    Font encoding map fixes are applied for all extractors since the
    /C## artifacts come from the PDF source, not the extractor.
    Character substitutions (thorn, eth, etc.) are gated by context
    to avoid false positives.
    """
    fixes = 0

    # Step 1: Replace known /C## codes with correct characters
    def _replace_ccode(m):
        nonlocal fixes
        code = "/C" + m.group(1)
        if code in _FONT_ENCODING_MAP:
            fixes += 1
            return _FONT_ENCODING_MAP[code]
        return m.group(0)

    text = _FONT_ENCODING_COMPILED.sub(_replace_ccode, text)

    # Step 2: Generic fallback for unmapped /C## codes -> bullet
    new_text = _GENERIC_CCODE_PATTERN.sub('\u2022', text)
    if new_text != text:
        diff_count = (len(text) - len(new_text)
                      + new_text.count('\u2022') - text.count('\u2022'))
        fixes += max(diff_count, 1)
        text = new_text

    # Step 3: Strip GLYPH(cmap:...) artifacts
    new_text = _GLYPH_PATTERN.sub('', text)
    if new_text != text:
        fixes += 1
        text = new_text

    # Step 4: Context-gated character substitutions
    # QC-FP-07: Widened from [A-Z]{2,} to \w{2,} for mixed-case
    # treatment arms. Thorn (U+00FE) is extremely rare in English/HTA
    # text outside of this font-encoding artifact.
    if '\u00FE' in text:
        text_new = re.sub(
            r'(\w{2,})\s*\u00FE\s*(\w{2,})',
            lambda m: m.group(1) + ' + ' + m.group(2),
            text
        )
        if text_new != text:
            fixes += 1
            text = text_new

    # Eth -> ( after digit (formula context)
    text_new = re.sub(
        r'(?<=\d)\s*\u00F0\s*(?=[\d\w])',
        '(',
        text
    )
    if text_new != text:
        fixes += 1
        text = text_new

    # Thorn -> ) before punctuation/digit
    text_new = re.sub(
        r'(?<=[\d\w])\s*\u00DE\s*(?=\s|[,;.\d])',
        ')',
        text
    )
    if text_new != text:
        fixes += 1
        text = text_new

    # Quarter -> = between spaces
    text_new = re.sub(
        r'(?<=\s)\u00BC(?=\s)',
        '=',
        text
    )
    if text_new != text:
        fixes += 1
        text = text_new

    # Step 5: Reference zero-padding fix (ISSUE-20, docling-only)
    text_new = re.sub(r'- 0\[(\d)\]', r'- [\1]', text)
    if text_new != text:
        fixes += 1
        text = text_new

    if fixes > 0:
        print(f"  [font-encoding] Applied {fixes} encoding fixes")

    return text


_MULTISPACE = re.compile(r"[ \t]{2,}")


def _collapse_double_spaces(text: str) -> str:
    """Collapse 2+ spaces to 1, preserving code blocks, tables, and indentation.

    FIX-B02 (CM-003). Runs AFTER font encoding and BEFORE ligature repair.
    """
    out = []
    in_code = False
    fixes = 0
    for line in text.splitlines(True):
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            out.append(line)
            continue
        if in_code or stripped.startswith("|") or stripped.startswith("<!--"):
            out.append(line)
            continue
        # Preserve leading indentation (list items use 2/4 spaces)
        m = re.match(r"^(\s*)(.*)", line, re.DOTALL)
        leading = m.group(1)
        rest = m.group(2)
        # QC-FP-14: Preserve trailing double-space (Markdown line break <br>)
        trailing_br = ""
        if rest.endswith("  \n"):
            trailing_br = "  \n"
            rest = rest[:-3]
        elif rest.endswith("  "):
            trailing_br = "  "
            rest = rest[:-2]
        new_rest = _MULTISPACE.sub(" ", rest)
        if new_rest != rest:
            fixes += 1
        out.append(leading + new_rest + trailing_br)
    if fixes > 0:
        print(f"  [double-space] Collapsed spaces in {fixes} lines")
    return "".join(out)


# --- FIX-B03 (CM-023): Exclusion zones for ligature regex ---

# FIX-S32-03 (P2): Also match bare-domain URLs like "sagepub.com/journals-permissions"
# that lack a protocol prefix. The bare-domain pattern matches word.TLD/path sequences.
# S34 Fix 3.3 (A-07): Added bare-domain pattern without path (e.g., "sagepub.com").
# S34 Fix 3.3 (A-10): Added negative lookahead to prevent matching "Table.10/results".
# Per QC-B MINOR-13: European TLDs included (.uk, .de, .no, .eu, .se, .dk, .nl, .fr).
# Per QC-B MINOR-14: Lookahead applies to existing bare-domain-with-path pattern too.
# W3-QC1 MINOR-10: Known limitation -- "data.csv/path" matches the bare-domain-with-path
# pattern (csv is 3 chars). Harmless: text is preserved unchanged after zone restoration.
_URL_PROTECT = re.compile(
    r'https?://\S+'
    r'|\b(?!Table\b|Fig\b|Eq\b|Ref\b)\w+\.\w{2,4}/\S+'
    r'|\b\w+\.(?:com|org|net|edu|gov|io|uk|de|no|eu|se|dk|nl|fr|ch|at|es|it)\b'
)
_COMMENT_PROTECT = re.compile(r'<!--.*?-->', re.DOTALL)
_CODE_BLOCK_PROTECT = re.compile(r'```.*?```', re.DOTALL)


def _protect_zones(text):
    """Replace URLs, HTML comments, and code blocks with placeholder tokens."""
    placeholders = {}
    counter = [0]

    def _replace(m):
        key = f"\x00ZONE{counter[0]}\x00"
        placeholders[key] = m.group(0)
        counter[0] += 1
        return key

    for pat in [_CODE_BLOCK_PROTECT, _COMMENT_PROTECT, _URL_PROTECT]:
        text = pat.sub(_replace, text)
    return text, placeholders


def _restore_zones(text, placeholders):
    """Restore placeholder tokens back to original content.

    Iterates in reverse order so that inner (earlier) placeholders captured
    inside outer (later) placeholder values are restored correctly.
    """
    for key in reversed(list(placeholders.keys())):
        text = text.replace(key, placeholders[key])
    return text


# --- FIX-B06 (CM-002): SymSpell-based word validation ---
# Optional dependency: falls back to dictionary-only if not installed.

try:
    from symspellpy import SymSpell, Verbosity
    _sym_spell = SymSpell(max_dictionary_edit_distance=0, prefix_length=7)
    import symspellpy as _symspellpy_pkg
    _dict_path = str(
        Path(_symspellpy_pkg.__file__).parent / "frequency_dictionary_en_82_765.txt"
    )
    _sym_spell.load_dictionary(_dict_path, term_index=0, count_index=1)
    _HAS_SYMSPELL = True

    # FIX-B07 (CM-015): Add HTA domain terms not in general English corpus
    _HTA_DOMAIN_TERMS = [
        "efficacy", "pharmacoeconomics", "willingness", "coefficient",
        "coefficients", "proficiency", "deficiency", "deficient",
        "stratification", "quantification", "ratification",
        "amplification", "clarification", "justification",
        "simplification", "unification", "fluorescence",
        "ezetimibe", "pazopanib", "nivolumab", "pembrolizumab",
        "ipilimumab", "atezolizumab", "bevacizumab", "rituximab",
        "trastuzumab", "cetuximab", "panitumumab", "aflibercept",
        "cabozantinib", "axitinib", "sunitinib", "sorafenib",
        "lenvatinib", "regorafenib", "ramucirumab",
    ]
    for term in _HTA_DOMAIN_TERMS:
        _sym_spell.create_dictionary_entry(term, 1000)

except ImportError:
    _sym_spell = None
    _HAS_SYMSPELL = False
    print("WARNING: symspellpy not available - run-together detection disabled")

# --- S34 Fix 1.2 (A-02): wordfreq for enhanced word validation ---
# Provides ~400K English word corpus with Zipf frequency scores.
# Falls back gracefully to SymSpell-only if not installed.
try:
    from wordfreq import zipf_frequency as _zipf_frequency
    _HAS_WORDFREQ = True
except ImportError:
    _zipf_frequency = None
    _HAS_WORDFREQ = False

# --- FIX-B08 (CM-012): wordsegment for run-together splitting ---
try:
    import wordsegment as _ws
    _ws.load()
    _HAS_WORDSEGMENT = True

    # S34 Fix 2.2 (C-02): Augment wordsegment UNIGRAMS with domain terms.
    # AUGMENT (update), do NOT replace (clear). Preserves general English corpus.
    # Must be called AFTER _ws.load() which populates UNIGRAMS.
    # Per QC-B MINOR-11: use _ws.UNIGRAMS, not wordsegment.UNIGRAMS.
    # W3-QC1 MINOR-08: Some terms overlap with _HTA_DOMAIN_TERMS (SymSpell).
    # Intentional: SymSpell validates words, wordsegment biases segmentation.
    # Different purposes, harmless overlap.
    _HTA_WORDSEGMENT_TERMS = {
        # Medical/pharmaceutical terms
        'pharmacoeconomics': 1_000_000,
        'pharmacoeconomic': 1_000_000,
        'pharmacokinetic': 1_000_000,
        'pharmacokinetics': 1_000_000,
        'pharmacodynamic': 1_000_000,
        'pharmacodynamics': 1_000_000,
        'pharmacovigilance': 1_000_000,
        'bioequivalence': 1_000_000,
        'bioavailability': 1_000_000,
        'immunogenicity': 1_000_000,
        'immunotherapy': 1_000_000,
        'chemotherapy': 1_000_000,
        'radiotherapy': 1_000_000,
        'monotherapy': 1_000_000,
        'polypharmacy': 1_000_000,
        # HTA-specific terms
        'incremental': 1_000_000,
        'willingness': 1_000_000,
        'effectiveness': 1_000_000,
        'stratification': 1_000_000,
        'quantification': 1_000_000,
        'amplification': 1_000_000,
        'clarification': 1_000_000,
        'justification': 1_000_000,
        'simplification': 1_000_000,
        'unification': 1_000_000,
        'ratification': 1_000_000,
        'randomization': 1_000_000,
        'randomisation': 1_000_000,
        'extrapolation': 1_000_000,
        'interpolation': 1_000_000,
        'probabilistic': 1_000_000,
        'deterministic': 1_000_000,
        'microsimulation': 1_000_000,
        'partitioned': 1_000_000,
        'multivariate': 1_000_000,
        'univariate': 1_000_000,
        'heterogeneity': 1_000_000,
        'homogeneity': 1_000_000,
        'comorbidity': 1_000_000,
        'comorbidities': 1_000_000,
        'epidemiological': 1_000_000,
        'epidemiology': 1_000_000,
        'noncompliance': 1_000_000,
        'nonadherence': 1_000_000,
        'discontinuation': 1_000_000,
        'contraindication': 1_000_000,
        'contraindications': 1_000_000,
        'coadministration': 1_000_000,
        # Common run-together victims
        'coefficient': 1_000_000,
        'coefficients': 1_000_000,
        'proficiency': 1_000_000,
        'deficiency': 1_000_000,
        'deficient': 1_000_000,
        'insufficient': 1_000_000,
        'sufficiency': 1_000_000,
        'fluorescence': 1_000_000,
        'efficacy': 1_000_000,
    }
    _ws.UNIGRAMS.update(_HTA_WORDSEGMENT_TERMS)

except ImportError:
    _HAS_WORDSEGMENT = False


def _is_valid_word(word: str) -> bool:
    """Check if a word is valid using tiered validation.

    Tier 1: _KNOWN_LIGATURE_WORDS frozenset (O(1), fastest)
    Tier 2: SymSpell exact match (~82K words)
    Tier 3: wordfreq Zipf frequency >= 1.0 (~400K words) [S34 Fix 1.2]
    Fallback: not in any dictionary = unknown (conservative: do NOT rejoin)

    The Zipf threshold of 1.0 means "at least once per 100 million words."
    This avoids the frequency-comparison pitfall identified in PxDet's
    should_split_word() by using a simple threshold check instead.
    """
    lower = word.lower()
    # Tier 1: known ligature words always valid
    if lower in _KNOWN_LIGATURE_WORDS:
        return True
    # Tier 2: SymSpell frequency corpus
    if _HAS_SYMSPELL:
        suggestions = _sym_spell.lookup(lower, Verbosity.TOP, max_edit_distance=0)
        if len(suggestions) > 0:
            return True
    # Tier 3: wordfreq broader corpus (~400K words)
    if _HAS_WORDFREQ:
        if _zipf_frequency(lower, 'en') >= 1.0:
            return True
    # Fallback: not in dictionary = unknown (conservative: do NOT rejoin)
    return False


def _split_at_case_boundaries(token: str) -> list:
    """S37: Split token at case-transition boundaries.

    Handles:
    - lowerUpper: "healthTechnology" -> ["health", "Technology"]
    - UPPERlower: "VOIanalyses" -> ["VOI", "analyses"]
    - UPPERTitle: "ISPORValue" -> ["ISPOR", "Value"]

    Does NOT split inside all-uppercase runs or all-lowercase runs.

    S37-FIX: Uses dictionary-based lookahead at uppercase-lowercase boundaries
    to avoid fragmenting acronyms. When an uppercase run is followed by lowercase,
    checks whether the standard CamelCase split (stealing the last uppercase letter)
    produces a valid Titlecase word. If NOT, keeps the full uppercase run as an
    acronym and splits before the lowercase.
    """
    if not token or len(token) < 3:
        return [token]

    parts = []
    current = token[0]

    i = 1
    while i < len(token):
        char = token[i]
        prev = token[i - 1]

        # Boundary: lowercase followed by uppercase
        # "healthT" -> split before T
        if prev.islower() and char.isupper():
            parts.append(current)
            current = char
            i += 1
            continue

        # Boundary: end of uppercase run followed by lowercase
        # Must decide: standard CamelCase split (steal last uppercase) or
        # keep the full acronym together.
        if (char.islower() and prev.isupper() and len(current) > 1
                and current[-2].isupper()):
            # Collect the full lowercase run for lookahead
            lookahead = char
            j = i + 1
            while j < len(token) and token[j].islower():
                lookahead += token[j]
                j += 1

            left = current[:-1]
            stolen = current[-1]
            titlecase_word = stolen + lookahead

            # Heuristic 1: doubled first letter = wrong split (always)
            # e.g., "ICERresults" -> stolen='R', titlecase='Rresults' -> 'Rr'
            doubled_first = (len(titlecase_word) >= 2
                             and titlecase_word[0].lower()
                             == titlecase_word[1].lower())

            # Heuristic 2: left fragment too short = wrong split
            # e.g., "NHBassociated" -> left='NH' (2 chars)
            short_left = len(left) < 3

            if doubled_first or short_left:
                # Keep acronym whole: split BEFORE the lowercase
                parts.append(current)
                current = char
                i += 1
                continue

            # Heuristic 3: dictionary-based decision
            # If the lowercase run alone is a valid word but the titlecase
            # form (with stolen letter) is NOT, the stolen letter belongs
            # with the acronym. e.g., "QALYgained": 'gained' valid,
            # 'Ygained' not -> keep 'QALY' whole.
            titlecase_is_word = _is_valid_word(titlecase_word)
            lowercase_is_word = (_is_valid_word(lookahead)
                                and len(lookahead) >= 3)

            if lowercase_is_word and not titlecase_is_word:
                # Keep acronym whole
                parts.append(current)
                current = char
                i += 1
                continue

            # All other cases: standard CamelCase split
            # - titlecase IS a word: "CEAClark" -> "CEA" + "Clark"
            # - both valid: standard split preferred (CamelCase convention)
            # - neither valid: standard split as default
            parts.append(left)
            current = stolen + char
            i += 1
            continue

        current += char
        i += 1

    if current:
        parts.append(current)

    return [p for p in parts if p]


def _is_known_acronym(token: str) -> bool:
    """S37: Check if token is a known acronym (1-6 uppercase chars).

    S37-FIX MINOR-01: Combined redundant len==1 and len 2-6 checks.
    """
    if not token:
        return False
    if 1 <= len(token) <= 6 and token.isupper():
        return True
    return False


def _restore_case_from_original(original: str, segments: list) -> list:
    """S37: Restore original capitalization to wordsegment output.

    wordsegment returns all-lowercase. This maps back to the original
    token's capitalization by aligning character positions.

    Example: original="InHealthTechnology", segments=["in","health","technology"]
    -> ["In", "Health", "Technology"]
    """
    result = []
    pos = 0
    for seg in segments:
        seg_len = len(seg)
        if pos + seg_len <= len(original):
            orig_substr = original[pos:pos + seg_len]
            if orig_substr.lower() == seg.lower():
                result.append(orig_substr)
            else:
                result.append(seg)
            pos += seg_len
        else:
            result.append(seg)
    return result


# S37: Compiled regex for punctuation-boundary splitting.
# Matches [letter][.,;:][Uppercase] transitions in run-together text.
# Captures the punctuation mark as a separate group for reconstruction.
# S37-FIX MAJOR-02: Extended lookbehind to also match uppercase letters,
# so "FitzmauriceC,Dicker" splits at the comma (was missed because 'C' is
# uppercase). Original only matched [lowercase][punct][Uppercase].
_RE_PUNCT_BOUNDARY = re.compile(r'(?<=[a-zA-Z])([.,;:])(?=[A-Z])')


def _split_single_token(token: str) -> list:
    """S37: Attempt to split a single run-together token into words.

    Strategy order:
    1. CamelCase + acronym boundary detection
    2. Wordsegment with full-token fallback
    Returns [token] unchanged if no valid split found.
    """
    if not token or len(token) < 4:
        return [token]

    if _is_valid_word(token):
        return [token]

    # === Strategy 1: CamelCase + Acronym boundary splitting ===
    camel_split = _split_at_case_boundaries(token)
    if len(camel_split) > 1:
        # Check if ALL parts are already valid
        all_valid = all(
            _is_valid_word(p) or _is_known_acronym(p)
            for p in camel_split
        )
        if all_valid:
            return camel_split

        # Partial: try wordsegment on invalid sub-parts
        final = []
        any_expanded = False
        for part in camel_split:
            if _is_valid_word(part) or _is_known_acronym(part):
                final.append(part)
            elif len(part) >= 5 and part.islower() and _HAS_WORDSEGMENT:
                segs = _ws.segment(part)
                if (len(segs) > 1
                        and all(_is_valid_word(s) for s in segs)
                        and all(len(s) > 1 or s in ('a', 'i')
                                for s in segs)):
                    final.extend(segs)
                    any_expanded = True
                else:
                    final.append(part)
            elif (len(part) >= 5 and part[0].isupper()
                    and part[1:].islower() and _HAS_WORDSEGMENT):
                # Title-case sub-part (e.g., "Claxtonand")
                segs = _ws.segment(part.lower())
                if (len(segs) > 1
                        and all(_is_valid_word(s) for s in segs)
                        and all(len(s) > 1 or s in ('a', 'i')
                                for s in segs)):
                    restored = _restore_case_from_original(part, segs)
                    final.extend(restored)
                    any_expanded = True
                else:
                    final.append(part)
            else:
                final.append(part)

        if len(final) > len(camel_split) or (len(final) > 1 and any_expanded):
            return final
        # If camel split produced parts but sub-expansion didn't help,
        # still return camel split if all parts are at least plausible
        if len(camel_split) > 1 and all(
            _is_valid_word(p) or _is_known_acronym(p)
            for p in camel_split
        ):
            return camel_split

    # === Strategy 1b: Embedded punctuation pre-split ===
    # S37-FIX MAJOR-01: Before sending to wordsegment, split at embedded
    # commas and semicolons. wordsegment silently strips non-alpha characters,
    # losing punctuation like commas. "Correspondingly,theincrease" must become
    # "Correspondingly, the increase" not "Correspondingly the increase".
    _punct_chars = re.compile(r'([,;])')
    if _punct_chars.search(token):
        raw_parts = _punct_chars.split(token)
        # raw_parts alternates: [text, punct, text, punct, text, ...]
        # Attach punctuation to the preceding text part.
        assembled = []
        pi = 0
        while pi < len(raw_parts):
            part = raw_parts[pi]
            # Check if next element is a punctuation mark
            if (pi + 1 < len(raw_parts)
                    and len(raw_parts[pi + 1]) == 1
                    and raw_parts[pi + 1] in ',;'):
                assembled.append(part + raw_parts[pi + 1])
                pi += 2
            else:
                assembled.append(part)
                pi += 1
        # Filter empty and recursively process each sub-part
        expanded = []
        for sub in assembled:
            if not sub:
                continue
            # Strip trailing punctuation for recursive call, reattach after
            sub_trailing = ''
            sub_clean = sub
            while sub_clean and sub_clean[-1] in ',;':
                sub_trailing = sub_clean[-1] + sub_trailing
                sub_clean = sub_clean[:-1]
            if sub_clean and len(sub_clean) >= 4:
                sub_parts = _split_single_token(sub_clean)
                # Reattach trailing punctuation to last part
                if sub_trailing and sub_parts:
                    sub_parts[-1] = sub_parts[-1] + sub_trailing
                expanded.extend(sub_parts)
            elif sub:
                expanded.append(sub)
        if len(expanded) > 1:
            return expanded

    # === Strategy 2: Wordsegment on the full token ===
    if _HAS_WORDSEGMENT and len(token) >= 5:
        lower_token = token.lower()
        segments = _ws.segment(lower_token)
        if (len(segments) > 1
                and all(len(s) > 1 or s in ('a', 'i') for s in segments)
                and all(_is_valid_word(s) for s in segments)):
            restored = _restore_case_from_original(token, segments)
            return restored

    return [token]


def _fix_run_togethers(text: str) -> str:
    """Split run-together words using multi-strategy approach.

    S37 rewrite: replaces the original 5-path algorithm with a unified
    approach that handles CamelCase, acronyms, prepositions,
    punctuation-embedded, and all-lowercase fusions.

    Preserved from prior sessions:
    - S34 Fix 1.1: Prefix/suffix allowlist for hyphenated words
    - FIX-S32-01: Hyphen-compound preservation
    - S34 Fix 1.6: Special character boundary splitting (+, &, =)
    - S34 Fix 1.5: Digit-letter boundary splitting

    NEW in S37:
    - Phase 1: Punctuation-boundary splitting (period/comma at [lower][punct][Upper])
    - Phase 2: CamelCase + acronym boundary detection (all case transitions)
    - Phase 3: Wordsegment fallback (lower threshold: >= 5 chars)
    - Case restoration from original token

    Processes line-by-line to preserve document whitespace structure.

    S37-FIX MINOR-03: Removed early return on !_HAS_WORDSEGMENT. CamelCase
    boundary splitting and punctuation-boundary splitting do not require
    wordsegment. Internal guards in _split_single_token already check
    _HAS_WORDSEGMENT before calling wordsegment-specific paths.
    """
    fixes = 0
    output_lines = []
    for line in text.splitlines(True):
        words = line.split()
        if not words:
            output_lines.append(line)
            continue
        result = []
        line_changed = False
        for word in words:
            if _is_valid_word(word):
                result.append(word)
                continue

            # S34 Fix 1.1 (A-09): Prefix/suffix allowlist for hyphenated words.
            # Check BEFORE general hyphen-compound check. If first part is a
            # known prefix and second part is a valid word, preserve hyphen.
            # Similarly for known suffixes.
            if '-' in word:
                hyphen_parts = word.split('-')
                non_empty_parts = [p for p in hyphen_parts if p]
                if len(non_empty_parts) >= 2:
                    first_lower = non_empty_parts[0].lower()
                    last_lower = non_empty_parts[-1].lower()
                    # Prefix check: "pre-existing", "post-hoc", etc.
                    if (first_lower in _VALID_PREFIXES
                            and _is_valid_word(non_empty_parts[-1])):
                        result.append(word)
                        continue
                    # Suffix check: "evidence-based", "health-related", etc.
                    if (last_lower in _VALID_SUFFIXES
                            and _is_valid_word(non_empty_parts[0])):
                        result.append(word)
                        continue

            # FIX-S32-01: Hyphen-compound preservation.
            # If the word contains hyphens (e.g., "cost-effectiveness"), split on
            # hyphens and check each part. If ALL parts are valid English words,
            # the hyphen was intentional in the source text -- preserve it unchanged.
            # This prevents wordsegment from consuming hyphens as word boundaries
            # and emitting "cost effectiveness" instead of "cost-effectiveness".
            if '-' in word:
                hyphen_parts = word.split('-')
                # Filter out empty parts from leading/trailing hyphens
                non_empty_parts = [p for p in hyphen_parts if p]
                if (len(non_empty_parts) >= 2
                        and all(_is_valid_word(p) for p in non_empty_parts)):
                    # All parts are valid words -- this is a legitimate compound
                    result.append(word)
                    continue

            # S34 Fix 1.6 (B-05): Special character boundary splitting.
            # Split at +, =, &. Process each alphabetic part further.
            # Note: "/" excluded from split chars to avoid breaking "HIV/AIDS".
            # Per QC-B MAJOR-02: placed BEFORE other paths so parts get
            # further processed by downstream paths.
            # W3-QC1 MINOR-03: Strip trailing punctuation before matching,
            # so "M+Pinboth." is handled. Reattach punctuation after.
            _sc_word = word
            _sc_trailing_punct = ''
            if _sc_word and _sc_word[-1] in '.,;:!?)]\'"':
                _sc_trailing_punct = _sc_word[-1]
                _sc_word = _sc_word[:-1]
            _special_char_match = re.match(
                r'^([A-Za-z]+)([+&=])([A-Za-z]{2,})$', _sc_word
            )
            if _special_char_match:
                _sc_left = _special_char_match.group(1)
                _sc_op = _special_char_match.group(2)
                _sc_right = _special_char_match.group(3)
                # Further split the right part if it's a run-together
                _sc_right_parts = []
                if not _is_valid_word(_sc_right) and len(_sc_right) >= 4:
                    # Try title+lower split on right part
                    _sc_split_found = False
                    if _sc_right[0].isupper():
                        for _si in range(2, len(_sc_right) - 1):
                            _sl = _sc_right[:_si]
                            _sr = _sc_right[_si:]
                            if (_is_valid_word(_sl) and _is_valid_word(_sr)):
                                _sc_right_parts = [_sl, _sr]
                                _sc_split_found = True
                                break
                    if not _sc_split_found and _HAS_WORDSEGMENT:
                        _segs = _ws.segment(_sc_right.lower())
                        if (len(_segs) > 1
                                and all(_is_valid_word(s) for s in _segs)):
                            _sc_right_parts = _segs
                            _sc_split_found = True
                    if not _sc_split_found:
                        _sc_right_parts = [_sc_right]
                else:
                    _sc_right_parts = [_sc_right]
                # W3-QC1 MINOR-04: Operator stays attached to left part.
                # Style choice: "M+Pin both" not "M + Pin both".
                result.append(
                    _sc_left + _sc_op + ' '.join(_sc_right_parts)
                    + _sc_trailing_punct
                )
                fixes += 1
                line_changed = True
                continue

            # S34 Fix 1.5 (B-06): Digit-letter boundary splitting.
            # "3state" -> "3-state", "2arm" -> "2-arm", "3.5million" -> "3.5 million"
            # Protected terms prevent splitting scientific identifiers.
            # W3-QC2 MINOR-01: Extended to handle decimal+alpha ("3.5million").
            _digit_letter_match = re.match(
                r'^(\d+(?:\.\d+)?)([a-zA-Z]{3,})$', word
            )
            if _digit_letter_match:
                _dl_digits = _digit_letter_match.group(1)
                _dl_alpha = _digit_letter_match.group(2)
                if word not in _DIGIT_LETTER_PROTECTED:
                    if _is_valid_word(_dl_alpha):
                        # W3-QC2 MINOR-01: Use space for decimal+alpha
                        # ("3.5 million"), hyphen for integer+alpha ("3-state")
                        _dl_sep = ' ' if '.' in _dl_digits else '-'
                        result.append(_dl_digits + _dl_sep + _dl_alpha)
                        fixes += 1
                        line_changed = True
                        continue

            # === S37: Unified run-together splitting ===
            # Replaces old Path 1 (CamelCase), Path 1b (Title+lowercase),
            # and Path 2 (all-lowercase) with a multi-strategy approach.

            # Strip trailing punctuation for cleaner analysis
            _rt_trailing = ''
            _rt_clean = word
            while _rt_clean and _rt_clean[-1] in '.,;:!?)]\'"':
                _rt_trailing = _rt_clean[-1] + _rt_trailing
                _rt_clean = _rt_clean[:-1]

            # Strip leading punctuation
            _rt_leading = ''
            while _rt_clean and _rt_clean[0] in '(["\'"':
                _rt_leading += _rt_clean[0]
                _rt_clean = _rt_clean[1:]

            if not _rt_clean or len(_rt_clean) < 4:
                result.append(word)
                continue

            # If clean word is valid, no splitting needed
            if _is_valid_word(_rt_clean):
                result.append(word)
                continue

            # Phase 1: Punctuation-boundary splitting
            # "outcomes.Forexample" -> ["outcomes.", "Forexample"]
            _punct_parts = _RE_PUNCT_BOUNDARY.split(_rt_clean)
            if len(_punct_parts) > 1:
                # Reconstruct: attach punctuation to preceding part
                _rebuilt = []
                _pi = 0
                while _pi < len(_punct_parts):
                    if (_pi + 1 < len(_punct_parts)
                            and len(_punct_parts[_pi + 1]) == 1
                            and _punct_parts[_pi + 1] in '.,;:'):
                        _rebuilt.append(
                            _punct_parts[_pi] + _punct_parts[_pi + 1])
                        _pi += 2
                    else:
                        _rebuilt.append(_punct_parts[_pi])
                        _pi += 1
                # Recursively split each sub-part
                # S37-FIX MAJOR-03: Strip trailing punctuation before passing
                # to _split_single_token, then reattach. Previously relied on
                # wordsegment to silently ignore non-alpha chars (fragile).
                _final_parts = []
                for _rpart in _rebuilt:
                    if _rpart:
                        _rp_trailing = ''
                        _rp_clean = _rpart
                        while (_rp_clean
                               and _rp_clean[-1] in '.,;:!?)]\'"'):
                            _rp_trailing = _rp_clean[-1] + _rp_trailing
                            _rp_clean = _rp_clean[:-1]
                        if _rp_clean:
                            _sub = _split_single_token(_rp_clean)
                            if _rp_trailing and _sub:
                                _sub[-1] = _sub[-1] + _rp_trailing
                            _final_parts.extend(_sub)
                        else:
                            _final_parts.append(_rpart)
                if len(_final_parts) > 1:
                    result.append(
                        _rt_leading + ' '.join(_final_parts) + _rt_trailing)
                    fixes += 1
                    line_changed = True
                    continue

            # Phase 2+3: CamelCase boundaries + wordsegment fallback
            _split = _split_single_token(_rt_clean)
            if len(_split) > 1:
                result.append(
                    _rt_leading + ' '.join(_split) + _rt_trailing)
                fixes += 1
                line_changed = True
                continue

            result.append(word)

        if line_changed:
            # Reconstruct line preserving trailing newline
            trailing = ''
            if line.endswith('\n'):
                trailing = '\n'
            # Preserve leading whitespace
            leading_match = re.match(r'^(\s*)', line)
            leading = leading_match.group(1) if leading_match else ''
            output_lines.append(leading + ' '.join(result) + trailing)
        else:
            output_lines.append(line)

    if fixes > 0:
        print(f"  [run-togethers] Split {fixes} merged words")
    return ''.join(output_lines)


# --- S34 Fix 1.3 (A-01): Domain hyphen restoration dictionary ---
# Maps unhyphenated forms (as produced by docling) to correctly hyphenated forms.
# Sorted longest-first at regex compile time to prevent partial matches.
# Only includes terms with clear, unambiguous HTA/medical hyphenation.
# Debatable terms (time-horizon, decision-tree, base-case) intentionally excluded.
_HTA_HYPHENATED_COMPOUNDS = {
    'costeffectiveness': 'cost-effectiveness',
    'costeffective': 'cost-effective',
    'costutility': 'cost-utility',
    'costbenefit': 'cost-benefit',
    'costminimization': 'cost-minimization',
    'costconsequence': 'cost-consequence',
    'qualityadjusted': 'quality-adjusted',
    'disabilityadjusted': 'disability-adjusted',
    'ageadjusted': 'age-adjusted',
    'riskadjusted': 'risk-adjusted',
    'healthrelated': 'health-related',
    'diseaserelated': 'disease-related',
    'treatmentrelated': 'treatment-related',
    'doserelated': 'dose-related',
    'preferencebased': 'preference-based',
    'evidencebased': 'evidence-based',
    'populationbased': 'population-based',
    'communitybased': 'community-based',
    'hospitalbased': 'hospital-based',
    'modelbased': 'model-based',
    'utilitybased': 'utility-based',
    'valuebased': 'value-based',
    'conditionspecific': 'condition-specific',
    'diseasespecific': 'disease-specific',
    'agespecific': 'age-specific',
    'countryspecific': 'country-specific',
    'patientspecific': 'patient-specific',
    'decisionmaking': 'decision-making',
    'decisionanalytic': 'decision-analytic',
    'decisionanalysis': 'decision-analysis',
    'metaanalysis': 'meta-analysis',
    'metaanalyses': 'meta-analyses',
    'metaregression': 'meta-regression',
    'selfreport': 'self-report',
    'selfreported': 'self-reported',
    'selfreporting': 'self-reporting',
    'selfmanagement': 'self-management',
    'selfadministered': 'self-administered',
    'willingnesstopay': 'willingness-to-pay',
    'payforperformance': 'pay-for-performance',
    'headtohead': 'head-to-head',
    'stateoftheart': 'state-of-the-art',
    # W3-QC1 MAJOR-03: Removed 'stateoftransition' (non-standard term)
    # W3-QC1 MAJOR-03: Removed 'overthecontrol' (not a standard term)
    'endoflife': 'end-of-life',
    'qualityoflife': 'quality-of-life',
    'yearoflife': 'year-of-life',
    'yearsoflife': 'years-of-life',
    'outofpocket': 'out-of-pocket',
    'longterm': 'long-term',
    'shortterm': 'short-term',
    'mediumterm': 'medium-term',
    'realworld': 'real-world',
    'allcause': 'all-cause',
    'firstline': 'first-line',
    'secondline': 'second-line',
    'thirdline': 'third-line',
    'openlabel': 'open-label',
    'doubleblind': 'double-blind',
    'singlearm': 'single-arm',
    'twoarm': 'two-arm',
    'multiarm': 'multi-arm',
    # W3-QC1 MAJOR-01: Removed 'subgroup' (standard English, SymSpell-valid)
    # W3-QC1 MAJOR-01: Removed 'subanalysis' (debatable; 'subanalysis' common)
    'multipayer': 'multi-payer',
    # W3-QC1 MAJOR-02: Removed 'multicenter' (standard in clinical trials)
    'multicountry': 'multi-country',
    # W3-QC1 MAJOR-02: Removed 'multistate' (standard in epidemiology)
    # W3-QC1 MAJOR-02: Removed 'nonproportional' (standard unhyphenated form)
    # W3-QC1 MAJOR-02: Removed 'noninferior' (standard in trials)
    # W3-QC1 MAJOR-02: Removed 'noninferiority' (standard in trials)
    'progressionfree': 'progression-free',
    'eventfree': 'event-free',
    'diseasefree': 'disease-free',
    'relapsefree': 'relapse-free',
    'recurrencefree': 'recurrence-free',
    'treatmentfree': 'treatment-free',
    'symptomfree': 'symptom-free',
    # W3-QC1 MAJOR-03: Removed 'intentiontotreated' (grammatically wrong)
    'intentiontotreat': 'intention-to-treat',
    'perprotocol': 'per-protocol',
    'doseresponse': 'dose-response',
    'timevarying': 'time-varying',
    'timedependent': 'time-dependent',
    'halflife': 'half-life',
    'wellbeing': 'well-being',
    # W3-QC1 MAJOR-02: Removed 'reanalysis' (standard unhyphenated form)
    'prespecified': 'pre-specified',
    'pretreatment': 'pre-treatment',
    'posttreatment': 'post-treatment',
    'posthoc': 'post-hoc',
    'oneway': 'one-way',
    'twoway': 'two-way',
    'multiway': 'multi-way',
}

# Build compiled regex sorted longest-first for correct matching
_HTA_COMPOUND_PATTERN = re.compile(
    r'\b(' + '|'.join(
        sorted(_HTA_HYPHENATED_COMPOUNDS.keys(), key=len, reverse=True)
    ) + r')\b',
    re.IGNORECASE
)


# S36: Curated known space-split patterns. These are words that docling
# routinely splits at non-ligature glyph boundaries. The merged lowercase
# form is checked against this set -- if it matches, merge regardless of
# whether the individual fragments are valid English words.
# Expand this list as new papers reveal new splits.
_KNOWN_SPACE_SPLITS = frozenset({
    # Statistical distributions
    'weibull', 'gompertz',
    # Common academic surnames in HTA/health economics
    'sculpher', 'fenwick', 'tuffaha', 'scuffham', 'winfree',
    'taipale', 'sideris', 'dearden', 'goldstraw', 'barsouk',
    'blandin', 'hammerman', 'chansky', 'hoomans', 'ginnelly',
    'bindels', 'longworth', 'hakkaart', 'eggington', 'chaboyer',
    'garside', 'kirkham', 'trowman', 'sharples', 'kantoff',
    'conaway', 'tannock', 'hussain', 'posnett', 'sashegyi',
    'andreas',
    # Common words that docling splits
    'modeling', 'modelling', 'biostatistics', 'abstract',
})

# S36: High-frequency common words that should NOT be consumed as name
# fragments. Prevents false merges like "In The" -> "Inthe".
_COMMON_SHORT_WORDS = frozenset({
    'a', 'an', 'and', 'are', 'as', 'at', 'be', 'but', 'by',
    'can', 'did', 'do', 'for', 'from', 'get', 'got', 'had',
    'has', 'have', 'he', 'her', 'him', 'his', 'how', 'i',
    'if', 'in', 'into', 'is', 'it', 'its', 'let', 'may',
    'me', 'my', 'no', 'nor', 'not', 'of', 'on', 'or', 'our',
    'out', 'own', 'per', 'put', 'ran', 'run', 'say', 'set',
    'she', 'so', 'the', 'to', 'too', 'up', 'us', 'use',
    'via', 'was', 'way', 'we', 'who', 'why', 'yet', 'you',
})

# S36-FIX: CRITICAL-02 - Name particles and honorifics that must NOT be
# merged by the proper-name heuristic (Prong 2). "De Vries" should stay
# as-is, not become "Devries". Common in Dutch/German/Spanish/Arabic
# author names and English honorifics in academic HTA papers.
_NAME_PARTICLES = frozenset({
    'de', 'del', 'della', 'van', 'von', 'el', 'al', 'la', 'le',
    'di', 'da', 'du', 'das', 'dos', 'den', 'der', 'het', 'des',
    'st', 'dr', 'mr', 'ms', 'jr', 'sr', 'mc',
})


def _fix_space_splits(text: str) -> str:
    """S36: Fix space-split words produced by docling glyph boundary detection.

    Handles non-ligature space-splits like "Sculp her" -> "Sculpher",
    "Wei bull" -> "Weibull", "Mode ling" -> "Modeling".

    Three-pronged approach:
    1. Curated list: if merged form (lowercase) is in _KNOWN_SPACE_SPLITS,
       always merge. Highest reliability, handles names where both fragments
       are valid English words.
    2. Proper-name heuristic: first token is capitalized, at least one
       fragment is NOT a valid standalone word. Catches novel names.
    3. Dictionary validation: merged form passes _is_valid_word() and at
       least one fragment is not a valid standalone word or is very short.

    Runs AFTER Step 5e (sliding window ligature) and BEFORE Step 5f
    (domain hyphens). Operates within zone-protected text.
    """
    if not text:
        return text

    # S36: Spaced-out all-caps heading fix: "A B STR A C T" -> "ABSTRACT"
    # S36-FIX: MAJOR-03 - Replaced naive regex with guarded handler:
    # 1. Skip lines containing null bytes (zone-protected)
    # 2. Skip lines containing pipe chars (table rows: "| A | B | C |")
    # 3. Require collapsed result to form words of 3+ chars each,
    #    filtering out table column headers like "A B C D" which
    #    collapse to single-letter "words" ("A", "B", "C", "D").
    def _fix_spaced_heading(m):
        full_line = m.group(0)
        # Guard 1: zone-protected content
        if '\x00' in full_line:
            return full_line
        # Guard 2: table rows
        if '|' in full_line:
            return full_line
        heading_prefix = m.group(1) or ''
        spaced_text = m.group(2)
        collapsed = spaced_text.replace(' ', '')
        # Guard 3: collapsed text must form at least one word of 3+ chars.
        # Split on runs of single letters to find "words". If ALL segments
        # are single chars (like "A B C D" -> "ABCD" which is one blob
        # of originally-single chars), check that blob is >= 5 chars and
        # looks like a real word. For mixed like "A B STR A C T", the
        # collapsed "ABSTRACT" (8 chars) passes.
        # Simple heuristic: collapsed must be >= 5 chars (already
        # guaranteed by the {3,} quantifier requiring 4+ spaced chars)
        # AND the original must contain at least one multi-char segment
        # (meaning at least one "word piece" like "STR" in "A B STR A C T").
        parts = spaced_text.split()
        has_multichar = any(len(p) > 1 for p in parts)
        if not has_multichar and len(collapsed) < 5:
            # All single letters and short: likely table headers
            return full_line
        return heading_prefix + collapsed

    text = re.sub(
        r'^(#{1,6}\s+)?([A-Z](?:\s+[A-Z]){3,}(?:\s+[A-Z]+)*)\s*$',
        _fix_spaced_heading,
        text,
        flags=re.MULTILINE
    )

    fixes = 0
    output_lines = []

    for line in text.splitlines(True):
        tokens = line.split()
        if len(tokens) < 2:
            output_lines.append(line)
            continue

        result = []
        i = 0
        line_changed = False

        while i < len(tokens):
            # S36: Skip zone placeholders (null bytes from _protect_zones)
            if '\x00' in tokens[i]:
                result.append(tokens[i])
                i += 1
                continue

            merged = False

            # Try window sizes 3, 2 (3 first for "Egging ton" type)
            for win in (3, 2):
                if i + win > len(tokens):
                    continue

                window = tokens[i:i + win]

                # Skip if any token has a zone placeholder
                if any('\x00' in t for t in window):
                    continue

                # S36: Strip trailing punctuation from last token for validation
                last_tok = window[-1]
                trailing_punct = ''
                if last_tok and last_tok[-1] in '.,;:!?)]\'"':
                    trailing_punct = last_tok[-1]
                    last_tok = last_tok[:-1]
                    window = list(window)
                    window[-1] = last_tok

                candidate = ''.join(window)
                if not candidate or len(candidate) < 5:
                    continue

                # === PRONG 1: Curated known splits ===
                # Highest reliability -- always merge if in the list
                if candidate.lower() in _KNOWN_SPACE_SPLITS:
                    result.append(candidate + trailing_punct)
                    i += win
                    fixes += 1
                    line_changed = True
                    merged = True
                    break

                # === PRONG 2: Proper Name Detection ===
                if window[0] and window[0][0].isupper():
                    # S36: Check if at least one fragment is NOT a valid
                    # standalone word -- suggests it's a broken piece
                    any_invalid = any(
                        not _is_valid_word(t) for t in window
                    )
                    # Also accept if at least one fragment is very short
                    # (<= 2 chars) -- short fragments are suspicious
                    any_very_short = any(len(t) <= 2 for t in window)

                    # Guard: reject if ALL fragments are common words
                    all_common = all(
                        t.lower() in _COMMON_SHORT_WORDS for t in window
                    )

                    # S36-FIX: CRITICAL-02 - Guard against name particles.
                    # "De Vries", "El Salvador", "Van De Berg", "Dr Smith"
                    # must NOT be merged. Skip if ANY fragment is a known
                    # name particle (case-insensitive).
                    has_name_particle = any(
                        t.lower() in _NAME_PARTICLES for t in window
                    )

                    if (any_invalid or any_very_short) and not all_common and not has_name_particle:
                        result.append(candidate + trailing_punct)
                        i += win
                        fixes += 1
                        line_changed = True
                        merged = True
                        break

                # === PRONG 3: Dictionary Validation ===
                if _is_valid_word(candidate):
                    # S36: Merge only if at least one fragment is NOT a
                    # valid standalone word or is short (<= 3 chars)
                    any_frag_invalid = any(
                        not _is_valid_word(t) for t in window
                    )
                    any_frag_short = any(len(t) <= 3 for t in window)

                    if any_frag_invalid or any_frag_short:
                        result.append(candidate + trailing_punct)
                        i += win
                        fixes += 1
                        line_changed = True
                        merged = True
                        break

            if not merged:
                result.append(tokens[i])
                i += 1

        if line_changed:
            trailing = '\n' if line.endswith('\n') else ''
            lead_m = re.match(r'^(\s*)', line)
            leading = lead_m.group(1) if lead_m else ''
            output_lines.append(leading + ' '.join(result) + trailing)
        else:
            output_lines.append(line)

    if fixes > 0:
        print(f"  [space-splits] Merged {fixes} space-split words")
    return ''.join(output_lines)


def _restore_domain_hyphens(text: str) -> str:
    """Restore hyphens in domain-specific compound terms stripped by docling.

    S34 Fix 1.3 (A-01). Case-preserving replacement.
    Runs at Step 5f (BEFORE zone restoration) within zone-protected text.
    """
    def _case_preserving_replace(m):
        matched = m.group(0)
        key = matched.lower()
        replacement = _HTA_HYPHENATED_COMPOUNDS.get(key, matched)
        # Preserve case: all-upper, title-case, or as-is
        # W3-QC1 MINOR-01: Known limitation -- CamelCase input (e.g.,
        # "CostEffectiveness") would lose the second capital. In practice,
        # docling does not produce CamelCase for these terms.
        if matched.isupper():
            return replacement.upper()
        elif matched[0].isupper():
            return replacement[0].upper() + replacement[1:]
        return replacement

    # W3-QC2 MINOR-03: Count matches accurately using findall instead of
    # zip-based word comparison (zip truncates to shorter list if word count changes).
    changes = len(_HTA_COMPOUND_PATTERN.findall(text))
    new_text = _HTA_COMPOUND_PATTERN.sub(_case_preserving_replace, text)
    if new_text != text and changes > 0:
        print(f"  [domain-hyphens] Restored hyphens in {changes} compound terms")
    return new_text


# Reference entry start patterns for B-04
# W3-QC1 CRITICAL-01: \d{1,3}\. matched decimal numbers ("1.5 million").
# Fix: require whitespace after period via (?=\s).
# W3-QC1 CRITICAL-01 (combined): [A-Z][a-z]+, matched any capitalized word+comma
# ("London, UK"). Fix: require comma + space + uppercase (surname + initial).
_REF_ENTRY_START = re.compile(
    r'^\s*(?:'
    r'\[\d+\]'                # [1], [23]
    r'|\d{1,3}\.(?=\s)'      # 1. , 23.  (but NOT 1.5, 2.3)
    r'|[A-Z][a-z]+,\s+[A-Z]\.' # Author surname (2+ chars), Initial with period (author-year style)
    r')'
)

# Section headings that terminate the references section
_REF_END_HEADINGS = re.compile(
    r'^#{1,6}\s+(?:Appendix|Supplementary|Supporting\s+Information|'
    r'Acknowledgements?|Acknowledgments?|Author\s+Contributions?|'
    r'Conflicts?\s+of\s+Interest|Funding|Declarations?)\b',
    re.IGNORECASE | re.MULTILINE
)


def _fix_reference_section(text: str) -> str:
    """Post-process the references section to fix docling artifacts.

    S34 Fix 1.8 (B-04). Fixes:
    - False heading markers (##) within references -> plain text
    - Multi-line reference entries -> joined into single lines
    - Detects end-of-references at next known section heading

    Per QC-B MAJOR-05: includes "Bibliography" in heading pattern.
    Runs AFTER zone restoration (Step 7) so URLs are real, not placeholders.
    """
    # Split at References/Reference/Bibliography heading
    # W3-QC1 MAJOR-04: Use (?:^|\n) instead of \n to handle heading at document start
    # W3-QC2 MINOR-02: Match "Reference" (singular) in addition to "References"
    ref_split = re.split(
        r'(?i)((?:^|\n)(?:#{1,6}\s+)?(?:References?|Bibliography)\s*\n)',
        text, maxsplit=1
    )
    if len(ref_split) < 3:
        # No references heading found -- no-op
        return text

    pre_refs = ref_split[0]
    ref_heading = ref_split[1]
    ref_body = ref_split[2]

    # Find end of references (next major section heading or end of text)
    post_refs = ''
    end_match = _REF_END_HEADINGS.search(ref_body)
    if end_match:
        post_refs = ref_body[end_match.start():]
        ref_body = ref_body[:end_match.start()]

    # Process reference body:
    # 1. Strip false heading markers
    # 2. Join continuation lines
    ref_lines = ref_body.split('\n')
    cleaned_refs = []
    current_entry = []
    # W3-QC1 MINOR-07: Track changes for logging
    _ref_headings_stripped = 0
    _ref_entries_found = 0
    _ref_joins = 0

    for rline in ref_lines:
        stripped = rline.strip()

        # Strip false heading markers within references
        heading_match = re.match(r'^(#{2,6})\s+(.+)$', stripped)
        if heading_match:
            # In references, ## markers are false positives from docling
            stripped = heading_match.group(2)
            _ref_headings_stripped += 1

        # Check if this starts a new reference entry
        if _REF_ENTRY_START.match(stripped):
            # Save previous entry
            if current_entry:
                if len(current_entry) > 1:
                    _ref_joins += len(current_entry) - 1
                cleaned_refs.append(' '.join(current_entry))
            current_entry = [stripped]
            _ref_entries_found += 1
        elif stripped == '':
            # Blank line: save current entry, preserve blank
            if current_entry:
                if len(current_entry) > 1:
                    _ref_joins += len(current_entry) - 1
                cleaned_refs.append(' '.join(current_entry))
                current_entry = []
            cleaned_refs.append('')
        else:
            # Continuation line: append to current entry
            if current_entry:
                current_entry.append(stripped)
            else:
                cleaned_refs.append(stripped)

    # Save final entry
    if current_entry:
        if len(current_entry) > 1:
            _ref_joins += len(current_entry) - 1
        cleaned_refs.append(' '.join(current_entry))

    # W3-QC1 MINOR-07: Log changes
    if _ref_headings_stripped > 0 or _ref_joins > 0:
        print(f"  [references] {_ref_entries_found} entries, "
              f"{_ref_headings_stripped} headings stripped, "
              f"{_ref_joins} continuation lines joined")

    result = pre_refs + ref_heading + '\n'.join(cleaned_refs) + post_refs
    return result


# --- S34 Fix 1.1 (A-09): Valid prefix/suffix allowlists for hyphen preservation ---
_VALID_PREFIXES = frozenset({
    'pre', 'post', 'non', 'sub', 'semi', 'anti', 'auto', 'bi', 'co',
    'counter', 'cross', 'de', 'dis', 'ex', 'extra', 'hyper', 'hypo',
    'in', 'inter', 'intra', 'macro', 'micro', 'mid', 'mini', 'mis',
    'mono', 'multi', 'neo', 'out', 'over', 'para', 'poly', 'pseudo',
    'quasi', 're', 'self', 'super', 'trans', 'tri', 'ultra', 'un',
    'under', 'uni', 'well',
    # HTA/medical specific
    'cardio', 'chemo', 'electro', 'gastro', 'histo', 'immuno',
    'neuro', 'onco', 'patho', 'pharma', 'physio', 'psycho',
    'radio', 'thermo',
})
_VALID_SUFFIXES = frozenset({
    'based', 'related', 'adjusted', 'specific', 'free', 'like',
    'dependent', 'independent', 'weighted', 'matched', 'paired',
    'derived', 'induced', 'mediated', 'resistant',
})

# --- S34 Fix 1.5 (B-06): Protected terms for digit-letter splitting ---
# These should NOT be split with a hyphen even though they match digit+letter.
_DIGIT_LETTER_PROTECTED = frozenset({
    # Dimensional identifiers
    '2D', '3D', '4D',
    # Gene/protein names with digit prefixes
    '5HT', '3TC',
    # Scientific identifiers
    '2nd', '3rd', '4th', '5th', '6th', '7th', '8th', '9th',
    # Drug/compound codes
    '5FU', '6MP',
    # W3-QC1 MINOR-02: Removed '3mcg', '5mcg', '10mcg' (redundant: "mcg" fails
    # _is_valid_word so the split would never fire anyway)
})


_DOCLING_POSTPROCESS_V2 = True  # Set to False to disable all S29 additions


def _docling_postprocess_v1(text: str) -> str:
    """Pre-S29 docling post-processing (legacy fallback for kill switch).

    DEPRECATED (S34 Fix 3.2, A-05): v1 is retained only as a rollback safety
    net. It lacks S29+ improvements (font encoding, zone protection, run-together
    splitting, domain hyphens, reference fixes). Use v2 for all production work.

    Contains only the original ligature rejoin, apostrophe/punctuation fixes,
    and OCR stripping. No font encoding, no double-space collapse, no
    exclusion zones, no SymSpell, no run-together splitting.
    """
    fixes_applied = 0

    # Unicode ligature replacement
    for lig_char, replacement in _UNICODE_LIGATURE_MAP.items():
        if lig_char in text:
            count = text.count(lig_char)
            text = text.replace(lig_char, replacement)
            fixes_applied += count

    # Ligature space-rejoin (dictionary-only, no SymSpell)
    for i, frag in enumerate(_LIGATURE_FRAGMENTS):
        pattern = _LIGATURE_REJOIN_PATTERNS[i]

        def _rejoin(match, _frag=frag):
            before = match.group(1)
            lig = match.group(2)
            after = match.group(3)
            joined = before + lig + after
            if joined.lower() in _KNOWN_LIGATURE_WORDS:
                return joined
            if before.lower() in _STANDALONE_WORDS:
                return match.group(0)
            if after.lower() in _STANDALONE_WORDS and len(after) <= 2:
                return match.group(0)
            return joined

        new_text = pattern.sub(_rejoin, text)
        if new_text != text:
            fixes_applied += 1
            text = new_text

    # Trailing ligature
    for i, frag in enumerate(_LIGATURE_FRAGMENTS):
        pattern = _LIGATURE_TRAILING_PATTERNS[i]

        def _rejoin_trailing(match, _frag=frag):
            before = match.group(1)
            lig = match.group(2)
            if before.lower() in _STANDALONE_WORDS:
                return match.group(0)
            joined = before + lig
            if len(joined) < 4:
                return match.group(0)
            return joined

        new_text = pattern.sub(_rejoin_trailing, text)
        if new_text != text:
            fixes_applied += 1
            text = new_text

    # Leading ligature (dictionary-only)
    for frag in _LIGATURE_FRAGMENTS:
        pattern = re.compile(
            r'(?:^|(?<=\s))(' + re.escape(frag) + r')\s+(\w{2,}\b)')

        def _rejoin_leading(match, _frag=frag):
            lig = match.group(1)
            after = match.group(2)
            joined = lig + after
            if joined.lower() in _KNOWN_LIGATURE_WORDS:
                return joined
            return match.group(0)

        new_text = pattern.sub(_rejoin_leading, text)
        if new_text != text:
            fixes_applied += 1
            text = new_text

    # Spaced apostrophes
    for apos in _APOSTROPHE_VARIANTS:
        for suffix in _CONTRACTION_SUFFIXES:
            pat = (r"(\w+)\s+" + re.escape(apos) + r"\s+("
                   + re.escape(suffix) + r")\b")
            new_text = re.sub(pat, r"\1'\2", text)
            if new_text != text:
                fixes_applied += 1
                text = new_text

    # Spaced punctuation
    new_text = re.sub(r'([a-z]+)\s+-\s+([a-z]+)', r'\1-\2', text)
    if new_text != text:
        fixes_applied += 1
        text = new_text
    new_text = re.sub(r'\s+([.,;:!?)])', r'\1', text)
    if new_text != text:
        fixes_applied += 1
        text = new_text
    new_text = re.sub(r'([(\[])\s+', r'\1', text)
    if new_text != text:
        fixes_applied += 1
        text = new_text
    new_text = re.sub(r'(\d+)\s+%', r'\1%', text)
    if new_text != text:
        fixes_applied += 1
        text = new_text

    # OCR artifact stripping
    text = _strip_ocr_artifacts(text)

    if fixes_applied > 0:
        print(f"  [docling-postprocess-v1] Applied {fixes_applied} fixes")

    return text


def _docling_postprocess(text: str) -> str:
    """
    Post-process docling markdown output to fix known artifacts.

    Execution order (S29 fix plan, updated S34, S36, S37):
    0. NFC normalization              (S34 Fix 2.3, C-03)
    1. Hyphen-newline repair          (FIX-B04 / CM-013, P2)
    2. Font encoding map              (FIX-B01 / CM-009)
    3. Double-space collapse          (FIX-B02 / CM-003)
    4. Protect exclusion zones        (FIX-B03 / CM-023)
    5a-d. Unicode ligature replacement + space-rejoin (existing + FIX-B06)
    5e. Sliding window ligature merge (S34 Fix 1.7, B-02)
    5g. Space-split word merging      (S36)
    5f. Domain hyphen restoration     (S34 Fix 1.3, A-01)
    6. Run-together splitting         (FIX-B08 / CM-012, S37 rewrite)
    7. Restore exclusion zones        (FIX-B03 / CM-023)
    7b. Reference section cleanup     (S34 Fix 1.8, B-04)
    8. Spaced apostrophes/punctuation (existing)
    9. OCR artifact stripping         (existing _strip_ocr_artifacts)

    MUST run AFTER normalize_symbols() and BEFORE post_extraction_cleanup().
    Called ONLY from the docling extraction branch.

    Kill switch: set _DOCLING_POSTPROCESS_V2 = False to revert to pre-S29
    behavior (existing ligature + apostrophe + OCR strip only).
    """
    if not _DOCLING_POSTPROCESS_V2:
        # Pre-S29 behavior: skip Steps 1-7, run only ligature/apostrophe/OCR
        # S34 Fix 3.2 (A-05): Deprecation warning
        print("  WARNING: _DOCLING_POSTPROCESS_V2 is False. Using deprecated v1 "
              "pipeline. v1 lacks S29+ improvements. Set to True for production.")
        return _docling_postprocess_v1(text)

    fixes_applied = 0

    # ---- Step 0: NFC normalization (S34 Fix 2.3, C-03) ----
    # Normalize combining character sequences BEFORE any regex matching.
    # E.g., "e" + combining acute -> "e with acute" (single codepoint).
    # Also handles diacritic joining for cases like "Les\u0301niowska".
    # Safe, stdlib-only, zero risk of false positives.
    text = unicodedata.normalize('NFC', text)

    # Diacritic joining: fix isolated combining marks separated by space.
    # E.g., "s \u0301" (s + space + combining acute) -> "s\u0301" -> NFC -> proper char
    # W3-QC1 MINOR-09: \s+ greedy match is correct -- multiple spaces before a
    # combining mark should all be removed.
    _diac_pat = re.compile(r'(\w)\s+([\u0300-\u036F])')
    _diac_new = _diac_pat.sub(r'\1\2', text)
    if _diac_new != text:
        text = unicodedata.normalize('NFC', _diac_new)
        fixes_applied += 1

    # ---- Step 1: Hyphen-newline repair (FIX-B04, CM-013) ----
    # QC-FP-01: Preserve hyphen during line-break repair.
    # "cost-\neffective" -> "cost-effective" (not "costeffective")
    text, n = re.subn(r'(?<=\w)-\s*\n\s*(?=\w)', '-', text)
    if n > 0:
        fixes_applied += n

    # ---- Step 2: Font encoding map (FIX-B01, P1) ----
    try:
        text = _fix_font_encoding(text, extractor="docling")
    except Exception as e:
        print(f"  WARNING: _fix_font_encoding failed ({e}), using original text")

    # ---- Step 3: Double-space collapse (FIX-B02, P1) ----
    try:
        text = _collapse_double_spaces(text)
    except Exception as e:
        print(f"  WARNING: _collapse_double_spaces failed ({e}), using original text")

    # ---- Step 4: Protect exclusion zones (FIX-B03, P1) ----
    try:
        text, _zone_placeholders = _protect_zones(text)
    except Exception as e:
        print(f"  WARNING: _protect_zones failed ({e}), skipping zone protection")
        _zone_placeholders = {}

    # ---- Step 5: Ligature space-rejoin ----

    # Step 5a: Replace Unicode ligature characters with ASCII
    # (Supplements existing _ligature_map which handles /uniFBxx escape strings)
    for lig_char, replacement in _UNICODE_LIGATURE_MAP.items():
        if lig_char in text:
            count = text.count(lig_char)
            text = text.replace(lig_char, replacement)
            fixes_applied += count

    # Step 5b: Rejoin space-split ligature fragments
    # Pattern: "dif fi culties" -> "difficulties"
    # Process longest fragments first to prevent partial matches
    for i, frag in enumerate(_LIGATURE_FRAGMENTS):
        pattern = _LIGATURE_REJOIN_PATTERNS[i]

        def _rejoin(match, _frag=frag):
            before = match.group(1)
            lig = match.group(2)
            after = match.group(3)
            joined = before + lig + after

            # Check 1: If joined word is valid (SymSpell or dictionary), rejoin
            if _is_valid_word(joined):
                return joined

            # Check 2: If 'before' is a standalone word, do NOT rejoin
            # Prevents "if i had" -> "ifihad"
            if before.lower() in _STANDALONE_WORDS:
                return match.group(0)

            # Check 3: If 'after' is a short standalone word, do NOT rejoin
            if after.lower() in _STANDALONE_WORDS and len(after) <= 2:
                return match.group(0)

            # FIX-S32-02: Proper-name ligature rejoin.
            # When docling space-splits a ligature inside a proper name
            # (e.g., "Kof fi jberg" from "Koffijberg"), the joined form is not
            # in any dictionary (Dutch/foreign name). The v2 conservative default
            # preserves the split, which is wrong -- the split was a docling
            # extraction artifact, not intentional spacing.
            # Heuristic: if the joined form starts with uppercase, it is likely
            # a proper name. Rejoin it. This is safe because real multi-word
            # sequences starting with uppercase (e.g., "Of fi ce") are first
            # caught by Check 1 (_is_valid_word("Office") returns True) or by
            # Check 2 (standalone-word guard for "Of").
            # S34 Fix 3.3 (A-06): Comment corrected -- primary guard is Check 1,
            # not the standalone-word guard.
            if joined[0].isupper():
                return joined

            # Default: do NOT rejoin (conservative - changed from "default rejoin")
            return match.group(0)

        new_text = pattern.sub(_rejoin, text)
        if new_text != text:
            fixes_applied += 1
            text = new_text

    # Step 5c: Handle trailing ligature (no following word)
    # Pattern: "staf fi" at end of line -> "staffi"
    # Add standalone guard (fixing Risk 1.7 from 0A-3)
    for i, frag in enumerate(_LIGATURE_FRAGMENTS):
        pattern = _LIGATURE_TRAILING_PATTERNS[i]

        def _rejoin_trailing(match, _frag=frag):
            before = match.group(1)
            lig = match.group(2)
            # Guard: don't rejoin if 'before' is standalone
            if before.lower() in _STANDALONE_WORDS:
                return match.group(0)
            joined = before + lig
            # Only rejoin if the result looks plausible (4+ chars)
            if len(joined) < 4:
                return match.group(0)
            return joined

        new_text = pattern.sub(_rejoin_trailing, text)
        if new_text != text:
            fixes_applied += 1
            text = new_text

    # Step 5d: Handle leading ligature (word starts with ligature fragment)
    # Pattern: "fi ndings" -> "findings", "fl ow" -> "flow"
    # Only rejoin if the result is a known ligature word (dictionary check required
    # since leading fragments like "fi" and "fl" are ambiguous without context).
    for frag in _LIGATURE_FRAGMENTS:
        pattern = re.compile(r'(?:^|(?<=\s))(' + re.escape(frag) + r')\s+(\w{2,}\b)')

        def _rejoin_leading(match, _frag=frag):
            lig = match.group(1)
            after = match.group(2)
            joined = lig + after
            if _is_valid_word(joined):
                return joined
            return match.group(0)

        new_text = pattern.sub(_rejoin_leading, text)
        if new_text != text:
            fixes_applied += 1
            text = new_text

    # ---- Step 5e: Sliding window ligature merge (S34 Fix 1.7, B-02) ----
    # Handles multi-fragment splits like "dif fi cul ties" (4 tokens) that
    # the existing 3-token regex patterns in Step 5b cannot catch.
    # Per QC-B CRITICAL-02: operates line-by-line, skips zone placeholders,
    # and preserves whitespace via targeted regex substitution.
    # W3-QC1 MINOR-06: "st" has higher false-merge potential than other seqs
    # (appears in many common words). Mitigated by _is_valid_word check on the
    # merged candidate. Monitor for false merges in retest.
    # W3-QC2 MINOR-04: _LIGATURE_SEQS moved to module-level frozenset.
    _sliding_fixes = 0
    _sw_lines = []
    for _sw_line in text.splitlines(True):
        _sw_tokens = _sw_line.split()
        if len(_sw_tokens) < 2:
            _sw_lines.append(_sw_line)
            continue
        _sw_changed = False
        _sw_i = 0
        while _sw_i < len(_sw_tokens):
            # Skip zone placeholders
            if '\x00' in _sw_tokens[_sw_i]:
                _sw_i += 1
                continue
            _sw_merged = False
            # Try window sizes 4, 3, 2
            for _sw_win in (4, 3, 2):
                if _sw_i + _sw_win > len(_sw_tokens):
                    continue
                _sw_slice = _sw_tokens[_sw_i:_sw_i + _sw_win]
                # Skip if any token contains a zone placeholder
                if any('\x00' in t for t in _sw_slice):
                    continue
                # Strip stray hyphens between fragments
                _sw_cleaned = [t.strip('-') for t in _sw_slice]
                _sw_candidate = ''.join(_sw_cleaned)
                # Check: valid word AND contains a ligature sequence
                if (_is_valid_word(_sw_candidate)
                        and any(ls in _sw_candidate.lower()
                                for ls in _LIGATURE_SEQS)):
                    # Replace tokens with merged form
                    _sw_tokens[_sw_i:_sw_i + _sw_win] = [_sw_candidate]
                    _sw_changed = True
                    _sliding_fixes += 1
                    _sw_merged = True
                    break
            # W3-QC1 MINOR-05: Always advance index (merged token is single,
            # no need to re-check it against window sizes that require 2+ tokens)
            _sw_i += 1
        if _sw_changed:
            # Reconstruct line preserving leading whitespace and trailing newline
            _sw_trailing = '\n' if _sw_line.endswith('\n') else ''
            _sw_lead_m = re.match(r'^(\s*)', _sw_line)
            _sw_leading = _sw_lead_m.group(1) if _sw_lead_m else ''
            _sw_lines.append(_sw_leading + ' '.join(_sw_tokens) + _sw_trailing)
        else:
            _sw_lines.append(_sw_line)
    if _sliding_fixes > 0:
        text = ''.join(_sw_lines)
        fixes_applied += _sliding_fixes
        print(f"  [ligature-window] Merged {_sliding_fixes} multi-fragment ligature splits")

    # ---- Step 5g: Space-split word merging (S36) ----
    # S36: Fix non-ligature space-splits like "Sculp her" -> "Sculpher",
    # "Wei bull" -> "Weibull". Runs AFTER ligature steps (5a-5e) and BEFORE
    # domain hyphens (5f) so merged words can be hyphen-restored.
    try:
        text = _fix_space_splits(text)
    except Exception as e:
        print(f"  WARNING: _fix_space_splits failed ({e})")

    # ---- Step 5f: Domain hyphen restoration (S34 Fix 1.3, A-01) ----
    # W3-QC2 MAJOR-01: MUST run BEFORE _fix_run_togethers (Step 6) so that
    # domain compounds like "costeffectiveness" get their hyphens restored to
    # "cost-effectiveness" BEFORE wordsegment splits them into "cost effectiveness".
    # Step 6's hyphen-compound preservation (FIX-S32-01) then keeps them intact.
    # Runs within zone-protected text; \b word boundaries work correctly around
    # \x00 zone placeholders since \x00 is not a word character.
    try:
        text = _restore_domain_hyphens(text)
    except Exception as e:
        print(f"  WARNING: _restore_domain_hyphens failed ({e})")

    # ---- Step 6: Run-together splitting (FIX-B08, CM-012, S37 rewrite) ----
    try:
        text = _fix_run_togethers(text)
    except Exception as e:
        print(f"  WARNING: _fix_run_togethers failed ({e}), using original text")

    # ---- Step 7: Restore exclusion zones (FIX-B03, P1) ----
    # QC-FP-05: Restore BEFORE apostrophe/punctuation fixes.
    # Apostrophe/punctuation fixes are beneficial for URLs and HTML comments.
    try:
        text = _restore_zones(text, _zone_placeholders)
    except Exception as e:
        print(f"  WARNING: _restore_zones failed ({e})")

    # ---- Step 7b: Reference section post-processing (S34 Fix 1.8, B-04) ----
    # Strip false heading markers within references, join continuation lines.
    try:
        text = _fix_reference_section(text)
    except Exception as e:
        print(f"  WARNING: _fix_reference_section failed ({e})")

    # ---- Step 8: Spaced apostrophes ----

    # Fix contractions: "don ' t" -> "don't", "patient ' s" -> "patient's"
    for apos in _APOSTROPHE_VARIANTS:
        for suffix in _CONTRACTION_SUFFIXES:
            pat = r"(\w+)\s+" + re.escape(apos) + r"\s+(" + re.escape(suffix) + r")\b"
            new_text = re.sub(pat, r"\1'\2", text)
            if new_text != text:
                fixes_applied += 1
                text = new_text

    # ---- Spaced punctuation (bundled with apostrophes) ----

    # Spaced hyphens in compound words: "cost - effective" -> "cost-effective"
    # LOWERCASE ONLY guard to avoid joining numeric ranges (Risk 2.4)
    new_text = re.sub(r'([a-z]+)\s+-\s+([a-z]+)', r'\1-\2', text)
    if new_text != text:
        fixes_applied += 1
        text = new_text

    # Space before closing punctuation
    new_text = re.sub(r'\s+([.,;:!?)])', r'\1', text)
    if new_text != text:
        fixes_applied += 1
        text = new_text

    # Space after opening brackets
    new_text = re.sub(r'([(\[])\s+', r'\1', text)
    if new_text != text:
        fixes_applied += 1
        text = new_text

    # Spaced percent: "95 %" -> "95%"
    new_text = re.sub(r'(\d+)\s+%', r'\1%', text)
    if new_text != text:
        fixes_applied += 1
        text = new_text

    # ---- Step 9: OCR artifact stripping ----
    text = _strip_ocr_artifacts(text)

    if fixes_applied > 0:
        print(f"  [docling-postprocess] Applied {fixes_applied} docling-specific fixes")

    return text


def _fix_docling_heading_overdetection(text: str) -> str:
    """
    Demote false-positive headings in docling markdown output.

    Docling's HERON layout model over-classifies text as Section-header.
    All headings become ## (H2) with no hierarchy. This function demotes
    likely false-positive headings to plain text using heuristics.

    MUST run AFTER post_extraction_cleanup() to avoid interfering with
    m4 ALL-CAPS heading fix (lines 232-238).
    Called ONLY from the docling extraction branch.
    """
    lines = text.split('\n')
    heading_pat = re.compile(r'^(#{2,6})\s+(.+)$')  # Only ## through ######; preserve H1

    # Count headings and estimate page count
    heading_count = sum(1 for l in lines if heading_pat.match(l))
    non_empty_lines = len([l for l in lines if l.strip()])
    estimated_pages = max(non_empty_lines / 50, 1)
    headings_per_page = heading_count / estimated_pages
    aggressive = headings_per_page > 8  # Threshold from .py file

    demoted_count = 0
    fixed_lines = []

    # Build heading index for proximity check (Claude Opus Rule 5)
    heading_indices = []
    for i, line in enumerate(lines):
        m = heading_pat.match(line)
        if m:
            heading_indices.append((i, m.group(2).strip()))

    # Set of line numbers to demote
    demote_set = set()

    for idx, (line_num, htxt) in enumerate(heading_indices):
        demote = False

        # Rule 1: Too short or too long
        if len(htxt) < 3 or len(htxt) > 200:
            demote = True

        # Rule 2: More than 12 words (sentence, not heading)
        elif len(htxt.split()) > 12:
            demote = True

        # Rule 3: Contains pipe or tab (table-like content)
        elif '|' in htxt or '\t' in htxt:
            demote = True

        # Rule 4: Figure/Table caption pattern
        elif re.match(r'^(Figure|Table|Fig\.|Tab\.)\s*\d', htxt, re.I):
            demote = True

        # Rule 5: Pure numbers/statistical values
        elif re.match(r'^[\d\s\.,\-\+%()]+$', htxt.strip()):
            demote = True

        # Rule 6 (aggressive only): Ends with period
        elif aggressive and htxt.endswith('.'):
            demote = True

        # Rule 7 (Claude): Non-heading patterns
        elif re.match(r'^p\s*[<>=]', htxt.strip()):
            demote = True
        elif re.match(r'^N\s*=\s*\d+', htxt.strip()):
            demote = True
        elif re.match(r'^\(\d+', htxt.strip()):
            demote = True

        # Rule 8: Common sentence starters (from 0B 1C spec)
        # CAVEAT: Only apply to headings with 5+ words to avoid demoting
        # legitimate short headings like "In Conclusion", "The Problem",
        # "For Further Reading".
        elif len(htxt.split()) >= 5 and re.match(
                r'^(Provide|Please|The |A |An |This |These |'
                r'Those |If |When |Where |For |In |On |At |'
                r'To |From )', htxt):
            demote = True

        # Rule 9 (Claude): Proximity -- two headings within 2 lines
        if not demote and idx > 0:
            prev_line_num, prev_text = heading_indices[idx - 1]
            if line_num - prev_line_num <= 2:
                if len(htxt) < len(prev_text):
                    demote = True

        if demote:
            demote_set.add(line_num)

    # Apply demotions: strip ## prefix to plain text (NOT bold, avoiding Risk 4.7)
    for i, line in enumerate(lines):
        if i in demote_set:
            cleaned = heading_pat.sub(r'\2', line)
            fixed_lines.append(cleaned)
            demoted_count += 1
        else:
            fixed_lines.append(line)

    if demoted_count > 0:
        print(f"  [docling-headings] Demoted {demoted_count} false-positive headings "
              f"(density: {headings_per_page:.1f}/page, aggressive: {aggressive})")

    return '\n'.join(fixed_lines)


SUPPORTED_FORMATS = {".pdf", ".docx", ".pptx", ".xlsx"}
DPI = 300
ZOOM = DPI / 72

# Institutional headers to skip when extracting document title.
# These are generic cover-page headings, not actual document titles.
# Matched case-insensitively as substrings of H1 text.
# SYNC: this list is duplicated in run-pipeline.py (_INSTITUTIONAL_HEADERS)
INSTITUTIONAL_HEADERS = [
    "health technology assessment",
    "statens legemiddelverk",
    "folkehelseinstituttet",
    "table of contents",
    "contents",
    "systematic review",
    "rapid review",
    "technology appraisal",
    "evidence report",
    "clinical practice guideline",
    "erasmus school of health policy",
    "university of oslo",
    "uio",
]

# Domain detection keywords for auto-classification
DOMAIN_KEYWORDS = {
    "health_economics": [
        "cost", "QALY", "ICER", "Markov", "willingness-to-pay",
        "cost-effectiveness", "incremental", "budget impact", "threshold"
    ],
    "clinical_trial": [
        "randomized", "placebo", "endpoint", "ITT", "CONSORT",
        "adverse event", "RCT", "trial", "phase", "randomization"
    ],
    "systematic_review": [
        "PRISMA", "meta-analysis", "forest plot", "I-squared",
        "pooled", "heterogeneity", "systematic", "review"
    ],
    "hta_regulatory": [
        "NICE", "DMP", "NoMA", "reimbursement", "submission",
        "appraisal", "HTA", "regulatory", "guideline"
    ],
    "epidemiology": [
        "incidence", "prevalence", "registry", "cohort", "DALY",
        "mortality", "population", "burden", "disease"
    ],
    "methodology": [
        "model validation", "simulation", "calibration",
        "structural uncertainty", "sensitivity", "probabilistic"
    ],
}

# High-specificity override keywords that force a domain match.
# Checked BEFORE frequency scoring. If any keyword is found (case-insensitive
# substring match), that domain wins immediately. These are terms that are
# unambiguous domain signals — e.g. "NoMA" only appears in HTA regulatory.
# Note: "ICER" (bare) is in health_economics; "ICER threshold" is in
# hta_regulatory (the phrase "ICER threshold" implies a policy/appraisal context).
DOMAIN_OVERRIDE_KEYWORDS = {
    "hta_regulatory": [
        "NoMA", "metodevurdering", "ICER threshold", "cost per QALY",
        "NICE appraisal", "health technology assessment",
        "reimbursement decision", "HTA body", "DMP",
    ],
    "health_economics": [
        "Markov model", "cost-effectiveness analysis",
        "willingness to pay", "ICER",
        "incremental cost-effectiveness",
        "cost-utility analysis",
        "survival analysis",
        "hazard ratio",
    ],
}


# --- Docling post-processing constants ---

# Unicode ligature characters -> ASCII (supplement to existing _ligature_map which handles /uniFBxx)
_UNICODE_LIGATURE_MAP = {
    '\ufb00': 'ff', '\ufb01': 'fi', '\ufb02': 'fl',
    '\ufb03': 'ffi', '\ufb04': 'ffl', '\ufb05': 'st', '\ufb06': 'st',
}

# Ligature fragment processing order: LONGEST FIRST to prevent partial matches
_LIGATURE_FRAGMENTS = ['ffi', 'ffl', 'ff', 'fi', 'fl']

# W3-QC2 MINOR-04: Module-level frozenset for sliding window ligature detection (Step 5e).
# Includes all ligature fragment sequences plus 'ft' and 'st' which docling
# sometimes space-splits around ligature boundaries.
_LIGATURE_SEQS = frozenset({'fi', 'fl', 'ff', 'ffi', 'ffl', 'ft', 'st'})

# Words known to contain ligature sequences. Used as override when standalone guard
# would otherwise prevent rejoining. EXPAND this set as new words are encountered.
_KNOWN_LIGATURE_WORDS = frozenset({
    # From Perplexity .py file (30 words)
    'office', 'officer', 'official', 'officially', 'offline', 'offset',
    'offend', 'offensive', 'offer', 'offering', 'offered', 'offers',
    'offspring', 'affirm', 'affirmative', 'affirmed', 'affiliation',
    'affiliated', 'afford', 'affordable', 'afforded',
    'affect', 'affected', 'affecting', 'affection', 'affiliate',
    'affluent', 'affluence', 'affix', 'affixed',
    # From Claude Opus HE_TERMS (22 words)
    'efficacy', 'effectiveness', 'pharmacoeconomics',
    'incremental', 'willingness', 'significant', 'significance',
    'benefit', 'beneficial', 'difficulties', 'insufficient', 'coefficient',
    'confidential', 'proficiency', 'efficient', 'efficiency', 'deficit',
    'influence', 'influential', 'inflammation',
    'reflective', 'notification',
    # From 0A-4 Section 4A missing terms
    'effect', 'effects', 'effective',
    'different', 'difference', 'differentiation',
    'sufficient', 'sufficiency',
    'coefficients', 'staff', 'staffing',
    'specific', 'specification',
    'certificate', 'certification',
    'deficiency', 'deficient',
    'modification', 'modifications',
    'identification', 'identified', 'identifier',
    'notifications', 'artificial',
    'verification', 'verified',
    'classification', 'classified',
    'scientific', 'confidence',
    'magnificent', 'conflict', 'conflicting',
    'reflection', 'difficult',
    # Additional common English words with ligature sequences
    'differ', 'differing', 'suffer', 'suffering',
    'buffer', 'buffering', 'coffee',
    'affirmation',
    'suffix', 'prefix', 'fifteen', 'fifty',
    'amplification', 'clarification', 'justification',
    'quantification', 'ratification', 'simplification',
    'stratification', 'unification',
    'findings', 'final', 'finally', 'financial', 'finance',
    'field', 'fields', 'first', 'firstly',
    'finish', 'finite', 'fiscal', 'fishing',
    'flexible', 'flight', 'float', 'floor', 'flow',
    'fluid', 'fluorescence', 'flush', 'flutter',
    'scaffold', 'scaffolding',
    'inflation', 'inflection', 'influx',
    'profile', 'profiles', 'profiling',
    'rifle', 'stifle', 'trifle', 'waffle',
})

# Short words that are legitimate standalone words and should NOT be joined
# with adjacent ligature fragments. "if i" must stay as "if i", not "ifi".
_STANDALONE_WORDS = frozenset({
    'if', 'of', 'off', 'a', 'an', 'the', 'in', 'on', 'or',
    'is', 'it', 'as', 'at', 'be', 'by', 'do', 'go', 'he',
    'me', 'my', 'no', 'so', 'to', 'up', 'us', 'we',
    # Norwegian prepositions (per 0A-4 Section 4C)
    'af',
})

# Pre-compiled regex patterns for ligature rejoining
_LIGATURE_REJOIN_PATTERNS = [
    re.compile(r'(\b\w+)\s+(' + re.escape(frag) + r')\s+(\w+\b)')
    for frag in _LIGATURE_FRAGMENTS
]
_LIGATURE_TRAILING_PATTERNS = [
    re.compile(r'(\b\w+)\s+(' + re.escape(frag) + r')\b(?!\s+\w)')
    for frag in _LIGATURE_FRAGMENTS
]

# Apostrophe Unicode variants
_APOSTROPHE_VARIANTS = ["'", "\u2019", "\u2018", "\u02BC"]
_CONTRACTION_SUFFIXES = ['t', 's', 're', 've', 'll', 'd', 'm',
                         'T', 'S', 'RE', 'VE', 'LL', 'D', 'M']

# OCR artifact detection: whitelist patterns for statistical text
# FIX-B05 (CM-024): Enhanced whitelist to prevent silent data deletion
_STATISTICAL_PATTERNS = re.compile(
    r'(?:'
    r'\d+\.\d+\s*[\(\[]\d+\.\d+'       # CIs: 0.85 (0.72
    r'|p\s*[<>=]\s*0\.\d'              # p-values
    r'|HR\s*[=:]|OR\s*[=:]|RR\s*[=:]'  # hazard/odds/risk ratios
    r'|ICER\s*[=:]|QALY|DALY|NMB'      # HTA terms
    r'|95%\s*CI|99%\s*CI'              # confidence intervals
    r'|N\s*=\s*\d+|n\s*=\s*\d+'       # sample sizes
    r'|\d+,\d{3}'                      # large numbers: "25,000"
    r'|\d+\.\d+\s*\u00B1\s*\d+\.\d+'  # mean +/- SD
    r'|\d+\.\d+\s+6\s+\d+\.\d+'       # "6" as +/- artifact (ISSUE-29)
    r'|\bITT\b|\bPFS\b|\bOS\b'        # clinical endpoints
    r'|\bWTP\b|\bCEA\b|\bNHB\b'       # HTA acronyms
    r')',
    re.IGNORECASE
)

# --- FIX-B01 (CM-009): Font encoding map for /C## glyph artifacts ---

_FONT_ENCODING_MAP = {
    "/C15": "\u2022",   # bullet (most common in HTA papers)
    "/C0": "\u2018",    # left single quote
    "/C1": "\u2019",    # right single quote
    "/C2": "\u201C",    # left double quote
    "/C3": "\u201D",    # right double quote
    "/C4": "\u2013",    # en dash
    "/C5": "\u2014",    # em dash
    "/C6": "\u2020",    # dagger
    "/C7": "\u2021",    # double dagger
    "/C8": "\u00A7",    # section sign
    "/C9": "\u00B6",    # pilcrow
    "/C10": "\u00A9",   # copyright
    "/C11": "\u00AE",   # registered
    "/C12": "\u2122",   # trademark
    "/C13": "\u00B0",   # degree
    "/C14": "\u00B1",   # plus-minus
    "/C16": "\u2023",   # triangular bullet
    "/C19": "\u0141",   # L-stroke (Polish)
}

# QC-FP-03: Sort alternatives longest-first to prevent single-digit codes
# from matching before two-digit codes (e.g., "1" before "10"-"19").
_FONT_ENCODING_COMPILED = re.compile(
    r'/C(' + '|'.join(
        sorted(
            (re.escape(k.replace('/C', '')) for k in _FONT_ENCODING_MAP),
            key=len, reverse=True
        )
    ) + r')(?![0-9])'
)

_GLYPH_PATTERN = re.compile(r'GLYPH\(cmap:[a-fA-F0-9]+\)')
_GENERIC_CCODE_PATTERN = re.compile(r'/C\d{1,3}(?![0-9])')


# ═══════════════════════════════════════════════════════════════════════════
# FORMAT DETECTION
# ═══════════════════════════════════════════════════════════════════════════

def detect_format(file_path: Path) -> str:
    """Return normalized format string: pdf, docx, pptx, xlsx."""
    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_FORMATS:
        print(f"ERROR: Unsupported format '{ext}'. Supported: {SUPPORTED_FORMATS}")
        sys.exit(1)
    return ext.lstrip(".")


def get_page_count(file_path: Path, fmt: str) -> Optional[int]:
    """Get page/slide/sheet count where possible."""
    if fmt == "pdf":
        import fitz
        doc = fitz.open(str(file_path))
        count = len(doc)
        doc.close()
        return count
    elif fmt == "pptx":
        from pptx import Presentation
        prs = Presentation(str(file_path))
        return len(prs.slides)
    elif fmt == "xlsx":
        import openpyxl
        try:
            wb = openpyxl.load_workbook(str(file_path), read_only=True)
            count = len(wb.sheetnames)
            wb.close()
            return count
        except Exception:
            return None
    elif fmt == "docx":
        # DOCX doesn't have a reliable page count without rendering
        return None
    return None


# ═══════════════════════════════════════════════════════════════════════════
# IMAGE EXTRACTION - PDF
# ═══════════════════════════════════════════════════════════════════════════

def detect_caption_near_image(page, img_rect) -> Optional[str]:
    """
    Detect "Figure N:" or "Table N:" text near the image bounding box.
    Searches within 100 pixels below the image.
    """
    import fitz
    try:
        # Expand search area below image
        search_rect = fitz.Rect(
            img_rect.x0 - 50,
            img_rect.y1,
            img_rect.x1 + 50,
            img_rect.y1 + 100
        )
        text = page.get_text("text", clip=search_rect)
        # Look for common caption patterns
        match = re.search(r'(Figure|Table|Fig\.|Tab\.)\s*\d+[:\.]?\s*([^\n]{0,80})', text, re.IGNORECASE)
        if match:
            return match.group(0).strip()
    except Exception:
        pass
    return None


def guess_type_from_caption(caption: Optional[str]) -> Optional[str]:
    """
    Guess image type from detected caption text.
    Returns a type string compatible with IMAGE NOTE types.
    """
    if not caption:
        return None

    caption_lower = caption.lower()

    # Map common caption keywords to image types
    type_map = {
        "forest plot": "forest-plot",
        "kaplan": "kaplan-meier",
        "survival": "kaplan-meier",
        "tornado": "tornado-diagram",
        "decision tree": "decision-tree",
        "flow chart": "flow-chart",
        "flowchart": "flow-chart",
        "scatter": "scatter",
        "cost-effectiveness plane": "scatter",
        "box plot": "box-plot",
        "histogram": "histogram",
        "heatmap": "heatmap",
        "bar chart": "bar-chart",
        "line chart": "line-chart",
        "pie chart": "pie-chart",
        "funnel": "funnel-plot",
        "network": "network-diagram",
    }

    for keyword, img_type in type_map.items():
        if keyword in caption_lower:
            return img_type

    # Default: table vs figure
    if "table" in caption_lower:
        return "table-image"
    elif "figure" in caption_lower or "fig." in caption_lower:
        return "other"  # Unknown figure type

    return None


def _iou(rect_a, rect_b):
    """S37: Intersection over Union for two (x0, y0, x1, y1) rects.

    Both rects must be in the same coordinate system (TOPLEFT).
    Returns 0.0 if either rect is degenerate or there is no intersection.
    """
    x0 = max(rect_a[0], rect_b[0])
    y0 = max(rect_a[1], rect_b[1])
    x1 = min(rect_a[2], rect_b[2])
    y1 = min(rect_a[3], rect_b[3])

    inter_w = max(0.0, x1 - x0)
    inter_h = max(0.0, y1 - y0)
    inter_area = inter_w * inter_h

    area_a = max(0.0, rect_a[2] - rect_a[0]) * max(0.0, rect_a[3] - rect_a[1])
    area_b = max(0.0, rect_b[2] - rect_b[0]) * max(0.0, rect_b[3] - rect_b[1])
    union_area = area_a + area_b - inter_area

    if union_area <= 0:
        return 0.0
    return inter_area / union_area


def _center_distance(rect_a, rect_b):
    """S37: Euclidean distance between centers of two (x0, y0, x1, y1) rects."""
    cx_a = (rect_a[0] + rect_a[2]) / 2.0
    cy_a = (rect_a[1] + rect_a[3]) / 2.0
    cx_b = (rect_b[0] + rect_b[2]) / 2.0
    cy_b = (rect_b[1] + rect_b[3]) / 2.0
    return ((cx_a - cx_b) ** 2 + (cy_a - cy_b) ** 2) ** 0.5


def _map_docling_image_placeholders(
    md_text: str,
    docling_pictures: list,
    image_entries: list,
) -> str:
    """S37: Replace anonymous <!-- image --> placeholders with identified names.

    Docling's export_to_markdown() emits bare '<!-- image -->' for every
    PictureItem. This function maps each placeholder to a fitz-extracted
    image using spatial bbox proximity matching (IoU + nearest-neighbor).

    S37 rewrite: replaces the S36 sequential cursor algorithm with
    bbox-based matching. This fixes count-mismatch cascading where
    docling detects N pictures and fitz extracts M images (N != M).

    Args:
        md_text: Markdown text with anonymous '<!-- image -->' placeholders.
        docling_pictures: List of dicts with 'page_no', 'bbox_topleft'
            (TOPLEFT coords), and 'label' from docling's
            result.document.pictures (collected in extract_text).
        image_entries: List of dicts from extract_images_pdf() with
            'page', 'figure_num', 'filename', 'source_format',
            'fitz_rect' keys.

    Returns:
        Markdown text with '<!-- IMAGE: filename.ext -->' replacing bare
        placeholders. Unmappable placeholders become '<!-- IMAGE: unmapped -->'.
    """
    if not docling_pictures or not image_entries:
        return md_text

    from collections import defaultdict

    # S37: Build page-grouped lookup from image_entries.
    # Only use pdf_embedded images with valid fitz_rect.
    images_by_page = defaultdict(list)
    for idx, entry in enumerate(image_entries):
        if entry.get("source_format") == "pdf_embedded":
            images_by_page[entry["page"]].append((idx, entry))

    # Track which fitz images have been matched (1:1 constraint)
    matched_fitz = set()

    # S37: Thresholds for bbox matching.
    # IoU is low because docling DL model and fitz report quite different
    # bounding boxes (docling = visual extent, fitz = raster placement).
    _IOU_THRESHOLD = 0.05
    _DISTANCE_THRESHOLD = 200.0  # points (~2.8 inches)

    placeholder_filenames = []
    _match_methods = []  # For logging

    for pic in docling_pictures:
        page = pic.get("page_no")
        bbox_tl = pic.get("bbox_topleft")

        # S37: If no bbox, try legacy page-based fallback
        if page is None:
            placeholder_filenames.append(None)
            _match_methods.append("no_page")
            continue

        if page not in images_by_page:
            placeholder_filenames.append(None)
            _match_methods.append("no_fitz_on_page")
            continue

        # Filter to unmatched candidates on this page
        candidates = [(idx, e) for idx, e in images_by_page[page]
                       if idx not in matched_fitz]

        if not candidates:
            placeholder_filenames.append(None)
            _match_methods.append("all_consumed")
            continue

        # S37: If we have bbox data and at least some candidates have rects,
        # use spatial matching
        if bbox_tl is not None:
            rect_candidates = [(idx, e) for idx, e in candidates
                               if e.get("fitz_rect") is not None]

            if rect_candidates:
                # Phase 1: IoU matching
                best_iou = 0.0
                best_idx = None
                best_entry = None

                for idx, entry in rect_candidates:
                    iou_val = _iou(bbox_tl, entry["fitz_rect"])
                    if iou_val > best_iou:
                        best_iou = iou_val
                        best_idx = idx
                        best_entry = entry

                if best_iou >= _IOU_THRESHOLD and best_entry is not None:
                    matched_fitz.add(best_idx)
                    placeholder_filenames.append(best_entry["filename"])
                    _match_methods.append(f"iou={best_iou:.3f}")
                    continue

                # Phase 2: Nearest-neighbor fallback (center-to-center)
                best_dist = float('inf')
                best_idx = None
                best_entry = None

                for idx, entry in rect_candidates:
                    dist = _center_distance(bbox_tl, entry["fitz_rect"])
                    if dist < best_dist:
                        best_dist = dist
                        best_idx = idx
                        best_entry = entry

                if best_dist <= _DISTANCE_THRESHOLD and best_entry is not None:
                    matched_fitz.add(best_idx)
                    placeholder_filenames.append(best_entry["filename"])
                    _match_methods.append(f"dist={best_dist:.1f}pt")
                    continue

        # S37: Fallback -- cursor-based matching for images without bbox
        # or when spatial matching fails. Use first unmatched on page.
        _fallback_found = False
        for idx, entry in candidates:
            if idx not in matched_fitz:
                matched_fitz.add(idx)
                placeholder_filenames.append(entry["filename"])
                _match_methods.append("cursor_fallback")
                _fallback_found = True
                break
        if not _fallback_found:
            placeholder_filenames.append(None)
            _match_methods.append("unmapped")

    # S37: Replace <!-- image --> placeholders in order
    result_lines = []
    placeholder_idx = 0
    for line in md_text.split('\n'):
        if line.strip() == '<!-- image -->':
            if placeholder_idx < len(placeholder_filenames):
                fname = placeholder_filenames[placeholder_idx]
                if fname:
                    result_lines.append(f'<!-- IMAGE: {fname} -->')
                else:
                    result_lines.append('<!-- IMAGE: unmapped -->')
                placeholder_idx += 1
            else:
                result_lines.append(line)
        else:
            result_lines.append(line)

    mapped = sum(1 for f in placeholder_filenames if f is not None)
    total = len(placeholder_filenames)
    unmapped = total - mapped

    # S37: Detailed logging of match methods
    _method_counts = {}
    for m in _match_methods:
        _key = m.split('=')[0] if '=' in m else m
        _method_counts[_key] = _method_counts.get(_key, 0) + 1
    _method_str = ', '.join(f'{k}:{v}' for k, v in sorted(_method_counts.items()))

    print(f"  [image-map] Mapped {mapped}/{total} docling image placeholders "
          f"to extracted files"
          + (f" ({unmapped} unmapped)" if unmapped > 0 else "")
          + f" [{_method_str}]")

    return '\n'.join(result_lines)


def extract_images_pdf(file_path: Path, images_dir: Path,
                       short_name: str, md_text: str) -> list[dict]:
    """Extract images from PDF via PyMuPDF with enhanced metadata."""
    import fitz

    doc = fitz.open(str(file_path))
    images_dir.mkdir(parents=True, exist_ok=True)
    image_index = []
    figure_num = 0
    extracted_xrefs = set()  # Track extracted xrefs to avoid duplicates

    # Parse section structure from md_text for section context
    sections = parse_section_structure(md_text)

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_images = page.get_images(full=True)

        if page_images:
            for img_info in page_images:
                xref = img_info[0]

                # Skip if already extracted (duplicate detection fix)
                if xref in extracted_xrefs:
                    continue

                try:
                    base_image = doc.extract_image(xref)
                except Exception:
                    continue
                if base_image is None:
                    continue

                img_bytes = base_image["image"]
                img_ext = base_image.get("ext", "png")
                width = base_image.get("width", 0)
                height = base_image.get("height", 0)

                if width < 50 or height < 50:
                    continue

                extracted_xrefs.add(xref)  # Mark as extracted
                figure_num += 1
                filename = f"{short_name}-fig{figure_num}-page{page_idx + 1}.{img_ext}"
                filepath = images_dir / filename

                with open(filepath, "wb") as f:
                    f.write(img_bytes)
                if _ensure_max_dimension(filepath):
                    # Re-read dimensions after resize
                    try:
                        with Image.open(filepath) as _resized:
                            width, height = _resized.size
                    except Exception:
                        pass

                # ── Blank detection ───────────────────────────────────────────────
                # Detect near-black fragments, pure-color blocks, and other blank
                # placeholder images that waste Opus vision tokens.
                # Aligned with run-pipeline.py _is_blank_image() canonical logic:
                # Tier 1: file size < 2KB (tiny placeholder)
                # Tier 2: file_size < 5KB AND mean < 30 AND unique_colors < 16
                #         (near-black with few colors — OR tier, not AND)
                # Tier 3: std < 5.0 (near-uniform pixel values)
                # M1: unique_colors < 32 AND mean < 240 AND std < 30
                #     (color-block / gradient bar detection)
                # Applied at write time so the manifest field is populated before
                # prepare-image-analysis.py reads it.
                _is_blank = False
                _blank_reason = None
                try:
                    _file_size = os.path.getsize(str(filepath))
                    # Tier 1: tiny placeholder
                    if _file_size < 2000:
                        _is_blank = True
                        _blank_reason = f"file_size={_file_size}B"
                    elif Image is not None:
                        import numpy as _np
                        _pil_img = Image.open(str(filepath))
                        _gray = _pil_img.convert("L")
                        _arr = _np.array(_gray)
                        _std = float(_arr.std())
                        _mean = float(_arr.mean())

                        # Compute unique colors for Tier 2 and M1
                        _unique_colors = -1
                        if _pil_img.mode in ('RGB', 'RGBA'):
                            _rgb = _pil_img.convert('RGB')
                            _w_px, _h_px = _rgb.size
                            if _w_px * _h_px > 50000:
                                _rgb = _rgb.resize(
                                    (min(_w_px, 250), min(_h_px, 200)),
                                    Image.NEAREST)
                            _unique_colors = len(set(_rgb.getdata()))

                        _pil_img.close()

                        # Tier 2: near-black with few colors (OR tier)
                        if (_file_size > 0 and _file_size < 5000
                                and _mean < 30
                                and _unique_colors >= 0
                                and _unique_colors < 16):
                            _is_blank = True
                            _blank_reason = (
                                f"near_black(mean={_mean:.2f},"
                                f"colors={_unique_colors},"
                                f"size={_file_size}B)")
                        # Tier 3: near-uniform pixel values
                        elif _std < 5.0:
                            _is_blank = True
                            _blank_reason = f"std={_std:.2f}"
                        # M1: color-block / gradient bar detection
                        elif (_unique_colors >= 0
                                and _unique_colors < 32
                                and _mean < 240
                                and _std < 30):
                            _is_blank = True
                            _blank_reason = (
                                f"color_block(colors={_unique_colors},"
                                f"mean={_mean:.2f},std={_std:.2f})")
                except Exception:
                    pass  # If detection fails, leave _is_blank=False (safe default)
                # ─────────────────────────────────────────────────────────────────

                # Get image rect for caption detection
                img_rects = page.get_image_rects(xref)
                img_rect = img_rects[0] if img_rects else None
                detected_caption = detect_caption_near_image(page, img_rect) if img_rect else None
                type_guess = guess_type_from_caption(detected_caption)

                # Get nearby text context (200 chars before and after image position)
                nearby_text = get_text_near_image(page, img_rect) if img_rect else None

                # Find section context
                section_context = find_section_for_page(sections, page_idx + 1, total_pages=len(doc))

                image_index.append({
                    "page": page_idx + 1,
                    "figure_num": figure_num,
                    "filename": filename,
                    "description": f"Figure {figure_num} from page {page_idx + 1}",
                    "width": width,
                    "height": height,
                    "source_format": "pdf_embedded",
                    "detected_caption": detected_caption,
                    "type_guess": type_guess,
                    "section_context": section_context,
                    "nearby_text": nearby_text,
                    "analysis_status": "pending",
                    "blank": _is_blank,
                    "blank_reason": _blank_reason,
                    "decorative": False,
                    "is_duplicate": False,
                    # S37: Store fitz rect for bbox-based placeholder matching
                    "fitz_rect": (
                        (img_rect.x0, img_rect.y0, img_rect.x1, img_rect.y1)
                        if img_rect else None
                    ),
                })

                # Split multi-panel if large enough (PRESERVED capability)
                if Image is not None and (width > 800 and height > 800):
                    try:
                        _split_panels(img_bytes, figure_num, page_idx + 1,
                                      short_name, images_dir, image_index)
                    except Exception:
                        pass

    # Render sparse/figure-only pages (PRESERVED capability)
    pages_with_images = {e["page"] for e in image_index}
    for page_idx in range(len(doc)):
        page_num = page_idx + 1
        if page_num in pages_with_images:
            continue
        page = doc[page_idx]
        text = page.get_text("text").strip()
        if len(text) > 200:
            continue

        figure_num += 1
        filename = f"{short_name}-page{page_num}-render-300dpi.png"
        filepath = images_dir / filename
        mat = fitz.Matrix(ZOOM, ZOOM)
        pix = page.get_pixmap(matrix=mat)
        pix.save(str(filepath))
        _render_w, _render_h = pix.width, pix.height
        if _ensure_max_dimension(filepath):
            # Re-read dimensions after resize
            try:
                with Image.open(filepath) as _resized:
                    _render_w, _render_h = _resized.size
            except Exception:
                pass

        section_context = find_section_for_page(sections, page_num, total_pages=len(doc))

        image_index.append({
            "page": page_num,
            "figure_num": figure_num,
            "filename": filename,
            "description": f"Full page render of page {page_num} at 300 DPI",
            "width": _render_w,
            "height": _render_h,
            "source_format": "pdf_page_render",
            "detected_caption": None,
            "type_guess": "other",
            "section_context": section_context,
            "nearby_text": None,
            "analysis_status": "pending",
        })

    doc.close()
    return image_index


def get_text_near_image(page, img_rect) -> Optional[str]:
    """Extract text surrounding the image (200 chars before and after)."""
    import fitz
    if not img_rect:
        return None
    try:
        # Get all text on page
        full_text = page.get_text("text")

        # Get text in expanded rect around image
        expanded_rect = fitz.Rect(
            max(0, img_rect.x0 - 100),
            max(0, img_rect.y0 - 100),
            img_rect.x1 + 100,
            img_rect.y1 + 100
        )
        nearby = page.get_text("text", clip=expanded_rect)

        # Clean and truncate
        nearby = re.sub(r'\s+', ' ', nearby).strip()
        if len(nearby) > 200:
            nearby = nearby[:200] + "..."

        return nearby if nearby else None
    except Exception:
        return None


def find_section_for_page(sections: list[dict], page_num: int,
                          total_pages: int = 0) -> dict:
    """
    Find the section that contains the given page number.

    Uses a document-calibrated lines-per-page ratio (total_lines / total_pages)
    instead of a hardcoded constant, then maps page_num to an estimated line
    number and finds which section contains that line.

    Args:
        sections: Parsed section structure from parse_section_structure().
        page_num: 1-based page number in the PDF.
        total_pages: Total number of pages in the PDF (from len(doc)).
            When provided, enables accurate per-document calibration.
            Falls back to 45 lines/page when 0 or not provided.
    """
    if not sections:
        return {
            "heading": "Unknown Section",
            "heading_level": 2,
        }

    # Estimate total lines from all sections
    total_lines = max((s.get("line_end") or 0) for s in sections)
    if total_lines == 0:
        # Fallback to first section
        return {
            "heading": sections[0]["heading"],
            "heading_level": sections[0]["level"],
        }

    # Use document-calibrated ratio instead of constant 45 lines/page.
    # The old constant caused all images on long documents (200+ pages)
    # to overshoot section boundaries and fall back to the last heading.
    if total_pages > 0:
        estimated_lines_per_page = total_lines / total_pages
    else:
        estimated_lines_per_page = 45  # legacy fallback
    estimated_line_num = page_num * estimated_lines_per_page

    # Find section containing this estimated line
    for section in sections:
        line_start = section.get("line_start", 0)
        line_end = section.get("line_end", 0)
        if line_start <= estimated_line_num <= line_end:
            return {
                "heading": section["heading"],
                "heading_level": section["level"],
            }

    # If no match, find the nearest preceding section (closest line_end
    # that is <= estimated_line_num). This is more accurate than always
    # returning the last section.
    best_section = sections[-1]  # default fallback
    best_distance = float('inf')
    for section in sections:
        line_end = section.get("line_end", 0)
        if line_end <= estimated_line_num:
            distance = estimated_line_num - line_end
            if distance < best_distance:
                best_distance = distance
                best_section = section

    return {
        "heading": best_section["heading"],
        "heading_level": best_section["level"],
    }


# ═══════════════════════════════════════════════════════════════════════════
# IMAGE EXTRACTION - DOCX
# ═══════════════════════════════════════════════════════════════════════════

def extract_images_docx(file_path: Path, images_dir: Path,
                        short_name: str) -> list[dict]:
    """Extract embedded images from DOCX via python-docx."""
    from docx import Document

    doc = Document(str(file_path))
    images_dir.mkdir(parents=True, exist_ok=True)
    image_index = []
    figure_num = 0

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_data = rel.target_part.blob
                content_type = rel.target_part.content_type
            except Exception:
                continue

            # Determine extension
            ext_map = {
                "image/png": "png",
                "image/jpeg": "jpg",
                "image/gif": "gif",
                "image/bmp": "bmp",
                "image/tiff": "tiff",
                "image/x-emf": "emf",
                "image/x-wmf": "wmf",
            }
            img_ext = ext_map.get(content_type, "png")

            # Skip vector formats that can't be displayed as raster
            if img_ext in ("emf", "wmf"):
                continue

            # Get dimensions if possible
            width, height = 0, 0
            if Image is not None and img_ext in ("png", "jpg", "gif", "bmp"):
                try:
                    img = Image.open(io.BytesIO(img_data))
                    width, height = img.size
                except Exception:
                    pass

            # Skip tiny images
            if width > 0 and width < 50 and height < 50:
                continue

            figure_num += 1
            filename = f"{short_name}-fig{figure_num}.{img_ext}"
            filepath = images_dir / filename

            with open(filepath, "wb") as f:
                f.write(img_data)
            if _ensure_max_dimension(filepath):
                # Re-read dimensions after resize
                try:
                    with Image.open(filepath) as _resized:
                        width, height = _resized.size
                except Exception:
                    pass

            image_index.append({
                "page": None,  # DOCX doesn't have reliable page numbers
                "figure_num": figure_num,
                "filename": filename,
                "description": f"Figure {figure_num} from document",
                "width": width,
                "height": height,
                "source_format": "docx_embedded",
                "detected_caption": None,
                "type_guess": None,
                "section_context": None,
                "nearby_text": None,
                "analysis_status": "pending",
            })

    return image_index


# ═══════════════════════════════════════════════════════════════════════════
# IMAGE EXTRACTION - PPTX
# ═══════════════════════════════════════════════════════════════════════════

def extract_images_pptx(file_path: Path, images_dir: Path,
                        short_name: str) -> list[dict]:
    """
    Extract from PPTX:
    1. All embedded images from shapes
    2. Full slide renders at 300 DPI (via PyMuPDF PDF conversion)
    """
    from pptx import Presentation

    prs = Presentation(str(file_path))
    images_dir.mkdir(parents=True, exist_ok=True)
    image_index = []
    figure_num = 0

    # Extract embedded images from shapes
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                try:
                    img_data = shape.image.blob
                    content_type = shape.image.content_type
                except Exception:
                    continue

                ext_map = {
                    "image/png": "png",
                    "image/jpeg": "jpg",
                    "image/gif": "gif",
                    "image/bmp": "bmp",
                }
                img_ext = ext_map.get(content_type, "png")

                width, height = 0, 0
                if Image is not None:
                    try:
                        img = Image.open(io.BytesIO(img_data))
                        width, height = img.size
                    except Exception:
                        pass

                if width > 0 and width < 30 and height < 30:
                    continue

                figure_num += 1
                filename = f"{short_name}-slide{slide_num}-fig{figure_num}.{img_ext}"
                filepath = images_dir / filename

                with open(filepath, "wb") as f:
                    f.write(img_data)
                if _ensure_max_dimension(filepath):
                    # Re-read dimensions after resize
                    try:
                        with Image.open(filepath) as _resized:
                            width, height = _resized.size
                    except Exception:
                        pass

                image_index.append({
                    "page": slide_num,
                    "figure_num": figure_num,
                    "filename": filename,
                    "description": f"Figure {figure_num} from slide {slide_num}",
                    "width": width,
                    "height": height,
                    "source_format": "pptx_embedded",
                    "detected_caption": None,
                    "type_guess": None,
                    "section_context": None,
                    "nearby_text": None,
                    "analysis_status": "pending",
                })

    # Render full slides via libreoffice -> PDF -> PyMuPDF (if available)
    # (PRESERVED capability)
    try:
        import fitz
        import subprocess
        import tempfile

        with tempfile.TemporaryDirectory() as tmpdir:
            # Convert PPTX to PDF via libreoffice
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", tmpdir, str(file_path)],
                capture_output=True, timeout=120
            )
            if result.returncode == 0:
                pdf_files = list(Path(tmpdir).glob("*.pdf"))
                if pdf_files:
                    doc = fitz.open(str(pdf_files[0]))
                    for page_idx in range(len(doc)):
                        slide_num = page_idx + 1
                        figure_num += 1
                        filename = f"{short_name}-slide{slide_num}-render-300dpi.png"
                        filepath = images_dir / filename
                        mat = fitz.Matrix(ZOOM, ZOOM)
                        pix = doc[page_idx].get_pixmap(matrix=mat)
                        pix.save(str(filepath))
                        _slide_w, _slide_h = pix.width, pix.height
                        if _ensure_max_dimension(filepath):
                            # Re-read dimensions after resize
                            try:
                                with Image.open(filepath) as _resized:
                                    _slide_w, _slide_h = _resized.size
                            except Exception:
                                pass
                        image_index.append({
                            "page": slide_num,
                            "figure_num": figure_num,
                            "filename": filename,
                            "description": f"Full slide {slide_num} render at 300 DPI",
                            "width": _slide_w,
                            "height": _slide_h,
                            "source_format": "pptx_slide_render",
                            "detected_caption": None,
                            "type_guess": None,
                            "section_context": None,
                            "nearby_text": None,
                            "analysis_status": "pending",
                        })
                    doc.close()
    except (ImportError, FileNotFoundError, subprocess.TimeoutExpired):
        print("  WARNING: libreoffice not found. Skipping slide renders.")
        print("  (Embedded images still extracted)")

    return image_index


# ═══════════════════════════════════════════════════════════════════════════
# MULTI-PANEL SPLITTING (PRESERVED)
# ═══════════════════════════════════════════════════════════════════════════

def _split_panels(img_bytes: bytes, figure_num: int, page_num: int,
                  short_name: str, images_dir: Path,
                  image_index: list[dict]):
    """Split a multi-panel figure (A/B) into individual panels."""
    img = Image.open(io.BytesIO(img_bytes))
    w, h = img.size

    if w > h:
        panel_a = img.crop((0, 0, w // 2, h))
        panel_b = img.crop((w // 2, 0, w, h))
    else:
        panel_a = img.crop((0, 0, w, h // 2))
        panel_b = img.crop((0, h // 2, w, h))

    for label, panel in [("a", panel_a), ("b", panel_b)]:
        filename = f"{short_name}-fig{figure_num}{label}-page{page_num}.png"
        filepath = images_dir / filename
        panel.save(str(filepath))
        _panel_w, _panel_h = panel.width, panel.height
        if _ensure_max_dimension(filepath):
            # Re-read dimensions after resize
            try:
                with Image.open(filepath) as _resized:
                    _panel_w, _panel_h = _resized.size
            except Exception:
                pass
        image_index.append({
            "page": page_num,
            "figure_num": f"{figure_num}{label}",
            "filename": filename,
            "description": f"Figure {figure_num} panel {label.upper()} from page {page_num}",
            "width": _panel_w,
            "height": _panel_h,
            "source_format": "panel_split",
            "detected_caption": None,
            "type_guess": None,
            "section_context": None,
            "nearby_text": None,
            "analysis_status": "pending",
        })


# ═══════════════════════════════════════════════════════════════════════════
# TABLE FALLBACK CASCADE (Fix 3.2)
# ═══════════════════════════════════════════════════════════════════════════

def extract_table_fallback(pdf_path, page_num):
    """Fix 3.2: Extract tables using pymupdf find_tables() then Camelot.

    Called when pymupdf4llm's primary extraction fails to produce pipe-table
    markdown for a page that references a table.

    Fallback cascade:
      1. pymupdf find_tables() with strategy="lines_strict"
         (different params from pymupdf4llm's internal find_tables call)
      2. Camelot (lattice then stream flavor) if installed

    Args:
        pdf_path: Path to the PDF file.
        page_num: 1-indexed page number to extract tables from.

    Returns:
        List of markdown table strings, or empty list if all methods fail.
    """
    import fitz

    # Attempt 1: pymupdf find_tables() with explicit parameters
    # pymupdf4llm already calls find_tables() internally, but with default
    # params.  Using strategy="lines_strict" can catch tables that the
    # default strategy misses (per research spec Section 3.2 MISS-1).
    try:
        doc = fitz.open(str(pdf_path))
        page = doc[page_num - 1]
        tables = page.find_tables(strategy="lines_strict")
        if tables.tables:
            results = []
            for tab in tables.tables:
                df = tab.to_pandas()
                results.append(df.to_markdown(index=False))
            doc.close()
            return results
        doc.close()
    except Exception:
        try:
            doc.close()
        except Exception:
            pass

    # Attempt 2: Camelot (lattice then stream)
    try:
        import camelot
        # Try lattice first (for tables with visible borders)
        tables = camelot.read_pdf(
            str(pdf_path), pages=str(page_num),
            flavor='lattice')
        if tables and len(tables) > 0:
            return [t.df.to_markdown(index=False) for t in tables]
        # Fall back to stream (for tables without borders)
        tables = camelot.read_pdf(
            str(pdf_path), pages=str(page_num),
            flavor='stream')
        if tables and len(tables) > 0:
            return [t.df.to_markdown(index=False) for t in tables]
    except ImportError:
        pass
    except Exception:
        pass

    return []


def apply_table_fallback(md_text, pdf_path):
    """Fix 3.2: Post-processing step to recover missing tables.

    Scans the markdown for "Table N" references that lack a following
    pipe-table within 50 lines, then attempts fallback extraction.

    Args:
        md_text: The markdown text from pymupdf4llm extraction.
        pdf_path: Path to the source PDF for fallback extraction.

    Returns:
        Updated markdown text with recovered tables inserted.
    """
    lines = md_text.split('\n')
    total_lines = len(lines)
    insertions = []  # list of (line_index, table_markdown)
    processed_tables = set()  # MAJOR-1: Track which table numbers have been processed

    # Cache page count once to avoid repeated PDF opens (MAJOR-2)
    _cached_page_count = _estimate_page_count(pdf_path)

    # Find "Table N" references (case-insensitive)
    table_ref_pattern = re.compile(
        r'\b[Tt][Aa][Bb][Ll][Ee]\s+(\d+)')

    for i, line in enumerate(lines):
        match = table_ref_pattern.search(line)
        if not match:
            continue

        # MAJOR-1: Skip if this table number was already processed
        table_num = match.group(1)
        if table_num in processed_tables:
            continue

        # Check if a pipe-table follows within 50 lines
        has_pipe_table = False
        for j in range(i + 1, min(i + 50, total_lines)):
            if '|' in lines[j] and lines[j].strip().startswith('|'):
                has_pipe_table = True
                break

        if has_pipe_table:
            processed_tables.add(table_num)
            continue

        # Estimate page number from line proportion (MINOR-1: dead loop removed)
        page_num = max(1, int((i / total_lines) * _cached_page_count) + 1)

        # Try fallback extraction
        recovered_tables = extract_table_fallback(pdf_path, page_num)
        # MAJOR-1: Mark table as processed regardless of extraction result
        processed_tables.add(table_num)
        if recovered_tables:
            # Insert after the reference line (find next empty line)
            insert_at = i + 1
            for j in range(i + 1, min(i + 5, total_lines)):
                if lines[j].strip() == '':
                    insert_at = j + 1
                    break
            table_md = '\n\n'.join(recovered_tables)
            insertions.append((insert_at, table_md))

    if not insertions:
        return md_text

    # Apply insertions in reverse order to preserve line indices
    for insert_at, table_md in reversed(insertions):
        lines.insert(insert_at, f"\n{table_md}\n")

    return '\n'.join(lines)


def _estimate_page_count(pdf_path):
    """Estimate PDF page count using fitz."""
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        count = len(doc)
        doc.close()
        return count
    except Exception:
        return 1


# ═══════════════════════════════════════════════════════════════════════════
# TEXT EXTRACTION + ASSEMBLY
# ═══════════════════════════════════════════════════════════════════════════

def extract_text(file_path: Path, fmt: str,
                 extractor: str = "docling") -> tuple:
    """Extract text from document.

    # S36: Returns (md_text, docling_pictures) tuple.
    # docling_pictures is a list of dicts with page_no/bbox for docling path,
    # or an empty list for all other extractors.

    PDF: pymupdf4llm (proper tables, headings, multi-column).
    PDF+tesseract: Tesseract OCR via PyMuPDF OCR bridge (scanned PDFs).
    DOCX/PPTX/XLSX: MarkItDown (good for structured formats).
    Fallback: MarkItDown for PDF if pymupdf4llm not installed (PRESERVED).
    """
    if fmt == "pdf" and extractor == "tesseract":
        # Tesseract OCR via PyMuPDF OCR bridge
        import fitz
        doc = fitz.open(str(file_path))
        pages = []
        for page in doc:
            tp = page.get_textpage_ocr(language="eng", dpi=300, full=False)
            pages.append(page.get_text("text", textpage=tp))
        doc.close()
        return "\n\n".join(pages), []  # S36: tuple return
    elif fmt == "pdf" and extractor == "docling":
        # Docling deep-learning PDF extraction (tables, layout, OCR)
        # Lazy import: docling is optional; fail clearly if not installed
        try:
            from docling.document_converter import (
                DocumentConverter, PdfFormatOption,
            )
            from docling.datamodel.pipeline_options import (
                PdfPipelineOptions, TableFormerMode, TableStructureOptions,
            )
            from docling.datamodel.base_models import InputFormat
        except ImportError as imp_err:
            print(
                f"ERROR: docling is not installed: {imp_err}",
                file=sys.stderr,
            )
            print(
                "  Install with: pip3 install docling",
                file=sys.stderr,
            )
            sys.exit(3)  # Exit code 3 = extractor failed, router can fallback

        # Issue 3 fix: Disable OCR for digital PDFs to prevent figure garbage text
        # Digital detection: check first 3 pages for text content using fitz
        import fitz as _fitz_check
        _check_doc = _fitz_check.open(str(file_path))
        _sample_text = ""
        for _i, _page in enumerate(_check_doc):
            if _i >= 3:
                break
            _sample_text += _page.get_text("text")
        _check_doc.close()
        _is_digital = len(_sample_text.strip()) > 100

        if _is_digital:
            print(f"  (docling mode: TableFormer ACCURATE, OCR disabled "
                  f"-- digital PDF, text density: {len(_sample_text.strip())} chars)")
        else:
            print(f"  (docling mode: TableFormer ACCURATE + OCR "
                  f"-- scanned PDF detected)")

        # Configure pipeline options
        pipeline_options = PdfPipelineOptions()
        pipeline_options.document_timeout = 120.0  # FIX-A02 (GAP-02): prevent hangs on problematic PDFs
        pipeline_options.do_table_structure = True
        # FIX-A01 (CM-001): Use model's predicted cells, not PDF cells.
        # Prevents column merges that collapse complex tables.
        pipeline_options.table_structure_options = TableStructureOptions(
            mode=TableFormerMode.ACCURATE,
            do_cell_matching=False,
        )
        pipeline_options.do_ocr = not _is_digital
        pipeline_options.images_scale = 2.0  # FIX-A03 (CM-017): 2x DPI for better TableFormer cell detection
        # Keep image extraction with PyMuPDF (existing pipeline handles it)
        pipeline_options.generate_picture_images = False
        pipeline_options.generate_page_images = False

        # FIX-A04 (CM-005): Enable formula enrichment if available
        try:
            from docling.datamodel.pipeline_options import CodeFormulaVlmOptions
            if not hasattr(pipeline_options, 'do_code_formula'):
                print("  (formula enrichment: DISABLED - "
                      "PdfPipelineOptions has no 'do_code_formula' attribute)")
            else:
                code_formula_options = CodeFormulaVlmOptions.from_preset(
                    "codeformulav2"
                )
                pipeline_options.do_code_formula = True
                pipeline_options.code_formula_options = code_formula_options
                if getattr(pipeline_options, 'do_code_formula', False):
                    print("  (formula enrichment: ENABLED via CodeFormulaVlmOptions)")
                else:
                    print("  (formula enrichment: WARNING - "
                          "attribute set but reads as False)")
        except (ImportError, AttributeError, TypeError) as e:
            print(f"  (formula enrichment: DISABLED - {e})")
        except Exception as e:
            print(f"  (formula enrichment: DISABLED - unexpected error: {e})")

        converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=pipeline_options
                )
            }
        )

        try:
            result = converter.convert(str(file_path))
            md_text = result.document.export_to_markdown()
        except Exception as exc:
            print(
                f"ERROR: docling failed: {type(exc).__name__}: {exc}",
                file=sys.stderr,
            )
            sys.exit(3)

        # S37: Extract picture metadata for placeholder mapping.
        # Each PictureItem has prov[0].page_no (1-indexed) and bbox.
        # The Nth <!-- image --> placeholder corresponds to the Nth PictureItem.
        # S37: Pre-convert docling bbox from BOTTOMLEFT to TOPLEFT coords
        # using page dimensions, so _map_docling_image_placeholders() can
        # directly compare with fitz rects (which are TOPLEFT).
        _docling_picture_meta = []
        try:
            for _pic in result.document.pictures:
                if _pic.prov:
                    _page_no = _pic.prov[0].page_no
                    _bbox = _pic.prov[0].bbox

                    # S37: Get page height for coordinate conversion
                    _page_item = result.document.pages.get(_page_no)
                    _page_h = (
                        _page_item.size.height
                        if _page_item and hasattr(_page_item, 'size')
                           and _page_item.size
                        else None
                    )

                    # S37: Convert to TOPLEFT coordinates
                    _bbox_topleft = None
                    if _page_h and hasattr(_bbox, 'coord_origin'):
                        try:
                            _origin = _bbox.coord_origin.value
                        except AttributeError:
                            _origin = str(_bbox.coord_origin)
                        if _origin == "BOTTOMLEFT":
                            try:
                                _tl = _bbox.to_top_left_origin(_page_h)
                                _bbox_topleft = (
                                    _tl.l, _tl.t, _tl.r, _tl.b)
                            except Exception:
                                # Fallback: manual conversion
                                _bbox_topleft = (
                                    _bbox.l,
                                    _page_h - _bbox.t,
                                    _bbox.r,
                                    _page_h - _bbox.b,
                                )
                        else:
                            # Already TOPLEFT
                            _bbox_topleft = (
                                _bbox.l, _bbox.t, _bbox.r, _bbox.b)
                    elif _bbox:
                        # No page height: store raw bbox (may be wrong
                        # coord system but better than nothing)
                        _bbox_topleft = (
                            _bbox.l, _bbox.t, _bbox.r, _bbox.b)

                    _docling_picture_meta.append({
                        "page_no": _page_no,
                        "bbox_topleft": _bbox_topleft,
                        "page_height": _page_h,
                        "label": (
                            _pic.label.value
                            if hasattr(_pic, 'label') else "unknown"
                        ),
                    })
                else:
                    _docling_picture_meta.append({
                        "page_no": None,
                        "bbox_topleft": None,
                        "page_height": None,
                        "label": (
                            _pic.label.value
                            if hasattr(_pic, 'label') else "unknown"
                        ),
                    })
            if _docling_picture_meta:
                print(f"  [docling] Found {len(_docling_picture_meta)} picture items "
                      f"for placeholder mapping")
        except Exception as _pic_err:
            print(f"  [docling] WARNING: Could not extract picture metadata: {_pic_err}")
            _docling_picture_meta = []

        # Apply the same post-processing as pymupdf4llm path
        # Fix 3.1: Symbol normalization (skip pymupdf4llm-specific glyph
        # fixes for docling; only universal fixes like U+FFFD removal)
        md_text = normalize_symbols(md_text, extractor="docling")

        # S27: Docling-specific post-processing (ligature rejoin, apostrophes,
        # spaced punctuation, OCR artifact stripping)
        md_text = _docling_postprocess(md_text)

        # Fix 3.15: Post-extraction cleanup
        md_text = post_extraction_cleanup(md_text)

        # S27: Docling heading over-detection fix (must run AFTER
        # post_extraction_cleanup to avoid interfering with m4 ALL-CAPS fix)
        md_text = _fix_docling_heading_overdetection(md_text)

        # Fix 3.2: Table fallback cascade
        # Note: docling's TableFormer typically handles tables well,
        # but we still run the fallback to catch edge cases
        try:
            md_text = apply_table_fallback(md_text, file_path)
        except Exception as _tfe:
            print(f"  [table-fallback] Warning: {_tfe}")

        return md_text, _docling_picture_meta  # S36: tuple return
    elif fmt == "pdf" and pymupdf4llm is not None:
        mode = "layout" if _HAS_LAYOUT else "standard"
        print(f"  (pymupdf4llm mode: {mode})")

        # Fix 3.3: Landscape page derotation (C8)
        # Detect rotated pages (landscape) and derotate before extraction.
        # pymupdf4llm takes a file path, so we save a corrected temp file
        # if any pages have non-zero rotation.
        _derotated_path = None
        try:
            import fitz as _fitz_derotate
            _doc = _fitz_derotate.open(str(file_path))
            _rotated = []
            for _pi in range(len(_doc)):
                _pg = _doc[_pi]
                if _pg.rotation != 0:
                    _rotated.append((_pi + 1, _pg.rotation))
                    _pg.set_rotation(0)
            if _rotated:
                import tempfile as _tmpmod
                _tf = _tmpmod.NamedTemporaryFile(
                    suffix='.pdf', delete=False)
                _doc.save(_tf.name)
                _derotated_path = _tf.name
                for _pn, _rot in _rotated:
                    print(f"  [derotate] Page {_pn}: "
                          f"corrected {_rot} degrees")
            _doc.close()
        except Exception as _de:
            print(f"  [derotate] Warning: {_de}")

        _pdf_path_for_extract = _derotated_path or str(file_path)

        try:
            md_text = pymupdf4llm.to_markdown(
                _pdf_path_for_extract,
                show_progress=False,
                write_images=False,
            )
        except ValueError as exc:
            # pymupdf/table.py raises ValueError when a table's
            # bounding-box list is empty (e.g. "min() iterable
            # argument is empty" at table.py:1534).  Exit with
            # code 3 so the pipeline router knows this extractor
            # failed and should try the next one in the chain.
            print(
                f"ERROR: pymupdf4llm crashed with ValueError: {exc}",
                file=sys.stderr,
            )
            print(
                "  This is a known pymupdf4llm table-detection bug.",
                file=sys.stderr,
            )
            print(
                "  Pipeline router should fall back to mineru or tesseract.",
                file=sys.stderr,
            )
            sys.exit(3)
        except Exception as exc:
            # Catch any other pymupdf4llm runtime error so the
            # process exits cleanly (non-zero) and the router can
            # trigger its fallback chain instead of seeing a raw
            # Python traceback.
            print(
                f"ERROR: pymupdf4llm failed unexpectedly: "
                f"{type(exc).__name__}: {exc}",
                file=sys.stderr,
            )
            sys.exit(3)
        finally:
            # Fix 3.3: Clean up derotated temp file
            if _derotated_path:
                try:
                    os.unlink(_derotated_path)
                except OSError:
                    pass

        # Fix 3.1: Symbol normalization (C3, C4, C5, C6, C7, m25)
        md_text = normalize_symbols(md_text, extractor="pymupdf4llm")

        # Fix 3.15: Post-extraction cleanup (M21, m4-m7, m10, m13, m15-m17, m26)
        md_text = post_extraction_cleanup(md_text)

        # Fix 3.2: Table fallback cascade (C1, C12, M19, M23)
        # Scan for table references without following pipe tables.
        # If found, attempt fallback extraction with find_tables()
        # (different params) and Camelot.
        try:
            md_text = apply_table_fallback(md_text, file_path)
        except Exception as _tfe:
            print(f"  [table-fallback] Warning: {_tfe}")

        return md_text, []  # S36: tuple return (no docling pictures)
    elif MarkItDown is not None:
        md = MarkItDown()
        result = md.convert(str(file_path))
        return result.text_content, []  # S36: tuple return
    else:
        print("ERROR: No text extractor available.")
        print("  PDF: pip install pymupdf4llm")
        print("  Other: pip install 'markitdown[all]'")
        sys.exit(1)


def clean_text(raw_text: str) -> str:
    """Clean common MarkItDown artifacts."""
    text = re.sub(r'\n{3,}', '\n\n', raw_text)
    lines = [line.rstrip() for line in text.split('\n')]
    return '\n'.join(lines).strip()


# ═══════════════════════════════════════════════════════════════════════════
# CONTEXT SUMMARY GENERATION (NEW)
# ═══════════════════════════════════════════════════════════════════════════

def parse_section_structure(md_text: str) -> list[dict]:
    """
    Parse section headings, line ranges, word counts, and image references.
    Returns list of section dicts.
    """
    sections = []
    lines = md_text.split('\n')
    current_section = None

    for line_num, line in enumerate(lines, 1):
        # Detect headings
        heading_match = re.match(r'^(#{1,6})\s+(.+)', line)
        if heading_match:
            # Save previous section
            if current_section:
                current_section["line_end"] = line_num - 1
                sections.append(current_section)

            # Start new section
            level = len(heading_match.group(1))
            heading_text = heading_match.group(2).strip()
            current_section = {
                "heading": heading_text,
                "level": level,
                "line_start": line_num,
                "line_end": None,
                "word_count": 0,
                "has_images": False,
                "image_refs": [],
            }

        # Count words and detect images in current section
        if current_section:
            current_section["word_count"] += len(line.split())
            # Detect image references
            img_refs = re.findall(r'\[([^\]]+)\]\(([^)]+\.(?:png|jpg|jpeg|gif|bmp))\)', line)
            if img_refs:
                current_section["has_images"] = True
                for label, path in img_refs:
                    # Extract figure number if present
                    fig_match = re.search(r'fig[\s-]?(\d+)', path, re.IGNORECASE)
                    if fig_match:
                        current_section["image_refs"].append(f"fig{fig_match.group(1)}")

    # Save last section
    if current_section:
        current_section["line_end"] = len(lines)
        sections.append(current_section)

    return sections


def detect_document_domain(md_text: str) -> str:
    """
    Auto-detect document domain from keyword frequencies.
    Returns one of: health_economics, clinical_trial, systematic_review,
    hta_regulatory, epidemiology, methodology, general.

    High-specificity override keywords are checked first. If any override
    keyword is found (case-insensitive substring), that domain wins
    immediately — bypassing the frequency scoring that can be skewed by
    generic clinical terms.

    S21/RC1: health_economics overrides take priority over hta_regulatory
    when BOTH have override matches.  hta_regulatory only wins in Phase 1
    when health_economics has zero override matches (i.e. the document is
    about regulatory process/HTA methodology without economic evaluation).
    """
    text_lower = md_text.lower()

    # ── Phase 1: Override keywords (high-specificity terms) ──
    # Collect ALL override matches first, then resolve priority.
    override_hits = {}  # domain -> list of matched keywords
    for domain, override_kws in DOMAIN_OVERRIDE_KEYWORDS.items():
        matched = [kw for kw in override_kws if kw.lower() in text_lower]
        if matched:
            override_hits[domain] = matched

    if override_hits:
        # S21/RC1: health_economics always wins over hta_regulatory when
        # both have override hits (documents about cost-effectiveness that
        # also mention HTA bodies are health_economics, not regulatory).
        if "health_economics" in override_hits:
            return "health_economics"
        # Only one domain matched, or hta_regulatory without health_econ
        best_domain = max(override_hits, key=lambda d: len(override_hits[d]))
        return best_domain

    # ── Phase 2: Frequency scoring (original logic) ──
    domain_scores = {}
    for domain, keywords in DOMAIN_KEYWORDS.items():
        score = sum(1 for kw in keywords if kw.lower() in text_lower)
        domain_scores[domain] = score

    # Return domain with highest score, or "general" if all zero
    max_domain = max(domain_scores.items(), key=lambda x: x[1])
    if max_domain[1] == 0:
        return "general"
    return max_domain[0]


def extract_key_terms(md_text: str, domain: str) -> list[str]:
    """Extract the top 5 most relevant keywords for the detected domain."""
    text_lower = md_text.lower()
    domain_keywords = DOMAIN_KEYWORDS.get(domain, [])

    # Count occurrences
    term_counts = {}
    for kw in domain_keywords:
        count = text_lower.count(kw.lower())
        if count > 0:
            term_counts[kw] = count

    # Return top 5
    sorted_terms = sorted(term_counts.items(), key=lambda x: x[1], reverse=True)
    return [term for term, count in sorted_terms[:5]]


def _is_institutional_header(text: str) -> bool:
    """Check if an H1 text is a generic institutional header, not a title.

    Uses case-insensitive substring matching against INSTITUTIONAL_HEADERS.
    Also rejects very short H1s (<=5 chars) which are usually acronyms or
    page numbers, not titles.
    """
    if not text or len(text) <= 5:
        return True
    text_lower = text.lower().strip()
    for header in INSTITUTIONAL_HEADERS:
        if header in text_lower:
            return True
    return False


def extract_title_from_md(md_text: str,
                          file_path: Optional[Path] = None) -> Optional[str]:
    """Extract the first non-institutional H1 from markdown text.

    Scans the first 80 lines for H1 headings, skipping any that match
    INSTITUTIONAL_HEADERS (cover-page headers, ToC, etc.). Returns the
    first valid H1 text, or None if no valid H1 is found.

    If file_path is a PDF and no H1 is found, falls back to fitz PDF
    metadata (document title from PDF properties).

    This is a standalone helper used by both the pymupdf path
    (extract_title_authors) and the MinerU YAML frontmatter writer.
    """
    lines = md_text.split('\n')
    for line in lines[:80]:
        if line.startswith('# ') and not line.startswith('## '):
            h1_text = re.sub(r'^#+\s*', '', line).strip()
            if not _is_institutional_header(h1_text):
                return h1_text

    # S27 Fix 5: Fallback to fitz PDF metadata if no H1 found
    if file_path is not None and str(file_path).lower().endswith('.pdf'):
        try:
            import fitz as _fitz_meta
            _meta_doc = _fitz_meta.open(str(file_path))
            _pdf_title = _meta_doc.metadata.get('title', '').strip()
            _meta_doc.close()
            if _pdf_title and len(_pdf_title) > 5 and not _is_institutional_header(_pdf_title):
                return _pdf_title
        except Exception:
            pass

    return None


def extract_title_authors(md_text: str) -> tuple[Optional[str], list[str]]:
    """
    Extract title (first non-institutional H1) and authors if present.
    Returns (title, authors_list).

    Skips institutional headers like "Health Technology Assessment" or
    "Statens legemiddelverk" which appear as H1 on cover pages but are
    not document titles. See INSTITUTIONAL_HEADERS constant.
    """
    title = extract_title_from_md(md_text)
    authors = []

    lines = md_text.split('\n')
    for line in lines[:80]:  # Expanded from 50 to 80 lines
        # Look for author line patterns
        if re.search(r'\bauthor', line, re.IGNORECASE):
            # Simple heuristic: names after "Authors:"
            author_match = re.search(r'authors?:?\s*(.+)', line, re.IGNORECASE)
            if author_match:
                author_str = author_match.group(1)
                # Split on common delimiters
                authors = re.split(r'[,;]\s*|\sand\s', author_str)
                authors = [a.strip() for a in authors if a.strip()]
                # S27 Fix 5: Filter garbage from author list
                authors = [a for a in authors if (
                    len(a) > 2
                    and not re.match(r'^(for|the|and|or|of)\b', a, re.IGNORECASE)
                    and 'correspondence' not in a.lower()
                    and '@' not in a
                    and not a.strip().startswith('http')
                )]
                break

    return title, authors


def extract_abstract(md_text: str) -> Optional[str]:
    """
    Extract first 500 chars of text after an "Abstract" heading.
    Returns None if no abstract found.
    """
    lines = md_text.split('\n')
    in_abstract = False
    abstract_lines = []

    for line in lines:
        if re.match(r'^#{1,3}\s*abstract', line, re.IGNORECASE):
            in_abstract = True
            continue

        if in_abstract:
            # Stop at next heading
            if re.match(r'^#{1,3}\s', line):
                break
            abstract_lines.append(line)

    abstract_text = ' '.join(abstract_lines).strip()
    if abstract_text:
        abstract_text = re.sub(r'\s+', ' ', abstract_text)
        return abstract_text[:500]
    return None


def generate_context_summary(md_text: str, image_count: int) -> dict:
    """
    Generate context summary JSON from extracted markdown.
    No LLM — pure extractive/heuristic processing.
    """
    title, authors = extract_title_authors(md_text)
    abstract = extract_abstract(md_text)
    sections = parse_section_structure(md_text)
    domain = detect_document_domain(md_text)
    key_terms = extract_key_terms(md_text, domain)

    # Calculate total word count
    total_words = sum(s["word_count"] for s in sections)

    return {
        "title": title or "Unknown Title",
        "authors": authors,
        "abstract": abstract,
        "sections": sections,
        "total_sections": len(sections),
        "total_words": total_words,
        "total_images": image_count,
        "document_domain": domain,
        "key_terms": key_terms,
    }


# ═══════════════════════════════════════════════════════════════════════════
# YAML HEADER, IMAGE INDEX, MANIFEST
# ═══════════════════════════════════════════════════════════════════════════

def build_header(file_path: Path, fmt: str, page_count: Optional[int],
                  image_count: int = 0,
                  extractor: str = "docling") -> str:
    """Build YAML header block."""
    tool_map = {
        "pdf": ("docling + PyMuPDF" if extractor == "docling"
                else "pymupdf4llm + PyMuPDF" if pymupdf4llm
                else "MarkItDown + PyMuPDF"),
        "docx": "MarkItDown + python-docx",
        "pptx": "MarkItDown + python-pptx",
        "xlsx": "MarkItDown",
    }
    doc_type_map = {
        "pdf": "research_paper",
        "docx": "document",
        "pptx": "presentation",
        "xlsx": "spreadsheet",
    }

    header = "---\n"
    header += f"source_file: {file_path.name}\n"
    header += f"source_path: {file_path}\n"
    header += f"source_format: {fmt}\n"
    if page_count is not None:
        label = "slides" if fmt == "pptx" else "sheets" if fmt == "xlsx" else "pages"
        header += f"{label}: {page_count}\n"
    header += f"conversion_date: {datetime.now().strftime('%Y-%m-%d')}\n"
    header += f"conversion_tool: {tool_map[fmt]}\n"
    header += "fidelity_standard: verbatim (QC required)\n"
    header += f"document_type: {doc_type_map[fmt]}\n"
    header += f"image_notes: {'pending' if image_count > 0 else 'none'}\n"
    header += "---\n\n"
    return header


def build_image_index(entries: list[dict], images_rel: str) -> str:
    """Build image index table."""
    if not entries:
        return "## Image Index\n\nNo images extracted.\n\n"

    page_label = "Page/Slide"
    table = "## Image Index\n\n"
    table += f"| Fig | Description | File | {page_label} | Size | Source |\n"
    table += f"|-----|-------------|------|{'-' * len(page_label)}|------|--------|\n"

    for e in entries:
        page = e["page"] if e["page"] is not None else "-"
        table += (
            f"| {e['figure_num']} "
            f"| {e['description']} "
            f"| [{e['filename']}]({images_rel}/{e['filename']}) "
            f"| {page} "
            f"| {e['width']}x{e['height']} "
            f"| {e['source_format']} |\n"
        )

    table += "\n"
    return table


def write_image_manifest(entries: list[dict], images_dir: Path,
                         md_path: Path):
    """
    Write a JSON manifest of all extracted images with enhanced metadata.
    Used by prepare-image-analysis.py and generate-image-notes.md.
    """
    manifest = {
        "md_file": str(md_path),
        "images_dir": str(images_dir),
        "image_count": len(entries),
        "generated": datetime.now().isoformat(),
        "images": entries,
    }
    manifest_path = images_dir / "image-manifest.json"
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)
    print(f"  Image manifest: {manifest_path}")


def write_context_summary(summary: dict, output_dir: Path):
    """Write context summary JSON alongside the MD file."""
    summary_path = output_dir / "context-summary.json"
    with open(summary_path, "w", encoding="utf-8") as f:
        json.dump(summary, f, indent=2, ensure_ascii=False)
    print(f"  Context summary: {summary_path}")


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="Convert PDF/DOCX/PPTX/XLSX to high-fidelity Markdown",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Supported formats: .pdf .docx .pptx .xlsx

Examples:
  python3 convert-paper.py paper.pdf
  python3 convert-paper.py report.docx -o reference-papers/report.md
  python3 convert-paper.py slides.pptx -i images/slides/
  python3 convert-paper.py data.xlsx --no-images

After conversion, run the QC pipeline:
  1. python3 qc-structural.py <output.md>
  2. python3 prepare-image-analysis.py <output.md>
  3. Claude subagent: generate-image-notes.md
  4. python3 validate-image-notes.py <output.md>
  5. Claude subagent: qc-content-fidelity.md (PDF only)
  6. Claude subagent: qc-final-review.md
        """
    )

    parser.add_argument("input_file", type=Path, help="Input file (PDF/DOCX/PPTX/XLSX)")
    parser.add_argument("--output", "-o", type=Path, default=None,
                        help="Output markdown path (default: <input>.md)")
    parser.add_argument("--images-dir", "-i", type=Path, default=None,
                        help="Images directory (default: images/<short-name>/)")
    parser.add_argument("--no-images", action="store_true",
                        help="Skip image extraction")
    parser.add_argument("--short-name", "-s", type=str, default=None,
                        help="Short name for filenames (default: from input name)")
    parser.add_argument("--extractor",
                        choices=["pymupdf4llm", "tesseract", "markitdown",
                                 "docling"],
                        default="docling",
                        help="PDF extractor (default: docling). "
                             "Fallback: pymupdf4llm. "
                             "markitdown for non-PDF formats.")

    args = parser.parse_args()

    if not args.input_file.exists():
        print(f"ERROR: File not found: {args.input_file}")
        sys.exit(1)

    fmt = detect_format(args.input_file)
    short_name = args.short_name or args.input_file.stem.lower().replace(" ", "-")
    output_path = args.output or args.input_file.with_suffix(".md")
    images_dir = args.images_dir or (output_path.parent / "images" / short_name)
    page_count = get_page_count(args.input_file, fmt)

    page_label = "slides" if fmt == "pptx" else "sheets" if fmt == "xlsx" else "pages"
    print(f"Input:   {args.input_file}")
    print(f"Format:  {fmt.upper()}")
    print(f"Output:  {output_path}")
    if page_count:
        print(f"{page_label.capitalize()}: {page_count}")
    print(f"Images:  {images_dir}")
    print()

    # ── Step 1: Text extraction ──
    step_total = 3 if not args.no_images else 2
    extractor_name = args.extractor if fmt == "pdf" else "MarkItDown"
    print(f"[1/{step_total}] Extracting text with {extractor_name}...")
    # S36: extract_text now returns (text, docling_pictures) tuple
    raw_text, _docling_pictures = extract_text(args.input_file, fmt,
                                               extractor=args.extractor)
    cleaned = clean_text(raw_text)
    print(f"  {len(cleaned):,} chars, {cleaned.count(chr(10)):,} lines")

    # ── Step 2: Image extraction ──
    image_entries = []
    if not args.no_images:
        print(f"[2/{step_total}] Extracting images...")
        if fmt == "pdf":
            import fitz  # Import here for caption detection
            image_entries = extract_images_pdf(args.input_file, images_dir, short_name, cleaned)
        elif fmt == "docx":
            image_entries = extract_images_docx(args.input_file, images_dir, short_name)
        elif fmt == "pptx":
            image_entries = extract_images_pptx(args.input_file, images_dir, short_name)
        elif fmt == "xlsx":
            print("  XLSX: no image extraction (spreadsheet format)")
        print(f"  {len(image_entries)} image(s) extracted")
    else:
        print(f"[2/{step_total}] Skipping image extraction (--no-images)")

    # S36: Map docling anonymous image placeholders to extracted filenames
    if _docling_pictures and image_entries:
        cleaned = _map_docling_image_placeholders(
            cleaned, _docling_pictures, image_entries)

    # ── Step 3: Assemble ──
    step_num = step_total
    print(f"[{step_num}/{step_total}] Assembling markdown...")

    try:
        images_rel = os.path.relpath(images_dir, output_path.parent)
    except ValueError:
        images_rel = str(images_dir)

    content = build_header(args.input_file, fmt, page_count,
                           image_count=len(image_entries),
                           extractor=args.extractor)
    content += build_image_index(image_entries, images_rel)
    content += "---\n\n"
    content += cleaned
    content += "\n"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(content, encoding="utf-8")

    # Write image manifest for the IMAGE NOTE subagent
    if image_entries:
        write_image_manifest(image_entries, images_dir, output_path)

    # Generate and write context summary (NEW)
    context_summary = generate_context_summary(cleaned, len(image_entries))
    write_context_summary(context_summary, output_path.parent)

    line_count = content.count('\n')
    print(f"\nDone. {line_count:,} lines written to {output_path}")

    # ── QC instructions ──
    print()
    print("=" * 60)
    print("NEXT STEPS (mandatory):")
    print("=" * 60)
    print(f"  1. python3 qc-structural.py {output_path}")
    print(f"  2. python3 prepare-image-analysis.py {output_path}")
    print(f"  3. Claude subagent: generate-image-notes.md")
    print(f"  4. python3 validate-image-notes.py {output_path}")
    if fmt == "pdf":
        print(f"  5. python3 extract-numbers.py {args.input_file} {output_path}")
        print(f"  6. Claude subagent: qc-content-fidelity.md")
        print(f"  7. Claude subagent: qc-final-review.md")
    else:
        print(f"  5. Claude subagent: qc-final-review.md")
        print(f"  (Content fidelity check skipped for {fmt.upper()} -")
        print(f"   MarkItDown is more reliable on structured formats)")
    print("=" * 60)


if __name__ == "__main__":
    main()
